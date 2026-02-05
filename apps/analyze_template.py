#!/usr/bin/env python3
"""
Template Analyzer CLI Tool

Analyzes PPTX templates and generates manifest files for the dynamic template system.

Usage:
    # Basic analysis (print summary)
    python analyze_template.py --template path/to/template.pptx
    
    # Generate manifest file
    python analyze_template.py --template path/to/template.pptx --output manifest.json
    
    # Register template to registry
    python analyze_template.py --template path/to/template.pptx --register my_template
    
    # Analyze template directory (legacy format)
    python analyze_template.py --directory path/to/template/dir
    
    # Show layout mapping suggestions
    python analyze_template.py --template path/to/template.pptx --show-mappings
"""

import argparse
import sys
import json
from pathlib import Path

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def basic_analyze(template_path: str) -> None:
    """Basic analysis without using the full template analyzer"""
    print(f"\nAnalyzing template: {template_path}")
    
    if not Path(template_path).exists():
        print(f"Error: File not found: {template_path}")
        return
    
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"Error loading presentation: {e}")
        return
    
    print(f"\nSlide Dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print(f"Slide Masters: {len(prs.slide_masters)}")
    print(f"Total Layouts: {sum(len(m.slide_layouts) for m in prs.slide_masters)}")
    
    print("\n" + "=" * 80)
    print(f"{'IDX':<5} | {'LAYOUT NAME':<35} | {'PLACEHOLDERS'}")
    print("=" * 80)
    
    for master_idx, master in enumerate(prs.slide_masters):
        if len(prs.slide_masters) > 1:
            print(f"\n[Master {master_idx}: {getattr(master, 'name', 'Unnamed')}]")
        
        for i, layout in enumerate(master.slide_layouts):
            placeholders = []
            for shape in layout.placeholders:
                ph_type = shape.placeholder_format.type
                ph_name = shape.name
                ph_idx = shape.placeholder_format.idx
                
                # Get position info
                pos = f"{shape.left.inches:.1f},{shape.top.inches:.1f}"
                placeholders.append(f"[{ph_idx}]{ph_type}@{pos}")
            
            ph_str = ", ".join(placeholders) if placeholders else "(no placeholders)"
            layout_name = layout.name[:35] if layout.name else f"Layout_{i}"
            print(f"{i:<5} | {layout_name:<35} | {ph_str}")
    
    print("=" * 80)


def full_analyze(template_path: str, output_path: str = None) -> dict:
    """Full analysis using the template analyzer module"""
    try:
        from app.services.template_analyzer import (
            TemplateAnalyzer,
            analyze_template,
            print_template_summary
        )
    except ImportError:
        print("Warning: Could not import template_analyzer module. Using basic analysis.")
        basic_analyze(template_path)
        return {}
    
    print(f"\nAnalyzing template: {template_path}")
    
    # Analyze template
    manifest = analyze_template(template_path, output_path)
    
    # Print summary
    print_template_summary(manifest)
    
    return manifest.to_dict() if hasattr(manifest, 'to_dict') else {}


def register_template(template_path: str, template_id: str) -> None:
    """Register template to the global registry"""
    try:
        from app.services.template_registry import get_registry
    except ImportError:
        print("Error: Could not import template_registry module.")
        return
    
    print(f"\nRegistering template '{template_id}' from: {template_path}")
    
    registry = get_registry()
    
    path = Path(template_path)
    if path.suffix == ".pptx":
        manifest = registry.register_from_pptx(str(path), template_id)
    elif path.is_dir():
        manifest = registry.register_from_directory(str(path), template_id)
    else:
        print(f"Error: Unsupported file type: {template_path}")
        return
    
    print(f"\nTemplate registered successfully!")
    print(f"  ID: {manifest.template_id}")
    print(f"  Name: {manifest.template_name}")
    print(f"  Layouts: {len(manifest.layouts)}")
    print(f"  Mode: {manifest.template_mode}")


def show_mappings(template_path: str) -> None:
    """Show intelligent layout mapping suggestions"""
    try:
        from app.services.template_analyzer import analyze_template
        from app.services.layout_mapper import LayoutMapper, print_all_mappings
    except ImportError:
        print("Error: Could not import required modules.")
        return
    
    print(f"\nAnalyzing layout mappings for: {template_path}")
    
    # Analyze template
    manifest = analyze_template(template_path)
    
    # Show mappings
    print_all_mappings(manifest)


def analyze_directory(dir_path: str) -> None:
    """Analyze a template directory (legacy format)"""
    dir_path = Path(dir_path)
    
    if not dir_path.exists():
        print(f"Error: Directory not found: {dir_path}")
        return
    
    print(f"\nAnalyzing template directory: {dir_path}")
    
    # Check for config files
    config_path = dir_path / "config.json"
    layouts_path = dir_path / "layouts.json"
    manifest_path = dir_path / "manifest.json"
    pptx_path = dir_path / "template.pptx"
    
    print("\nFiles found:")
    for f in [config_path, layouts_path, manifest_path, pptx_path]:
        status = "✓" if f.exists() else "✗"
        print(f"  {status} {f.name}")
    
    # Load and display config
    if config_path.exists():
        print("\n" + "=" * 60)
        print("CONFIG.JSON SUMMARY")
        print("=" * 60)
        
        with open(config_path, 'r', encoding='utf-8') as f:
            config = json.load(f)
        
        print(f"Template ID: {config.get('template_id', 'N/A')}")
        print(f"Template Name: {config.get('template_name', 'N/A')}")
        print(f"Template Mode: {config.get('template_mode', 'N/A')}")
        print(f"Version: {config.get('version', 'N/A')}")
        
        dims = config.get('slide_dimensions', {})
        print(f"Dimensions: {dims.get('width', 'N/A')}\" x {dims.get('height', 'N/A')}\"")
        
        # Content type mapping
        mapping = config.get('content_type_mapping', {})
        if mapping:
            print("\nContent Type Mapping:")
            for content_type, layout in mapping.items():
                if isinstance(layout, list):
                    layout = layout[0] if layout else "N/A"
                print(f"  {content_type:15} -> {layout}")
        
        # Background images
        backgrounds = config.get('background_images', {})
        if backgrounds:
            print(f"\nBackground Images: {len(backgrounds)} defined")
    
    # Load and display layouts
    if layouts_path.exists():
        print("\n" + "=" * 60)
        print("LAYOUTS.JSON SUMMARY")
        print("=" * 60)
        
        with open(layouts_path, 'r', encoding='utf-8') as f:
            layouts = json.load(f)
        
        print(f"Total Layouts: {len(layouts)}")
        print("\nLayout Names:")
        for name in layouts.keys():
            layout_info = layouts[name]
            bg = layout_info.get('background', 'N/A')
            elements = layout_info.get('elements', [])
            print(f"  - {name}: {len(elements)} elements, bg={bg}")
    
    # If PPTX exists, analyze it
    if pptx_path.exists():
        print("\n" + "=" * 60)
        print("TEMPLATE.PPTX ANALYSIS")
        print("=" * 60)
        basic_analyze(str(pptx_path))


def main():
    parser = argparse.ArgumentParser(
        description="Analyze PPTX templates for the dynamic template system",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python analyze_template.py --template template.pptx
  python analyze_template.py --template template.pptx --output manifest.json
  python analyze_template.py --template template.pptx --register my_template
  python analyze_template.py --directory templates/arweqah/
  python analyze_template.py --template template.pptx --show-mappings
        """
    )
    
    parser.add_argument(
        "--template", "-t",
        help="Path to .pptx template file"
    )
    
    parser.add_argument(
        "--directory", "-d",
        help="Path to template directory (for legacy format)"
    )
    
    parser.add_argument(
        "--output", "-o",
        help="Output path for manifest JSON file"
    )
    
    parser.add_argument(
        "--register", "-r",
        help="Register template with given ID to the registry"
    )
    
    parser.add_argument(
        "--show-mappings", "-m",
        action="store_true",
        help="Show intelligent layout mapping suggestions"
    )
    
    parser.add_argument(
        "--basic",
        action="store_true",
        help="Use basic analysis only (no module imports)"
    )
    
    args = parser.parse_args()
    
    if not args.template and not args.directory:
        parser.print_help()
        sys.exit(1)
    
    # Analyze directory
    if args.directory:
        analyze_directory(args.directory)
        return
    
    # Check template file exists
    if args.template and not Path(args.template).exists():
        print(f"Error: Template file not found: {args.template}")
        sys.exit(1)
    
    # Basic analysis only
    if args.basic:
        basic_analyze(args.template)
        return
    
    # Register to registry
    if args.register:
        register_template(args.template, args.register)
        return
    
    # Show mappings
    if args.show_mappings:
        show_mappings(args.template)
        return
    
    # Full analysis with optional output
    try:
        full_analyze(args.template, args.output)
    except Exception as e:
        print(f"Warning: Full analysis failed ({e}), using basic analysis")
        basic_analyze(args.template)


if __name__ == "__main__":
    main()
