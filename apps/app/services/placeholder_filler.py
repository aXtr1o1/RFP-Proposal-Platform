"""
Placeholder Filler Module
Handles content insertion into native PowerPoint placeholders.

This module provides utilities for filling content into PPTX placeholders
while maintaining proper formatting, RTL support, and editability.
"""

import logging
from typing import Dict, List, Optional, Any, Union
from pathlib import Path

from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.placeholder import PlaceholderPicture
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

from ..models.template_manifest import (
    TemplateManifest,
    LayoutDefinition,
    PlaceholderSlot,
    LanguageConfig,
)
from ..models.presentation import SlideContent, BulletPoint, TableData

logger = logging.getLogger("placeholder_filler")


# ============================================================================
# PLACEHOLDER TYPE CONSTANTS
# ============================================================================

# PowerPoint placeholder type values
PH_TYPE_TITLE = 1
PH_TYPE_BODY = 2
PH_TYPE_CENTER_TITLE = 3
PH_TYPE_SUBTITLE = 4
PH_TYPE_DATE = 5
PH_TYPE_SLIDE_NUMBER = 6
PH_TYPE_FOOTER = 7
PH_TYPE_PICTURE = 18
PH_TYPE_OBJECT = 9
PH_TYPE_CHART = 10
PH_TYPE_TABLE = 11


# ============================================================================
# PLACEHOLDER FILLER CLASS
# ============================================================================

class PlaceholderFiller:
    """
    Fills content into PowerPoint placeholders.
    
    This class handles the insertion of various content types into
    native PowerPoint placeholders, maintaining proper formatting
    and editability.
    
    Usage:
        filler = PlaceholderFiller(manifest, language="en")
        filler.fill_slide(slide, slide_content)
    """
    
    def __init__(
        self,
        manifest: Optional[TemplateManifest] = None,
        language: str = "en",
        lang_config: Optional[Dict[str, Any]] = None
    ):
        """
        Initialize the placeholder filler.
        
        Args:
            manifest: Template manifest with layout definitions
            language: Target language code (e.g., "en", "ar")
            lang_config: Language-specific configuration
        """
        self.manifest = manifest
        self.language = language
        self.lang_config = lang_config or {}
        self.is_rtl = self.lang_config.get("rtl", False)
        
        # Default fonts
        self.default_font = self.lang_config.get("default_font", "Calibri")
        self.heading_font = self.lang_config.get("heading_font", "Calibri Light")
    
    def fill_slide(
        self,
        slide: Slide,
        content: SlideContent,
        layout_def: Optional[LayoutDefinition] = None
    ) -> None:
        """
        Fill a slide with content using its placeholders.
        
        Args:
            slide: The PowerPoint slide to fill
            content: SlideContent object with the content to insert
            layout_def: Optional layout definition for placeholder info
        """
        logger.debug(f"Filling slide: {content.title[:30] if content.title else 'No title'}...")
        
        # Fill title placeholder
        if content.title:
            self._fill_title(slide, content.title)
        
        # Fill subtitle if present
        if content.subtitle:
            self._fill_subtitle(slide, content.subtitle)
        
        # Fill body content based on content type
        if content.bullets:
            self._fill_bullets(slide, content.bullets)
        elif content.paragraph:
            self._fill_paragraph(slide, content.paragraph)
        elif content.table_data:
            self._fill_table(slide, content.table_data)
        elif content.chart_data:
            self._fill_chart(slide, content.chart_data)
    
    def _fill_title(self, slide: Slide, title: str) -> None:
        """Fill the title placeholder"""
        for shape in slide.placeholders:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type in [PH_TYPE_TITLE, PH_TYPE_CENTER_TITLE]:
                    self._set_text_with_formatting(shape, title, is_title=True)
                    logger.debug(f"  Title filled: {title[:30]}...")
                    return
            except Exception as e:
                continue
        
        logger.warning("No title placeholder found")
    
    def _fill_subtitle(self, slide: Slide, subtitle: str) -> None:
        """Fill the subtitle placeholder"""
        for shape in slide.placeholders:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type == PH_TYPE_SUBTITLE:
                    self._set_text_with_formatting(shape, subtitle, is_title=False)
                    logger.debug(f"  Subtitle filled: {subtitle[:30]}...")
                    return
            except Exception as e:
                continue
        
        logger.debug("No subtitle placeholder found")
    
    def _fill_bullets(self, slide: Slide, bullets: List[BulletPoint]) -> None:
        """Fill the body placeholder with bullet points"""
        body_shape = self._find_body_placeholder(slide)
        if not body_shape:
            logger.warning("No body placeholder found for bullets")
            return
        
        tf = body_shape.text_frame
        tf.clear()
        
        for i, bullet in enumerate(bullets):
            # Add paragraph
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # Set bullet text
            p.text = bullet.text
            p.level = bullet.level if hasattr(bullet, 'level') else 0
            
            # Apply formatting
            self._apply_paragraph_formatting(p, is_bullet=True)
            
            # Handle sub-bullets
            if hasattr(bullet, 'sub_bullets') and bullet.sub_bullets:
                for sub_bullet in bullet.sub_bullets:
                    sub_p = tf.add_paragraph()
                    sub_p.text = sub_bullet.text if hasattr(sub_bullet, 'text') else str(sub_bullet)
                    sub_p.level = 1
                    self._apply_paragraph_formatting(sub_p, is_bullet=True)
        
        logger.debug(f"  Bullets filled: {len(bullets)} items")
    
    def _fill_paragraph(self, slide: Slide, paragraph: str) -> None:
        """Fill the body placeholder with paragraph text"""
        body_shape = self._find_body_placeholder(slide)
        if not body_shape:
            logger.warning("No body placeholder found for paragraph")
            return
        
        tf = body_shape.text_frame
        tf.clear()
        
        # Split by newlines for multiple paragraphs
        paragraphs = paragraph.split('\n')
        
        for i, para_text in enumerate(paragraphs):
            if not para_text.strip():
                continue
            
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = para_text.strip()
            self._apply_paragraph_formatting(p, is_bullet=False)
        
        logger.debug(f"  Paragraph filled: {len(paragraph)} chars")
    
    def _fill_table(self, slide: Slide, table_data: TableData) -> None:
        """
        Fill or create a table in the slide.
        
        For native PPTX templates, we try to use the table placeholder if available,
        otherwise we create a table shape at a standard position.
        """
        # Try to find table placeholder
        table_placeholder = None
        body_placeholder = None
        
        for shape in slide.placeholders:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type == PH_TYPE_TABLE:
                    table_placeholder = shape
                    break
                elif ph_type in [PH_TYPE_BODY, PH_TYPE_OBJECT]:
                    body_placeholder = shape
            except Exception:
                continue
        
        # Get table dimensions
        rows = len(table_data.rows) + (1 if table_data.headers else 0)
        cols = len(table_data.headers) if table_data.headers else (len(table_data.rows[0]) if table_data.rows else 0)
        
        if cols == 0 or rows == 0:
            logger.warning("Empty table data, skipping")
            return
        
        # Create table
        if table_placeholder:
            # Insert table into placeholder
            table = table_placeholder.insert_table(rows, cols).table
        elif body_placeholder:
            # Get position from body placeholder
            left = body_placeholder.left
            top = body_placeholder.top
            width = body_placeholder.width
            height = body_placeholder.height
            
            # Add table shape
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
        else:
            # Use default position
            left = Inches(1.0)
            top = Inches(1.8)
            width = Inches(8.0)
            height = Inches(5.0)
            
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
        
        # Fill headers
        row_idx = 0
        if table_data.headers:
            for col_idx, header in enumerate(table_data.headers):
                cell = table.cell(0, col_idx)
                cell.text = str(header)
                self._format_table_cell(cell, is_header=True)
            row_idx = 1
        
        # Fill data rows
        for data_row in table_data.rows:
            for col_idx, cell_value in enumerate(data_row):
                if col_idx < cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_value)
                    self._format_table_cell(cell, is_header=False)
            row_idx += 1
        
        logger.debug(f"  Table filled: {rows}x{cols}")
    
    def _fill_chart(self, slide: Slide, chart_data: Dict[str, Any]) -> None:
        """
        Fill or create a chart in the slide.
        
        Note: Chart creation requires chart data in the proper format.
        This is a placeholder for chart generation logic.
        """
        logger.debug("  Chart placeholder - implement with chart_service")
        # Chart creation will be handled by chart_service
        pass
    
    def _find_body_placeholder(self, slide: Slide) -> Optional[Any]:
        """Find the body/content placeholder in a slide"""
        for shape in slide.placeholders:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type in [PH_TYPE_BODY, PH_TYPE_OBJECT]:
                    return shape
            except Exception:
                continue
        return None
    
    def _set_text_with_formatting(
        self,
        shape,
        text: str,
        is_title: bool = False
    ) -> None:
        """Set text on a shape with proper formatting"""
        if not hasattr(shape, 'text_frame'):
            shape.text = text
            return
        
        tf = shape.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = text
        
        # Apply formatting
        self._apply_paragraph_formatting(p, is_bullet=False, is_title=is_title)
    
    def _apply_paragraph_formatting(
        self,
        paragraph,
        is_bullet: bool = False,
        is_title: bool = False
    ) -> None:
        """Apply RTL/LTR and font formatting to a paragraph"""
        # Set alignment based on RTL
        if self.is_rtl:
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = PP_ALIGN.LEFT if not is_title else PP_ALIGN.CENTER
        
        # Apply font
        font = paragraph.font
        if is_title:
            font.name = self.heading_font
        else:
            font.name = self.default_font
        
        # Apply RTL at run level
        if self.is_rtl:
            self._apply_rtl_to_paragraph(paragraph)
    
    def _apply_rtl_to_paragraph(self, paragraph) -> None:
        """Apply RTL text direction to paragraph XML"""
        try:
            p_elem = paragraph._p
            pPr = p_elem.find(qn('a:pPr'))
            if pPr is None:
                pPr = parse_xml('<a:pPr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>')
                p_elem.insert(0, pPr)
            pPr.set('rtl', '1')
        except Exception as e:
            logger.debug(f"Could not apply RTL: {e}")
    
    def _format_table_cell(self, cell, is_header: bool = False) -> None:
        """Format a table cell"""
        # Apply text formatting
        for paragraph in cell.text_frame.paragraphs:
            self._apply_paragraph_formatting(paragraph, is_bullet=False)
            
            # Set font
            if is_header:
                paragraph.font.bold = True
                paragraph.font.name = self.heading_font
            else:
                paragraph.font.name = self.default_font


# ============================================================================
# CONVENIENCE FUNCTIONS
# ============================================================================

def fill_slide_content(
    slide: Slide,
    content: SlideContent,
    language: str = "en",
    lang_config: Optional[Dict] = None
) -> None:
    """
    Convenience function to fill a slide with content.
    
    Args:
        slide: The PowerPoint slide
        content: SlideContent object
        language: Language code
        lang_config: Language configuration
    """
    filler = PlaceholderFiller(language=language, lang_config=lang_config)
    filler.fill_slide(slide, content)


def get_placeholder_by_type(slide: Slide, ph_type: int) -> Optional[Any]:
    """
    Get a placeholder from a slide by its type.
    
    Args:
        slide: The PowerPoint slide
        ph_type: Placeholder type constant
        
    Returns:
        Placeholder shape or None
    """
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.type == ph_type:
                return shape
        except Exception:
            continue
    return None


def get_placeholder_by_idx(slide: Slide, idx: int) -> Optional[Any]:
    """
    Get a placeholder from a slide by its index.
    
    Args:
        slide: The PowerPoint slide
        idx: Placeholder index
        
    Returns:
        Placeholder shape or None
    """
    try:
        return slide.placeholders[idx]
    except (KeyError, IndexError):
        return None


def list_placeholders(slide: Slide) -> List[Dict[str, Any]]:
    """
    List all placeholders in a slide with their properties.
    
    Args:
        slide: The PowerPoint slide
        
    Returns:
        List of placeholder info dictionaries
    """
    placeholders = []
    
    for shape in slide.placeholders:
        try:
            ph_format = shape.placeholder_format
            placeholders.append({
                "idx": ph_format.idx,
                "type": ph_format.type,
                "name": shape.name,
                "has_text_frame": hasattr(shape, 'text_frame'),
            })
        except Exception as e:
            logger.debug(f"Error getting placeholder info: {e}")
    
    return placeholders
