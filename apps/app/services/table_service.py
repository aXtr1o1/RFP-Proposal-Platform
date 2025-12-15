import json
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from typing import Dict, Union, Tuple
from pathlib import Path

from apps.app.config import settings

class TableService:
    """Table generation with multilingual RTL/LTR support - fully dynamic from constraints.json"""

    def __init__(self, template_id="arweqah", language="en"):
        self.template_id = template_id
        self.language = language
        
        # Load constraints dynamically
        self.constraints = self._load_constraints(template_id)
        
        # Get table configuration from constraints
        tbl = self.constraints.get("table", {})
        colors = self.constraints.get("colors", {})
        typography = self.constraints.get("typography", {})
        
        # Get language-specific fonts
        fonts = typography.get('fonts', {})
        lang_fonts = fonts.get(language, fonts.get('en', {}))
        
        # Colors from constraints
        self.header_color = self._get_color_rgb(tbl, 'header_bg', colors, "#01415C")
        self.header_text_color = self._get_color_rgb(tbl, 'header_text', colors, "#FFFCEC")
        self.alt_row_color = self._get_color_rgb(tbl, 'alternate_row_bg', colors, "#F7F4E7")
        self.text_color = self._get_color_rgb(tbl, 'body_text', colors, "#0D2026")
        self.border_color = self._get_color_rgb(tbl, 'border_color', colors, "#C6C3BE")
        
        # Configuration from constraints
        self.border_width = tbl.get("border_width", 1)
        self.max_rows_per_slide = tbl.get("max_rows", 20)
        self.cell_padding = tbl.get("cell_padding", 0.1)
        
        # Fonts from constraints (language-specific)
        self.header_font = lang_fonts.get('heading', 'Cairo')
        self.body_font = lang_fonts.get('body', 'Tajawal')
        
        # Font sizes from constraints
        font_sizes = typography.get('font_sizes', {})
        self.header_font_size = tbl.get("header_font_size", font_sizes.get('table_header', 16))
        self.body_font_size = tbl.get("body_font_size", font_sizes.get('table_body', 14))
        
        # Alignment from constraints (language-specific)
        self.header_alignment = tbl.get("header_alignment", "center")
        self.header_bold = tbl.get("header_bold", True)
        
        # Get language-specific body alignment from constraints
        alignment_config = self.constraints.get('alignment', {})
        lang_alignment = alignment_config.get(language, alignment_config.get('en', {}))
        
        if language == "ar":
            self.body_alignment = tbl.get("body_alignment_ar", lang_alignment.get('table_body', 'right'))
        else:
            self.body_alignment = tbl.get("body_alignment_en", lang_alignment.get('table_body', 'left'))
        
        # RTL support
        self.rtl_support = tbl.get("rtl_support", False)
        self.is_rtl = (language == "ar")
        
        logger_name = "table_service"
        import logging
        self.logger = logging.getLogger(logger_name)
        self.logger.info(f"✅ TableService initialized (template={template_id}, lang={language}, RTL={self.is_rtl})")

    def _get_color_rgb(self, config: Dict, key: str, colors_fallback: Dict, default: str) -> Tuple[int, int, int]:
        """Get color as RGB from config, with fallback to colors section"""
        # Try to get from table config first
        color_hex = config.get(key)
        
        # If not found, try colors section
        if not color_hex:
            color_hex = colors_fallback.get(key, default)
        
        return self._hex_to_rgb(color_hex)

    def add_table(self, slide, table_data: Union[Dict, object], position: Dict, size: Dict):
        """Add table with RTL/LTR support - fully styled from constraints (NO rounded background)"""
        try:
            # Extract data
            if isinstance(table_data, dict):
                headers = table_data.get('headers', [])
                rows = table_data.get('rows', [])
            else:
                headers = getattr(table_data, 'headers', [])
                rows = getattr(table_data, 'rows', [])

            # ✅ FIX: Remove duplicate headers if first row matches headers
            if headers and rows and len(rows) > 0:
                first_row = rows[0]
                # Check if first row is the same as headers (case-insensitive)
                if len(first_row) == len(headers):
                    first_row_clean = [str(cell).strip().lower() for cell in first_row]
                    headers_clean = [str(h).strip().lower() for h in headers]
                    
                    if first_row_clean == headers_clean:
                        self.logger.warning("   ⚠️  Duplicate headers detected in first row - removing")
                        rows = rows[1:]  # Skip the duplicate header row

            # ✅ ENHANCED VALIDATION - Prevent empty tables
            if not headers or len(headers) == 0:
                self.logger.error(f"❌ Table has no headers")
                self.logger.error(f"   Table data received: {table_data}")
                # Create error placeholder table
                headers = ["Column 1", "Column 2", "Column 3"]
                rows = [["No data", "Table not generated", "Check source content"]]
                self.logger.warning("   ⚠️  Creating error placeholder table")
            
            if not rows or len(rows) == 0:
                self.logger.error(f"❌ Table has no rows")
                self.logger.error(f"   Headers were: {headers}")
                # Create placeholder row
                rows = [["No data provided" for _ in headers]]
                self.logger.warning("   ⚠️  Creating placeholder row")
            
            # Clean headers
            headers = [str(h).strip() if h else f"Column {i+1}" for i, h in enumerate(headers)]
            
            # Validate rows (use max_rows from constraints)
            validated_rows = []
            num_cols = len(headers)
            
            for row_idx, row in enumerate(rows[:self.max_rows_per_slide]):
                # Convert row to list if needed
                if not isinstance(row, list):
                    row = list(row) if hasattr(row, '__iter__') else [str(row)]
                
                # Clean and pad/truncate row
                clean_row = [str(cell).strip() if cell else '' for cell in row]
                
                if len(clean_row) < num_cols:
                    self.logger.warning(f"   Row {row_idx} has only {len(clean_row)} cells, padding to {num_cols}")
                    clean_row.extend([''] * (num_cols - len(clean_row)))
                elif len(clean_row) > num_cols:
                    self.logger.warning(f"   Row {row_idx} has {len(clean_row)} cells, truncating to {num_cols}")
                    clean_row = clean_row[:num_cols]
                
                validated_rows.append(clean_row)
            
            rows = validated_rows
            
            self.logger.info(f"   ✅ Table validated: {num_cols} columns × {len(rows)} rows")

            num_cols = len(headers)
            num_rows = len(validated_rows) + 1
            left, top = Inches(position['left']), Inches(position['top'])
            width, height = Inches(size['width']), Inches(size['height'])

            # ❌ REMOVED: Rounded rectangle background
            # No background shape is added anymore

            # Create table
            table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
            table = table_shape.table
            
            # Set column widths
            col_width_emu = int(width / num_cols)
            for col_idx in range(num_cols):
                table.columns[col_idx].width = col_width_emu

            # Style header row (using constraints)
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.word_wrap = True
                cell.text_frame.margin_left = Inches(self.cell_padding)
                cell.text_frame.margin_right = Inches(self.cell_padding)
                cell.text_frame.margin_top = Inches(0.08)
                cell.text_frame.margin_bottom = Inches(0.08)
                
                # Header background (from constraints)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*self.header_color)
                
                # Header text (from constraints)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = self.header_bold
                paragraph.font.size = Pt(self.header_font_size)
                paragraph.font.color.rgb = RGBColor(*self.header_text_color)
                paragraph.font.name = self.header_font
                paragraph.alignment = self._get_alignment(self.header_alignment)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                
                self._remove_cell_borders(cell)

            # Style body rows (using constraints)
            for row_idx, row_data in enumerate(validated_rows, start=1):
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = cell_value
                    cell.text_frame.word_wrap = True
                    cell.text_frame.margin_left = Inches(self.cell_padding)
                    cell.text_frame.margin_right = Inches(self.cell_padding)
                    cell.text_frame.margin_top = Inches(0.05)
                    cell.text_frame.margin_bottom = Inches(0.05)
                    
                    # Alternating row colors (from constraints)
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*self.alt_row_color)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    
                    # Body text (from constraints)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(self.body_font_size)
                    paragraph.font.color.rgb = RGBColor(*self.text_color)
                    paragraph.font.name = self.body_font
                    paragraph.line_spacing = 1.2
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # Language-specific alignment (from constraints)
                    if col_idx == 0:
                        # First column follows body alignment (RTL/LTR aware)
                        paragraph.alignment = self._get_alignment(self.body_alignment)
                    else:
                        # Other columns center-aligned
                        paragraph.alignment = PP_ALIGN.CENTER
                    
                    self._add_subtle_inner_borders(cell, row_idx, col_idx, num_rows, num_cols)
            
            self.logger.info(f"✅ Table rendered: {num_cols} cols, {len(validated_rows)} rows (RTL={self.is_rtl})")
            return table
        except Exception as e:
            self.logger.error(f"  Table error: {e}")
            import traceback
            traceback.print_exc()
            return None

    def _get_alignment(self, alignment: str) -> PP_ALIGN:
        """Convert alignment string to PP_ALIGN enum"""
        alignment = alignment.lower()
        if alignment == 'center':
            return PP_ALIGN.CENTER
        elif alignment == 'right':
            return PP_ALIGN.RIGHT
        elif alignment == 'justify':
            return PP_ALIGN.JUSTIFY
        else:
            return PP_ALIGN.LEFT

    def _remove_cell_borders(self, cell):
        """Remove cell borders"""
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                existing = tcPr.find(f'.//{{{tcPr.nsmap.get("a", "http://schemas.openxmlformats.org/drawingml/2006/main")}}}{border_name}')
                if existing is not None:
                    tcPr.remove(existing)
                ln = OxmlElement(f'a:{border_name}')
                ln.set('w', '0')
                noFill = OxmlElement('a:noFill')
                ln.append(noFill)
                tcPr.append(ln)
        except:
            pass

    def _add_subtle_inner_borders(self, cell, row_idx, col_idx, num_rows, num_cols):
        """Add subtle inner borders (using border color from constraints)"""
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            if col_idx < num_cols - 1:
                self._set_border(tcPr, 'lnR', self.border_color, width=6350)
            if row_idx < num_rows - 1:
                self._set_border(tcPr, 'lnB', self.border_color, width=6350)
        except:
            pass

    def _set_border(self, tcPr, border_name: str, color_rgb: tuple, width: int = 12700):
        """Set cell border"""
        try:
            ln = OxmlElement(f'a:{border_name}')
            ln.set('w', str(width))
            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            srgbClr.set('val', '%02x%02x%02x' % color_rgb)
            solidFill.append(srgbClr)
            ln.append(solidFill)
            tcPr.append(ln)
        except:
            pass

    def _load_constraints(self, template_id):
        """Load constraints.json from template directory"""
        path = Path(settings.TEMPLATES_DIR) / template_id / "constraints.json"
        with open(path, "r", encoding="utf-8") as file:
            return json.load(file)

    def _hex_to_rgb(self, hex_color) -> Tuple[int, int, int]:
        """Convert hex to RGB"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))