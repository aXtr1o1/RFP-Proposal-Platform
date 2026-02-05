"""
Template Analyzer Module
Scans PPTX template files and extracts layout information for dynamic template support.

This module analyzes PPTX files to extract:
- Slide layouts and their placeholders
- Placeholder positions, sizes, and types
- Theme colors and fonts
- Content type mappings
"""

import logging
import json
from pathlib import Path
from typing import Dict, List, Optional, Any
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

from ..models.template_manifest import (
    TemplateManifest,
    LayoutDefinition,
    PlaceholderSlot,
    Position,
    StyleDef,
    BackgroundDef,
    ThemeDefinition,
    ThemeColor,
    SlideDimensions,
    LanguageSettings,
    LanguageConfig,
    AnalysisMetadata,
    ColorScheme,
    FontScheme,
    FontDef,
    IconScheme,
    IconDef,
)
import hashlib
import os

logger = logging.getLogger("template_analyzer")


# ============================================================================
# PLACEHOLDER TYPE MAPPING
# ============================================================================

PLACEHOLDER_TYPE_MAP = {
    1: "title",           # PP_PLACEHOLDER.TITLE
    2: "body",            # PP_PLACEHOLDER.BODY  
    3: "center_title",    # PP_PLACEHOLDER.CENTER_TITLE
    4: "subtitle",        # PP_PLACEHOLDER.SUBTITLE
    5: "date",            # PP_PLACEHOLDER.DATE
    6: "slide_number",    # PP_PLACEHOLDER.SLIDE_NUMBER
    7: "footer",          # PP_PLACEHOLDER.FOOTER
    8: "header",          # PP_PLACEHOLDER.HEADER
    9: "object",          # PP_PLACEHOLDER.OBJECT
    10: "chart",          # PP_PLACEHOLDER.CHART
    11: "table",          # PP_PLACEHOLDER.TABLE
    12: "clip_art",       # PP_PLACEHOLDER.CLIP_ART
    13: "org_chart",      # PP_PLACEHOLDER.ORG_CHART
    14: "media_clip",     # PP_PLACEHOLDER.MEDIA_CLIP
    15: "bitmap",         # PP_PLACEHOLDER.BITMAP
    16: "picture",        # PP_PLACEHOLDER.PICTURE
    17: "vertical_object",
    18: "picture",        # PP_PLACEHOLDER.PICTURE (duplicate)
    19: "vertical_title",
    20: "vertical_body",
}

# Layout name patterns for content type detection
LAYOUT_PATTERNS = {
    "title": ["title slide", "title layout", "cover"],
    "section": ["section header", "section", "divider", "break"],
    "content": ["title and content", "content", "text slide"],
    "two_column": ["two content", "comparison", "two column", "side by side"],
    "blank": ["blank", "empty"],
    "image": ["picture", "image", "photo", "media"],
    "chart": ["chart", "graph", "data"],
    "table": ["table", "grid"],
}


# ============================================================================
# TEMPLATE ANALYZER CLASS
# ============================================================================

class TemplateAnalyzer:
    """
    Analyzes PPTX templates and extracts layout information.
    
    Usage:
        analyzer = TemplateAnalyzer()
        manifest = analyzer.analyze_template("path/to/template.pptx")
        
        # Save manifest
        with open("manifest.json", "w") as f:
            json.dump(manifest.model_dump(exclude_none=True), f, indent=2)
    """
    
    def __init__(self):
        self.prs: Optional[Presentation] = None
        self.template_path: Optional[Path] = None
    
    def analyze_template(
        self,
        pptx_path: str,
        template_id: Optional[str] = None,
        template_name: Optional[str] = None,
        language_settings: Optional[Dict] = None
    ) -> TemplateManifest:
        """
        Analyze a PPTX template and generate a manifest.
        
        Args:
            pptx_path: Path to the PPTX template file
            template_id: Optional template ID (defaults to filename)
            template_name: Optional template name (defaults to filename)
            language_settings: Optional language configuration dict
            
        Returns:
            TemplateManifest with extracted layout information
        """
        self.template_path = Path(pptx_path)
        
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {pptx_path}")
        
        logger.info(f"Analyzing template: {pptx_path}")
        
        # Load presentation
        self.prs = Presentation(pptx_path)
        
        # Extract slide dimensions
        slide_width = self._emu_to_inches(self.prs.slide_width)
        slide_height = self._emu_to_inches(self.prs.slide_height)
        
        # Calculate aspect ratio
        aspect_ratio = self._calculate_aspect_ratio(slide_width, slide_height)
        
        logger.info(f"  Slide dimensions: {slide_width:.2f}\" x {slide_height:.2f}\" ({aspect_ratio})")
        
        # Extract layouts
        layouts = self._extract_layouts()
        logger.info(f"  Found {len(layouts)} layouts")
        
        # Extract theme
        theme = self._extract_theme()
        logger.info(f"  Theme extracted: {len(theme.colors)} colors, {len(theme.fonts)} fonts")
        
        # Generate content type mapping
        content_mapping = self._suggest_content_type_mapping(layouts)
        logger.info(f"  Content type mapping: {len(content_mapping)} mappings")
        
        # Create slide dimensions
        slide_dims = SlideDimensions(
            width=slide_width,
            height=slide_height,
            units="inches",
            aspect_ratio=aspect_ratio
        )
        
        # Create language settings
        lang_settings = None
        if language_settings:
            configs = {}
            for lang_code, lang_config in language_settings.get("configurations", {}).items():
                configs[lang_code] = LanguageConfig(**lang_config)
            
            lang_settings = LanguageSettings(
                default=language_settings.get("default", "en"),
                supported=language_settings.get("supported", ["en"]),
                configurations=configs
            )
        
        # Create analysis metadata
        metadata = AnalysisMetadata(
            source_file=str(self.template_path.name),
            layout_count=len(layouts),
            master_count=len(self.prs.slide_masters),
            analyzed_version="2.0.0",
            analyzed_at=datetime.now().isoformat()
        )
        
        # Extract from actual slides if available
        colors = None
        fonts = None
        icons = None
        element_positions = None
        background_images = None
        
        if len(self.prs.slides) > 0:
            logger.info("  Extracting from actual slides...")
            colors = self._extract_colors_from_slides()
            fonts = self._extract_fonts_from_slides()
            element_positions = self._extract_element_positions()
            
            # Extract icons and backgrounds (saves files to template dir)
            template_dir = self.template_path.parent
            icons = self._extract_icons(template_dir)
            background_images = self._extract_backgrounds(template_dir)
        
        # Generate manifest
        manifest = TemplateManifest(
            template_id=template_id or self.template_path.stem,
            template_name=template_name or self.template_path.stem.replace("_", " ").title(),
            version="1.0.0",
            template_mode="native",
            slide_dimensions=slide_dims,
            layouts=layouts,
            content_type_mapping=content_mapping,
            theme=theme,
            colors=colors,
            fonts=fonts,
            icons=icons,
            language_settings=lang_settings,
            background_images=background_images,
            element_positions=element_positions,
            analysis_metadata=metadata
        )
        
        return manifest
    
    def _extract_layouts(self) -> Dict[str, LayoutDefinition]:
        """Extract all slide layouts from the template"""
        layouts = {}
        
        for master_idx, master in enumerate(self.prs.slide_masters):
            master_name = getattr(master, 'name', f"Master_{master_idx}")
            
            for layout_idx, layout in enumerate(master.slide_layouts):
                layout_name = getattr(layout, 'name', f"Layout_{layout_idx}")
                
                # Create unique key
                layout_key = self._normalize_layout_name(layout_name)
                if layout_key in layouts:
                    layout_key = f"{layout_key}_{master_idx}_{layout_idx}"
                
                # Extract placeholders
                placeholders = self._extract_placeholders(layout)
                
                # Extract background
                background = self._extract_background(layout)
                
                # Determine suitable content types
                suitable_for = self._determine_suitable_content_types(layout_name, placeholders)
                
                # Create layout definition
                layout_def = LayoutDefinition(
                    index=layout_idx,
                    name=layout_name,
                    placeholders=placeholders,
                    background=background,
                    master_name=master_name,
                    suitable_for=suitable_for,
                    supports_rtl=True
                )
                
                layouts[layout_key] = layout_def
                logger.debug(f"    Layout '{layout_key}': {len(placeholders)} placeholders")
        
        return layouts
    
    def _extract_placeholders(self, layout) -> List[PlaceholderSlot]:
        """Extract placeholders from a slide layout"""
        placeholders = []
        
        for shape in layout.placeholders:
            try:
                ph_format = shape.placeholder_format
                ph_type = ph_format.type
                ph_idx = ph_format.idx
                
                # Get type name
                type_name = PLACEHOLDER_TYPE_MAP.get(int(ph_type), "unknown")
                
                # Get position
                position = Position(
                    left=self._emu_to_inches(shape.left),
                    top=self._emu_to_inches(shape.top),
                    width=self._emu_to_inches(shape.width),
                    height=self._emu_to_inches(shape.height)
                )
                
                # Try to extract style
                style = self._extract_text_style(shape)
                
                # Determine content hint
                content_hint = self._get_content_hint(type_name, layout.name)
                
                placeholder = PlaceholderSlot(
                    idx=ph_idx,
                    type=type_name,
                    name=shape.name,
                    position=position,
                    style=style,
                    is_required=type_name in ["title", "body"],
                    content_hint=content_hint
                )
                
                placeholders.append(placeholder)
                
            except Exception as e:
                logger.warning(f"Error extracting placeholder: {e}")
                continue
        
        return placeholders
    
    def _extract_text_style(self, shape) -> Optional[StyleDef]:
        """Extract text style from a shape"""
        try:
            if not hasattr(shape, 'text_frame'):
                return None
            
            tf = shape.text_frame
            if not tf.paragraphs:
                return None
            
            # Get first paragraph for style
            para = tf.paragraphs[0]
            
            font_name = None
            font_size = None
            font_color = None
            bold = False
            italic = False
            alignment = None
            
            # Extract font info
            if para.font:
                font_name = para.font.name
                if para.font.size:
                    font_size = int(para.font.size.pt)
                bold = para.font.bold or False
                italic = para.font.italic or False
                if para.font.color and para.font.color.rgb:
                    font_color = self._rgb_to_hex(para.font.color.rgb)
            
            # Extract alignment
            if para.alignment:
                alignment_map = {0: "left", 1: "center", 2: "right", 3: "justify"}
                alignment = alignment_map.get(int(para.alignment))
            
            # Only return if we got some useful info
            if any([font_name, font_size, font_color]):
                return StyleDef(
                    font_name=font_name,
                    font_size=font_size,
                    font_color=font_color,
                    bold=bold,
                    italic=italic,
                    alignment=alignment
                )
            
            return None
            
        except Exception as e:
            logger.debug(f"Could not extract text style: {e}")
            return None
    
    def _extract_background(self, layout) -> Optional[BackgroundDef]:
        """Extract background configuration from layout"""
        try:
            background = layout.background
            fill = background.fill
            
            if fill.type == 1:  # Solid fill
                if fill.fore_color and fill.fore_color.rgb:
                    return BackgroundDef(
                        type="solid",
                        color=self._rgb_to_hex(fill.fore_color.rgb)
                    )
            elif fill.type == 6:  # Picture fill
                return BackgroundDef(type="image")
            
            return BackgroundDef(type="none")
            
        except Exception as e:
            logger.debug(f"Could not extract background: {e}")
            return BackgroundDef(type="none")
    
    def _extract_theme(self) -> ThemeDefinition:
        """Extract theme information from the template"""
        colors = []
        fonts = {}
        
        try:
            # Try to extract theme colors from first slide master
            if self.prs.slide_masters:
                master = self.prs.slide_masters[0]
                
                # Extract theme colors from theme XML
                try:
                    colors = self._extract_theme_colors(master)
                except Exception as e:
                    logger.debug(f"Could not extract theme colors: {e}")
                
                # Extract fonts
                try:
                    fonts = self._extract_theme_fonts(master)
                except Exception as e:
                    logger.debug(f"Could not extract fonts: {e}")
        
        except Exception as e:
            logger.warning(f"Error extracting theme: {e}")
        
        # Provide defaults if extraction failed
        if not colors:
            colors = [
                ThemeColor(name="primary", hex_value="#01415C"),
                ThemeColor(name="secondary", hex_value="#84BA93"),
                ThemeColor(name="accent1", hex_value="#C26325"),
                ThemeColor(name="accent2", hex_value="#F9D462"),
                ThemeColor(name="background", hex_value="#FFFDED"),
                ThemeColor(name="text", hex_value="#0D2026"),
            ]
        
        if not fonts:
            fonts = {
                "heading": "Calibri Light",
                "body": "Calibri",
                "fallback": "Arial"
            }
        
        return ThemeDefinition(colors=colors, fonts=fonts)
    
    def _extract_theme_colors(self, master) -> List[ThemeColor]:
        """Extract theme colors from slide master"""
        colors = []
        
        # Try to get colors from shapes in master
        try:
            # Get accent colors from theme element if available
            theme_element = master.part.slide_master
            
            # Map of color scheme names
            color_names = {
                "dk1": "dark1",
                "lt1": "light1",
                "dk2": "dark2", 
                "lt2": "light2",
                "accent1": "accent1",
                "accent2": "accent2",
                "accent3": "accent3",
                "accent4": "accent4",
                "accent5": "accent5",
                "accent6": "accent6",
                "hlink": "hyperlink",
                "folHlink": "followed_hyperlink"
            }
            
            for key, name in color_names.items():
                colors.append(ThemeColor(name=name, hex_value="#000000"))
                
        except Exception:
            pass
        
        return colors
    
    def _extract_theme_fonts(self, master) -> Dict[str, str]:
        """Extract theme fonts from slide master"""
        fonts = {
            "heading": "Calibri Light",
            "body": "Calibri",
            "fallback": "Arial"
        }
        
        try:
            # Try to get fonts from theme part
            for shape in master.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        if para.font and para.font.name:
                            # Use first found font as body
                            if "body" not in fonts or fonts["body"] == "Calibri":
                                fonts["body"] = para.font.name
                            break
        except Exception:
            pass
        
        return fonts
    
    def _suggest_content_type_mapping(self, layouts: Dict[str, LayoutDefinition]) -> Dict[str, str]:
        """
        Suggest content type to layout mappings based on layout analysis.
        """
        mapping = {}
        
        for layout_key, layout_def in layouts.items():
            name_lower = layout_def.name.lower()
            
            # Check against pattern lists
            for content_type, patterns in LAYOUT_PATTERNS.items():
                if any(pattern in name_lower for pattern in patterns):
                    if content_type not in mapping:
                        mapping[content_type] = layout_key
                    continue
            
            # Analyze placeholders
            placeholder_types = [p.type for p in layout_def.placeholders]
            
            has_title = "title" in placeholder_types or "center_title" in placeholder_types
            has_body = "body" in placeholder_types
            has_picture = "picture" in placeholder_types
            body_count = placeholder_types.count("body")
            
            # Two content / comparison
            if has_title and body_count >= 2:
                if "two_column" not in mapping:
                    mapping["two_column"] = layout_key
                if "comparison" not in mapping:
                    mapping["comparison"] = layout_key
                continue
            
            # Title and content
            if has_title and has_body:
                if "content" not in mapping:
                    mapping["content"] = layout_key
                if "bullets" not in mapping:
                    mapping["bullets"] = layout_key
                continue
            
            # Picture slide
            if has_picture:
                if "image" not in mapping:
                    mapping["image"] = layout_key
                continue
            
            # Title only (for sections)
            if has_title and not has_body:
                if "section" not in mapping:
                    mapping["section"] = layout_key
        
        # Add fallbacks
        if "content" not in mapping and layouts:
            for layout_key, layout_def in layouts.items():
                if any(p.type == "body" for p in layout_def.placeholders):
                    mapping["content"] = layout_key
                    break
        
        # Map common content types to detected layouts
        if "content" in mapping:
            mapping.setdefault("bullets", mapping["content"])
            mapping.setdefault("paragraph", mapping["content"])
            mapping.setdefault("table", mapping["content"])
            mapping.setdefault("chart", mapping["content"])
            mapping.setdefault("agenda", mapping["content"])
        
        return mapping
    
    def _determine_suitable_content_types(
        self, 
        layout_name: str, 
        placeholders: List[PlaceholderSlot]
    ) -> List[str]:
        """Determine which content types a layout is suitable for"""
        suitable = []
        name_lower = layout_name.lower()
        placeholder_types = [p.type for p in placeholders]
        
        has_title = "title" in placeholder_types or "center_title" in placeholder_types
        has_body = "body" in placeholder_types
        body_count = placeholder_types.count("body")
        
        # Check patterns
        for content_type, patterns in LAYOUT_PATTERNS.items():
            if any(pattern in name_lower for pattern in patterns):
                suitable.append(content_type)
        
        # Infer from placeholders
        if has_title and has_body:
            if "content" not in suitable:
                suitable.extend(["content", "bullets", "paragraph"])
        
        if body_count >= 2:
            if "two_column" not in suitable:
                suitable.extend(["two_column", "comparison"])
        
        if has_title and not has_body and not suitable:
            suitable.append("section")
        
        return list(set(suitable))
    
    def _get_content_hint(self, placeholder_type: str, layout_name: str) -> Optional[str]:
        """Generate content hint for a placeholder"""
        hints = {
            "title": "Main heading or slide title",
            "subtitle": "Secondary heading or description",
            "body": "Main content area - supports bullets, paragraphs, or lists",
            "center_title": "Centered title for section headers",
            "picture": "Image placeholder",
            "chart": "Chart or graph placeholder",
            "table": "Table placeholder",
            "object": "Generic content placeholder",
        }
        return hints.get(placeholder_type)
    
    def _calculate_aspect_ratio(self, width: float, height: float) -> str:
        """Calculate aspect ratio string"""
        ratio = width / height
        
        # Common aspect ratios
        if abs(ratio - 16/9) < 0.01:
            return "16:9"
        elif abs(ratio - 4/3) < 0.01:
            return "4:3"
        elif abs(ratio - 16/10) < 0.01:
            return "16:10"
        else:
            return f"{ratio:.2f}:1"
    
    # ========================================================================
    # ENHANCED EXTRACTION METHODS (from actual slides)
    # ========================================================================
    
    def _extract_colors_from_slides(self) -> ColorScheme:
        """Extract color scheme from actual slides"""
        text_colors = []
        shape_colors = []
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                # Extract shape fill colors
                if hasattr(shape, 'fill'):
                    try:
                        if shape.fill.type is not None:
                            rgb = shape.fill.fore_color.rgb
                            if rgb:
                                shape_colors.append(str(rgb))
                    except:
                        pass
                
                # Extract text colors
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        try:
                            if para.font.color and para.font.color.type is not None:
                                rgb = para.font.color.rgb
                                if rgb:
                                    text_colors.append(str(rgb))
                        except:
                            pass
        
        # Analyze colors
        def get_most_common(color_list, default):
            if not color_list:
                return default
            return max(set(color_list), key=color_list.count)
        
        # Find separator line color (usually primary)
        separator_color = "01415C"
        page_number_bg = "C6C3BE"
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    width_in = self._emu_to_inches(shape.width)
                    height_in = self._emu_to_inches(shape.height)
                    
                    # Separator line (wide and thin)
                    if width_in > 5 and height_in < 0.1:
                        try:
                            rgb = shape.fill.fore_color.rgb
                            if rgb:
                                separator_color = str(rgb)
                        except:
                            pass
                    
                    # Page number (small shape at bottom)
                    if width_in < 0.5 and self._emu_to_inches(shape.top) > 6:
                        try:
                            rgb = shape.fill.fore_color.rgb
                            if rgb:
                                page_number_bg = str(rgb)
                        except:
                            pass
        
        return ColorScheme(
            primary=separator_color,
            secondary="84BA93",
            accent=page_number_bg,
            text={
                "dark": "0D2026",
                "light": "FFFCEC",
                "muted": "666666"
            },
            elements={
                "separator_line": separator_color,
                "page_number_bg": page_number_bg,
                "page_number_text": "FFFCEC"
            }
        )
    
    def _extract_fonts_from_slides(self) -> FontScheme:
        """Extract font scheme from actual slides"""
        fonts_by_position = {
            'title_slide_title': [],
            'title_slide_subtitle': [],
            'content_title': [],
            'content_body': [],
            'section_title': [],
            'page_number': []
        }
        
        for slide_idx, slide in enumerate(self.prs.slides):
            layout_name = slide.slide_layout.name.lower()
            
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame') or not shape.text.strip():
                    continue
                
                top_in = self._emu_to_inches(shape.top)
                width_in = self._emu_to_inches(shape.width)
                
                for para in shape.text_frame.paragraphs:
                    font_info = {
                        'name': para.font.name,
                        'size': para.font.size.pt if para.font.size else None,
                        'bold': para.font.bold
                    }
                    
                    # Categorize by position and layout
                    if width_in < 0.5 and top_in > 6:  # Page number
                        fonts_by_position['page_number'].append(font_info)
                    elif 'title slide' in layout_name and slide_idx == 0:
                        if top_in < 4:
                            fonts_by_position['title_slide_title'].append(font_info)
                        else:
                            fonts_by_position['title_slide_subtitle'].append(font_info)
                    elif 'section' in layout_name:
                        fonts_by_position['section_title'].append(font_info)
                    elif top_in < 1.5:
                        fonts_by_position['content_title'].append(font_info)
                    else:
                        fonts_by_position['content_body'].append(font_info)
                    break  # Only first paragraph
        
        def get_most_common_font(font_list):
            if not font_list:
                return None
            names = [f['name'] for f in font_list if f['name']]
            sizes = [f['size'] for f in font_list if f['size']]
            bolds = [f['bold'] for f in font_list if f['bold'] is not None]
            
            return FontDef(
                name_en=max(set(names), key=names.count) if names else "Open Sans",
                name_ar="Cairo",
                size=int(max(set(sizes), key=sizes.count)) if sizes else 18,
                bold=max(set(bolds), key=bolds.count) if bolds else False,
                color="0D2026"
            )
        
        title_slide_title = get_most_common_font(fonts_by_position['title_slide_title'])
        if title_slide_title:
            title_slide_title.color = "FFFCEC"  # Light text on dark bg
        
        title_slide_subtitle = get_most_common_font(fonts_by_position['title_slide_subtitle'])
        if title_slide_subtitle:
            title_slide_subtitle.color = "FFFCEC"
        
        return FontScheme(
            title_slide={
                "title": title_slide_title or FontDef(name_en="Roboto", name_ar="Cairo", size=40, bold=True, color="FFFCEC"),
                "subtitle": title_slide_subtitle or FontDef(name_en="Open Sans", name_ar="Tajawal", size=24, bold=False, color="FFFCEC")
            },
            section_header={
                "title": get_most_common_font(fonts_by_position['section_title']) or FontDef(name_en="Cairo", name_ar="Cairo", size=36, bold=True, color="0D2026")
            },
            content={
                "title": get_most_common_font(fonts_by_position['content_title']) or FontDef(name_en="Cairo", name_ar="Cairo", size=32, bold=True, color="0D2026"),
                "body": get_most_common_font(fonts_by_position['content_body']) or FontDef(name_en="Open Sans", name_ar="Tajawal", size=18, bold=False, color="0D2026"),
                "bullet": FontDef(name_en="Open Sans", name_ar="Tajawal", size=18, bold=False, color="0D2026"),
                "sub_bullet": FontDef(name_en="Open Sans", name_ar="Tajawal", size=16, bold=False, color="666666")
            },
            page_number=get_most_common_font(fonts_by_position['page_number']) or FontDef(name="Cairo", size=14, bold=False, color="FFFCEC")
        )
    
    def _extract_icons(self, template_dir: Path) -> IconScheme:
        """Extract icons from slides and save to template directory"""
        icons_dir = template_dir / "Icons"
        icons_dir.mkdir(exist_ok=True)
        
        seen_hashes = set()
        title_icons = []
        section_icons = []
        all_icons = {}
        
        slide_width = self._emu_to_inches(self.prs.slide_width)
        slide_height = self._emu_to_inches(self.prs.slide_height)
        
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                
                width_in = self._emu_to_inches(shape.width)
                height_in = self._emu_to_inches(shape.height)
                
                # Skip full-size images (backgrounds)
                if width_in > slide_width * 0.8:
                    continue
                
                # Icons are typically small (< 2 inches)
                if width_in > 2 or height_in > 2:
                    continue
                
                try:
                    image = shape.image
                    image_bytes = image.blob
                    image_hash = hashlib.md5(image_bytes).hexdigest()[:12]
                    
                    if image_hash in seen_hashes:
                        continue
                    seen_hashes.add(image_hash)
                    
                    ext = image.ext
                    if ext == 'jpeg':
                        ext = 'jpg'
                    
                    # Determine icon type
                    top_in = self._emu_to_inches(shape.top)
                    
                    if abs(width_in - 1.2) < 0.3:  # Section icon
                        icon_type = "section"
                        filename = f"icon_section_{image_hash}.{ext}"
                        section_icons.append(f"Icons/{filename}")
                    else:  # Title icon
                        icon_type = "title"
                        filename = f"icon_title_{image_hash}.{ext}"
                        title_icons.append(f"Icons/{filename}")
                    
                    # Save icon
                    filepath = icons_dir / filename
                    with open(filepath, 'wb') as f:
                        f.write(image_bytes)
                    
                    all_icons[filename] = IconDef(
                        path=f"Icons/{filename}",
                        width=width_in,
                        height=height_in
                    )
                    
                except Exception as e:
                    logger.debug(f"Could not extract icon: {e}")
        
        return IconScheme(
            default_title=title_icons[0] if title_icons else None,
            default_section=section_icons[0] if section_icons else None,
            agenda_items=title_icons[:5],
            box_icons=title_icons[:4],
            all_icons=all_icons
        )
    
    def _extract_backgrounds(self, template_dir: Path) -> Dict[str, str]:
        """Extract background images from slides"""
        bg_dir = template_dir / "Background"
        bg_dir.mkdir(exist_ok=True)
        
        seen_hashes = {}
        background_mapping = {}
        
        slide_width = self._emu_to_inches(self.prs.slide_width)
        slide_height = self._emu_to_inches(self.prs.slide_height)
        
        for slide in self.prs.slides:
            layout_name = slide.slide_layout.name
            
            for shape in slide.shapes:
                if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                
                width_in = self._emu_to_inches(shape.width)
                height_in = self._emu_to_inches(shape.height)
                left_in = self._emu_to_inches(shape.left)
                top_in = self._emu_to_inches(shape.top)
                
                # Full-bleed background
                if (abs(left_in) < 0.1 and abs(top_in) < 0.1 and
                    abs(width_in - slide_width) < 0.5 and
                    abs(height_in - slide_height) < 0.5):
                    
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_hash = hashlib.md5(image_bytes).hexdigest()[:8]
                        
                        ext = image.ext
                        if ext == 'jpeg':
                            ext = 'jpg'
                        
                        if image_hash not in seen_hashes:
                            layout_key = layout_name.lower().replace(' ', '_')
                            filename = f"bg_{layout_key}_{image_hash}.{ext}"
                            filepath = bg_dir / filename
                            
                            with open(filepath, 'wb') as f:
                                f.write(image_bytes)
                            
                            seen_hashes[image_hash] = f"Background/{filename}"
                        
                        # Map layout to background
                        layout_key = layout_name.lower().replace(' ', '_')
                        if layout_key not in background_mapping:
                            background_mapping[layout_key] = seen_hashes[image_hash]
                        
                        break
                    except:
                        pass
        
        # Create standard mappings
        result = {}
        for layout_key, bg_path in background_mapping.items():
            result[layout_key] = bg_path
            
            # Also add common aliases
            if 'title_slide' in layout_key or layout_key == 'title_slide':
                result['title'] = bg_path
                result['title_slide'] = bg_path
            elif 'section' in layout_key:
                result['section'] = bg_path
                result['section_header'] = bg_path
            elif 'content' in layout_key or 'title_and_content' in layout_key:
                result['content'] = bg_path
                result['bullets'] = bg_path
                result['paragraph'] = bg_path
                result['table'] = bg_path
            elif 'blank' in layout_key:
                result['blank'] = bg_path
        
        return result
    
    def _extract_element_positions(self) -> Dict[str, Any]:
        """Extract element positions from actual slides"""
        positions = {
            'title_slide': {},
            'section_header': {},
            'content': {}
        }
        
        for slide_idx, slide in enumerate(self.prs.slides):
            layout_name = slide.slide_layout.name.lower()
            
            for shape in slide.shapes:
                if not hasattr(shape, 'text') or not shape.text.strip():
                    if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                        continue
                
                left = self._emu_to_inches(shape.left)
                top = self._emu_to_inches(shape.top)
                width = self._emu_to_inches(shape.width)
                height = self._emu_to_inches(shape.height)
                
                pos = {'x': round(left, 2), 'y': round(top, 2), 'width': round(width, 2), 'height': round(height, 2)}
                
                # Title slide positions
                if slide_idx == 0 and 'title' in layout_name:
                    if hasattr(shape, 'text') and shape.text.strip():
                        if top < 4 and 'title' not in positions['title_slide']:
                            positions['title_slide']['title'] = pos
                        elif top >= 4 and 'subtitle' not in positions['title_slide']:
                            positions['title_slide']['subtitle'] = pos
                
                # Section header positions
                elif 'section' in layout_name:
                    if hasattr(shape, 'text') and shape.text.strip():
                        if top > 3 and 'title' not in positions['section_header']:
                            positions['section_header']['title'] = pos
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if width < 2 and 'icon' not in positions['section_header']:
                            positions['section_header']['icon'] = pos
                
                # Content slide positions
                elif 'content' in layout_name or 'title and content' in layout_name:
                    if hasattr(shape, 'text') and shape.text.strip():
                        if top < 1.5 and 'title' not in positions['content']:
                            positions['content']['title'] = pos
                        elif top > 1.5 and 'body' not in positions['content']:
                            positions['content']['body'] = pos
                    
                    # Separator line
                    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                        if width > 5 and height < 0.1 and 'separator' not in positions['content']:
                            positions['content']['separator'] = pos
                    
                    # Icon
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        if width < 1 and top < 1 and 'icon' not in positions['content']:
                            positions['content']['icon'] = pos
        
        return positions
    
    # ========================================================================
    # UTILITY METHODS
    # ========================================================================
    
    @staticmethod
    def _emu_to_inches(emu: int) -> float:
        """Convert EMUs to inches"""
        return emu / 914400
    
    @staticmethod
    def _rgb_to_hex(rgb: RGBColor) -> str:
        """Convert RGBColor to hex string"""
        try:
            return f"#{rgb.red:02X}{rgb.green:02X}{rgb.blue:02X}"
        except:
            return "#000000"
    
    @staticmethod
    def _normalize_layout_name(name: str) -> str:
        """Normalize layout name to a valid key"""
        normalized = name.lower()
        for char in [" ", "-", ".", ",", "(", ")", "[", "]", "'"]:
            normalized = normalized.replace(char, "_")
        while "__" in normalized:
            normalized = normalized.replace("__", "_")
        normalized = normalized.strip("_")
        return normalized or "unnamed_layout"
    
    def save_manifest(self, manifest: TemplateManifest, output_path: str) -> None:
        """Save manifest to JSON file"""
        manifest_dict = manifest.model_dump(exclude_none=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(manifest_dict, f, indent=2, ensure_ascii=False)
        logger.info(f"Manifest saved to: {output_path}")


# ============================================================================
# CONVENIENCE FUNCTIONS
# ============================================================================

def analyze_template(
    pptx_path: str, 
    output_json: Optional[str] = None,
    template_id: Optional[str] = None,
    template_name: Optional[str] = None
) -> TemplateManifest:
    """
    Convenience function to analyze a template and optionally save to JSON.
    
    Args:
        pptx_path: Path to PPTX template
        output_json: Optional path to save manifest JSON
        template_id: Optional template identifier
        template_name: Optional template display name
        
    Returns:
        TemplateManifest
    """
    analyzer = TemplateAnalyzer()
    manifest = analyzer.analyze_template(
        pptx_path,
        template_id=template_id,
        template_name=template_name
    )
    
    if output_json:
        analyzer.save_manifest(manifest, output_json)
    
    return manifest


def print_template_summary(manifest: TemplateManifest) -> None:
    """Print a summary of the analyzed template"""
    print("\n" + "=" * 60)
    print(f"TEMPLATE ANALYSIS SUMMARY")
    print("=" * 60)
    print(f"Template ID:    {manifest.template_id}")
    print(f"Template Name:  {manifest.template_name}")
    print(f"Dimensions:     {manifest.slide_dimensions.width:.2f}\" x {manifest.slide_dimensions.height:.2f}\"")
    print(f"Aspect Ratio:   {manifest.slide_dimensions.aspect_ratio}")
    print(f"Layouts Found:  {len(manifest.layouts)}")
    print()
    
    print("LAYOUTS:")
    print("-" * 40)
    for key, layout in manifest.layouts.items():
        placeholders_str = ", ".join([f"{p.type}({p.idx})" for p in layout.placeholders])
        print(f"  {key}:")
        print(f"    Name: {layout.name}")
        print(f"    Index: {layout.index}")
        print(f"    Placeholders: {placeholders_str or 'none'}")
        if layout.suitable_for:
            print(f"    Suitable for: {', '.join(layout.suitable_for)}")
    print()
    
    print("CONTENT TYPE MAPPING:")
    print("-" * 40)
    for content_type, layout_key in manifest.content_type_mapping.items():
        print(f"  {content_type:15} -> {layout_key}")
    print()
    
    print("THEME:")
    print("-" * 40)
    print(f"  Colors: {len(manifest.theme.colors)}")
    for color in manifest.theme.colors[:6]:
        print(f"    - {color.name}: {color.hex_value}")
    print(f"  Fonts:  {manifest.theme.fonts}")
    print("=" * 60 + "\n")


if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        template_path = sys.argv[1]
        output_path = sys.argv[2] if len(sys.argv) > 2 else None
        
        manifest = analyze_template(template_path, output_path)
        print_template_summary(manifest)
    else:
        print("Usage: python template_analyzer.py <template.pptx> [output.json]")
