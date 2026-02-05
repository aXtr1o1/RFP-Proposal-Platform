"""
Slide Builder Module
Orchestrates slide creation using native PPTX layouts.

This module provides high-level slide building functionality that:
- Selects appropriate layouts based on content type
- Creates slides using native PowerPoint layouts
- Fills content using the PlaceholderFiller
- Adds background images when needed
- Handles icons, page numbers, and branding
"""

import logging
from pathlib import Path
from typing import Dict, List, Optional, Any, Tuple
from io import BytesIO

from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..models.template_manifest import (
    TemplateManifest,
    LayoutDefinition,
    PlaceholderSlot,
)
from ..models.presentation import PresentationData, SlideContent, BulletPoint, TableData
from .placeholder_filler import PlaceholderFiller, fill_slide_content
from .layout_mapper import LayoutMapper

logger = logging.getLogger("slide_builder")


# ============================================================================
# CONTENT TYPE TO LAYOUT MAPPING
# ============================================================================

DEFAULT_CONTENT_TYPE_MAPPING = {
    "title": "title_slide",
    "section": "section_header",
    "content": "title_and_content",
    "bullets": "title_and_content",
    "paragraph": "title_and_content",
    "table": "title_and_content",
    "chart": "title_and_content",
    "two_column": "two_content",
    "comparison": "comparison",
    "image": "picture_with_caption",
    "blank": "blank",
    "agenda": "title_and_content",
}


# ============================================================================
# SLIDE BUILDER CLASS
# ============================================================================

class SlideBuilder:
    """
    Builds slides using native PowerPoint layouts.
    
    This class handles the complete slide creation process:
    1. Load the template PPTX
    2. Select appropriate layout for content type
    3. Create slide from layout
    4. Fill placeholders with content
    5. Add optional elements (background, icons, page numbers)
    
    Usage:
        builder = SlideBuilder(template_path, manifest, language="en")
        prs = builder.build_presentation(presentation_data)
        prs.save("output.pptx")
    """
    
    def __init__(
        self,
        template_path: str,
        manifest: Optional[TemplateManifest] = None,
        language: str = "en",
        lang_config: Optional[Dict[str, Any]] = None
    ):
        """
        Initialize the slide builder.
        
        Args:
            template_path: Path to the PPTX template file
            manifest: Optional template manifest (auto-loaded if not provided)
            language: Target language code
            lang_config: Language-specific configuration
        """
        self.template_path = Path(template_path)
        self.manifest = manifest
        self.language = language
        self.lang_config = lang_config or {}
        self.is_rtl = self.lang_config.get("rtl", False)
        
        # Initialize presentation
        self.prs: Optional[Presentation] = None
        
        # Initialize placeholder filler
        self.filler = PlaceholderFiller(
            manifest=manifest,
            language=language,
            lang_config=lang_config
        )
        
        # Initialize layout mapper if manifest available
        self.layout_mapper: Optional[LayoutMapper] = None
        if manifest:
            self.layout_mapper = LayoutMapper(manifest)
        
        # Background images directory
        self.backgrounds_dir = self.template_path.parent / "Background"
        
        logger.info(f"SlideBuilder initialized: {self.template_path.name}")
    
    def build_presentation(
        self,
        presentation_data: PresentationData,
        add_page_numbers: bool = True,
        add_backgrounds: bool = True
    ) -> Presentation:
        """
        Build a complete presentation from presentation data.
        
        Args:
            presentation_data: PresentationData object with slides
            add_page_numbers: Whether to add page numbers
            add_backgrounds: Whether to add background images
            
        Returns:
            PowerPoint Presentation object
        """
        # Load template
        self.prs = Presentation(str(self.template_path))
        
        # Remove any existing slides from template
        while len(self.prs.slides) > 0:
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]
        
        logger.info(f"Building presentation with {len(presentation_data.slides)} slides")
        
        # Build each slide
        page_number = 0
        for idx, slide_content in enumerate(presentation_data.slides):
            slide = self.build_slide(slide_content)
            
            # Add background if enabled
            if add_backgrounds:
                self._add_background(slide, slide_content.layout_type)
            
            # Add page number if enabled (skip title slides)
            if add_page_numbers and slide_content.layout_type not in ["title", "title_slide"]:
                page_number += 1
                self._add_page_number(slide, page_number)
        
        logger.info("Presentation built successfully")
        return self.prs
    
    def build_slide(self, content: SlideContent) -> Slide:
        """
        Build a single slide from content.
        
        Args:
            content: SlideContent object
            
        Returns:
            Created Slide object
        """
        # Determine content type
        content_type = self._determine_content_type(content)
        
        # Get layout for content type
        layout = self._get_layout_for_content(content_type)
        
        # Create slide
        slide = self.prs.slides.add_slide(layout)
        
        # Fill placeholders
        self.filler.fill_slide(slide, content)
        
        logger.debug(f"Slide built: {content_type} - {content.title[:30] if content.title else 'No title'}")
        
        return slide
    
    def _determine_content_type(self, content: SlideContent) -> str:
        """Determine content type from slide content"""
        # Check explicit layout type first
        if content.layout_type:
            return content.layout_type
        
        # Infer from content
        if content.table_data:
            return "table"
        elif content.chart_data:
            return "chart"
        elif content.two_column:
            return "two_column"
        elif content.bullets:
            return "bullets"
        elif content.paragraph:
            return "paragraph"
        elif not content.title:
            return "blank"
        else:
            return "content"
    
    def _get_layout_for_content(self, content_type: str):
        """Get the appropriate layout for a content type"""
        # Try manifest mapping first
        if self.manifest and self.manifest.content_type_mapping:
            layout_key = self.manifest.content_type_mapping.get(content_type)
            if layout_key:
                layout_def = self.manifest.layouts.get(layout_key)
                if layout_def:
                    try:
                        return self.prs.slide_layouts[layout_def.index]
                    except IndexError:
                        pass
        
        # Try layout mapper
        if self.layout_mapper:
            match = self.layout_mapper.find_best_layout(content_type)
            if match:
                try:
                    return self.prs.slide_layouts[match.layout_def.index]
                except IndexError:
                    pass
        
        # Fallback to default mapping
        layout_name = DEFAULT_CONTENT_TYPE_MAPPING.get(content_type, "title_and_content")
        
        # Try to find layout by name
        for layout in self.prs.slide_layouts:
            if self._normalize_name(layout.name) == self._normalize_name(layout_name):
                return layout
        
        # Ultimate fallback - use first content layout (usually index 1)
        try:
            return self.prs.slide_layouts[1]
        except IndexError:
            return self.prs.slide_layouts[0]
    
    def _add_background(self, slide: Slide, layout_type: str) -> None:
        """Add background image to slide if available"""
        if not self.backgrounds_dir.exists():
            return
        
        # Get background path from manifest
        bg_path = None
        
        if self.manifest and self.manifest.background_images:
            bg_filename = self.manifest.background_images.get(layout_type)
            if bg_filename:
                bg_path = self.template_path.parent / bg_filename
        
        # Try common background naming
        if not bg_path or not bg_path.exists():
            possible_names = [
                f"bg_{layout_type}.png",
                f"bg_{layout_type}.jpg",
                f"bg_{layout_type}.jpeg",
                f"{layout_type}.png",
                f"{layout_type}.jpg",
            ]
            for name in possible_names:
                test_path = self.backgrounds_dir / name
                if test_path.exists():
                    bg_path = test_path
                    break
        
        if bg_path and bg_path.exists():
            self._set_slide_background(slide, str(bg_path))
            logger.debug(f"  Background added: {bg_path.name}")
    
    def _set_slide_background(self, slide: Slide, image_path: str) -> None:
        """Set slide background from image file"""
        try:
            # Get slide dimensions
            slide_width = self.prs.slide_width
            slide_height = self.prs.slide_height
            
            # Add picture at slide size, positioned at origin
            pic = slide.shapes.add_picture(
                image_path,
                Inches(0),
                Inches(0),
                width=slide_width,
                height=slide_height
            )
            
            # Send to back
            spTree = slide.shapes._spTree
            sp = pic._element
            spTree.remove(sp)
            spTree.insert(2, sp)
            
        except Exception as e:
            logger.warning(f"Could not set background: {e}")
    
    def _add_page_number(self, slide: Slide, number: int) -> None:
        """Add page number to slide"""
        # Get position from config or use default
        if self.is_rtl:
            left = Inches(0.2)
        else:
            left = Inches(self._emu_to_inches(self.prs.slide_width) - 0.7)
        
        top = Inches(self._emu_to_inches(self.prs.slide_height) - 0.7)
        width = Inches(0.5)
        height = Inches(0.5)
        
        # Create shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            left, top, width, height
        )
        
        # Style shape
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xC6, 0xC3, 0xBE)
        shape.line.fill.background()
        
        # Add text
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(number)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFC, 0xEC)
        
        # Center vertically
        tf.word_wrap = False
        tf.auto_size = None
        shape.text_frame.paragraphs[0].font.name = "Arial"
    
    @staticmethod
    def _normalize_name(name: str) -> str:
        """Normalize layout name for comparison"""
        return name.lower().replace(" ", "_").replace("-", "_")
    
    @staticmethod
    def _emu_to_inches(emu: int) -> float:
        """Convert EMUs to inches"""
        return emu / 914400


# ============================================================================
# PRESENTATION BUILDER CLASS (HIGH-LEVEL)
# ============================================================================

class PresentationBuilder:
    """
    High-level presentation builder that handles template loading
    and manifest management.
    
    Usage:
        builder = PresentationBuilder(template_id="arweqah", language="en")
        prs = builder.build(presentation_data)
        builder.save("output.pptx")
    """
    
    def __init__(
        self,
        template_dir: str,
        language: str = "en",
        lang_config: Optional[Dict[str, Any]] = None
    ):
        """
        Initialize the presentation builder.
        
        Args:
            template_dir: Path to template directory
            language: Target language code
            lang_config: Language-specific configuration
        """
        self.template_dir = Path(template_dir)
        self.language = language
        self.lang_config = lang_config or {}
        
        # Load template files
        self.template_pptx = self.template_dir / "template.pptx"
        if not self.template_pptx.exists():
            # Try to find any PPTX file
            pptx_files = list(self.template_dir.glob("*.pptx"))
            if pptx_files:
                self.template_pptx = pptx_files[0]
            else:
                raise FileNotFoundError(f"No template.pptx found in {template_dir}")
        
        # Load manifest
        self.manifest = self._load_manifest()
        
        # Create slide builder
        self.slide_builder = SlideBuilder(
            template_path=str(self.template_pptx),
            manifest=self.manifest,
            language=language,
            lang_config=lang_config
        )
        
        self.prs: Optional[Presentation] = None
        
        logger.info(f"PresentationBuilder initialized: {self.template_dir.name}")
    
    def _load_manifest(self) -> Optional[TemplateManifest]:
        """Load template manifest"""
        manifest_path = self.template_dir / "manifest.json"
        
        if manifest_path.exists():
            try:
                import json
                with open(manifest_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                return TemplateManifest(**data)
            except Exception as e:
                logger.warning(f"Could not load manifest: {e}")
        
        return None
    
    def build(
        self,
        presentation_data: PresentationData,
        add_page_numbers: bool = True,
        add_backgrounds: bool = True
    ) -> Presentation:
        """
        Build presentation from data.
        
        Args:
            presentation_data: PresentationData object
            add_page_numbers: Whether to add page numbers
            add_backgrounds: Whether to add background images
            
        Returns:
            Built Presentation object
        """
        self.prs = self.slide_builder.build_presentation(
            presentation_data,
            add_page_numbers=add_page_numbers,
            add_backgrounds=add_backgrounds
        )
        return self.prs
    
    def save(self, output_path: str) -> str:
        """Save the presentation to file"""
        if not self.prs:
            raise RuntimeError("No presentation to save. Call build() first.")
        
        self.prs.save(output_path)
        logger.info(f"Presentation saved: {output_path}")
        return output_path
    
    def save_to_bytes(self) -> bytes:
        """Save the presentation to bytes"""
        if not self.prs:
            raise RuntimeError("No presentation to save. Call build() first.")
        
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.getvalue()


# ============================================================================
# CONVENIENCE FUNCTIONS
# ============================================================================

def build_presentation(
    template_dir: str,
    presentation_data: PresentationData,
    language: str = "en",
    output_path: Optional[str] = None
) -> Presentation:
    """
    Convenience function to build a presentation.
    
    Args:
        template_dir: Path to template directory
        presentation_data: PresentationData object
        language: Target language
        output_path: Optional output file path
        
    Returns:
        Built Presentation object
    """
    builder = PresentationBuilder(template_dir, language=language)
    prs = builder.build(presentation_data)
    
    if output_path:
        builder.save(output_path)
    
    return prs
