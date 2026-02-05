"""
PPTX Generator Module - Native Template Mode
Generates PowerPoint presentations using native PPTX templates.

This generator creates slides that match the sample template layout by:
1. Using the blank layout as base
2. Adding background images
3. Creating text boxes at exact positions from config
4. Adding separators, icons, and page numbers
"""

import hashlib
import logging
import re
import json
from datetime import datetime
from pathlib import Path
from typing import Tuple, Optional, Dict, Any, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

from ..config import settings
from ..models.presentation import PresentationData, SlideContent, TableData, BulletPoint
from ..models.template_manifest import TemplateManifest
from .chart_service import ChartService
from .icon_service import IconService
from ..utils.content_validator import validate_presentation

logger = logging.getLogger("pptx_generator")


# Icon keyword mappings for intelligent icon selection
ICON_KEYWORDS = {
    # Business & Management
    'strategy': ['strategy', 'strategic', 'plan', 'planning', 'roadmap'],
    'team': ['team', 'staff', 'people', 'workforce', 'employees', 'personnel', 'group'],
    'leadership': ['leader', 'leadership', 'management', 'executive', 'director'],
    'goals': ['goal', 'goals', 'objective', 'objectives', 'target', 'targets', 'kpi'],
    'growth': ['growth', 'expand', 'expansion', 'scale', 'scaling', 'increase'],
    
    # Technical
    'technology': ['technology', 'tech', 'digital', 'software', 'system', 'systems'],
    'data': ['data', 'analytics', 'analysis', 'statistics', 'metrics', 'insights'],
    'security': ['security', 'secure', 'protection', 'safety', 'compliance'],
    'process': ['process', 'workflow', 'procedure', 'operation', 'operations'],
    'infrastructure': ['infrastructure', 'network', 'architecture', 'platform'],
    
    # Communication
    'communication': ['communication', 'communicate', 'messaging', 'contact'],
    'presentation': ['presentation', 'present', 'overview', 'summary', 'introduction'],
    'report': ['report', 'reporting', 'documentation', 'document'],
    
    # Financial
    'finance': ['finance', 'financial', 'budget', 'cost', 'investment', 'revenue'],
    'money': ['money', 'payment', 'pricing', 'fee', 'profit', 'earnings'],
    
    # Project
    'project': ['project', 'initiative', 'program', 'delivery'],
    'timeline': ['timeline', 'schedule', 'milestone', 'deadline', 'phase'],
    'quality': ['quality', 'excellence', 'standard', 'standards', 'benchmark'],
    
    # Service
    'service': ['service', 'services', 'support', 'assistance', 'help'],
    'customer': ['customer', 'client', 'user', 'guest', 'visitor', 'pilgrim'],
    'experience': ['experience', 'journey', 'satisfaction', 'feedback'],
    
    # Specific domains
    'training': ['training', 'education', 'learning', 'development', 'skill'],
    'health': ['health', 'medical', 'wellness', 'care', 'safety'],
    'environment': ['environment', 'sustainability', 'green', 'eco'],
    'location': ['location', 'place', 'site', 'venue', 'facility', 'facilities'],
    'transport': ['transport', 'transportation', 'logistics', 'travel', 'mobility'],
    'food': ['food', 'catering', 'meal', 'dining', 'nutrition'],
    'accommodation': ['accommodation', 'housing', 'hotel', 'lodging', 'stay'],
    
    # Generic
    'solution': ['solution', 'solutions', 'approach', 'method', 'methodology'],
    'benefit': ['benefit', 'benefits', 'advantage', 'value', 'impact'],
    'challenge': ['challenge', 'challenges', 'issue', 'problem', 'risk'],
    'success': ['success', 'achievement', 'result', 'results', 'outcome'],
    'innovation': ['innovation', 'innovative', 'new', 'modern', 'advanced'],
}


class PptxGenerator:
    """
    Native PPTX Generator that creates slides matching the sample template layout.
    
    Uses element_positions from config to place content at exact positions.
    """

    def __init__(self, template_id: str, language: str = "en"):
        """Initialize the PPTX generator."""
        self.template_id = template_id
        self.target_language = self._normalize_language_code(language)
        
        # Template paths
        self.template_dir = Path(settings.TEMPLATES_DIR) / template_id
        logger.info(f"Initializing generator: {self.template_dir}")
        
        if not self.template_dir.exists():
            raise FileNotFoundError(f"Template directory not found: {self.template_dir}")
        
        self.backgrounds_dir = self.template_dir / "Background"
        if not self.backgrounds_dir.exists():
            self.backgrounds_dir = self.template_dir / "backgrounds"
        
        # Load configurations
        self.config = self._load_json("config.json")
        self.constraints = self._load_json("constraints.json") if (self.template_dir / "constraints.json").exists() else {}
        self.theme = self._load_json("theme.json") if (self.template_dir / "theme.json").exists() else self._get_default_theme()
        
        # Load manifest if exists
        self.manifest: Optional[TemplateManifest] = None
        manifest_path = self.template_dir / "manifest.json"
        if manifest_path.exists():
            try:
                with open(manifest_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                self.manifest = TemplateManifest(**data)
                logger.info(f"  Manifest loaded: {len(self.manifest.layouts)} layouts")
            except Exception as e:
                logger.warning(f"  Could not load manifest: {e}")
        
        # Initialize services
        try:
            self.icon_service = IconService(template_id=template_id)
        except Exception as e:
            logger.warning(f"  IconService failed: {e}")
            self.icon_service = None
        
        try:
            self.chart_service = ChartService(template_id=template_id)
        except Exception as e:
            logger.warning(f"  ChartService failed: {e}")
            self.chart_service = None
        
        # Runtime state
        self.prs: Optional[Presentation] = None
        self.lang_config: Dict[str, Any] = {}
        
        # Load element positions (prefer config, fallback to manifest)
        self.element_positions: Dict[str, Any] = self.config.get('element_positions', {})
        if not self.element_positions and self.manifest and self.manifest.element_positions:
            self.element_positions = self.manifest.element_positions
        
        # Load fonts config (prefer config, fallback to manifest)
        self.fonts_config: Dict[str, Any] = self.config.get('fonts', {})
        if not self.fonts_config and self.manifest and self.manifest.fonts:
            self.fonts_config = self.manifest.fonts.model_dump() if hasattr(self.manifest.fonts, 'model_dump') else {}
        
        # Load icons config (prefer config, fallback to manifest)
        self.icons_config: Dict[str, Any] = self.config.get('icons', {})
        if not self.icons_config and self.manifest and self.manifest.icons:
            self.icons_config = {
                'default_title': self.manifest.icons.default_title,
                'default_section': self.manifest.icons.default_section,
                'agenda_items': self.manifest.icons.agenda_items,
                'box_icons': self.manifest.icons.box_icons
            }
        
        # Load colors config (prefer config, fallback to manifest)
        self.colors_config: Dict[str, Any] = self.config.get('colors', {})
        if not self.colors_config and self.manifest and self.manifest.colors:
            self.colors_config = self.manifest.colors.model_dump() if hasattr(self.manifest.colors, 'model_dump') else {}
        
        # Build available icons list (section falls back to title icons if no section icons)
        self.available_title_icons = self._get_available_icons('title')
        self.available_section_icons = self._get_available_icons('section')
        if not self.available_section_icons and self.available_title_icons:
            self.available_section_icons = list(self.available_title_icons)
        self.icon_index = 0  # Fallback counter for cycling when no keyword match
        # Load icon keyword rules (new format: rules with priority/keywords/icons) or legacy category_to_icon
        self.icon_keyword_rules: List[Dict] = []  # [{priority, keywords, icons: {title?, section?, alt?}}, ...]
        self.category_to_icon: Dict[str, str] = {}
        icon_kw_path = self.template_dir / "icon_keywords.json"
        if icon_kw_path.exists():
            try:
                kw_data = self._load_json("icon_keywords.json")
                rules = kw_data.get("rules") or []
                if rules:
                    # Sort by priority desc (higher wins); Python sort is stable so order preserved for ties
                    self.icon_keyword_rules = sorted(rules, key=lambda r: -int(r.get("priority", 0)))
                    # Normalize: precompute keyword list and icons dict for fast matching
                    for r in self.icon_keyword_rules:
                        r["_keywords"] = [str(k).strip().lower() for k in (r.get("keywords") or []) if k]
                        r["_icons"] = r.get("icons") or {}
                    logger.info(f"  Icon keyword rules: {len(self.icon_keyword_rules)} rules (priority-based)")
                else:
                    raw = kw_data.get("category_to_icon") or {}
                    for cat, path in raw.items():
                        if path and (self.template_dir / path).exists():
                            self.category_to_icon[cat] = path
                    if self.category_to_icon:
                        logger.info(f"  Icon keyword mapping: {len(self.category_to_icon)} categories (legacy)")
            except Exception as e:
                logger.debug(f"Could not load icon_keywords.json: {e}")
        
        logger.info(f"Generator initialized: {template_id}, lang={self.target_language}")
        logger.info(f"  Available icons: {len(self.available_title_icons)} title, {len(self.available_section_icons)} section")
    
    # ========================================================================
    # CONFIGURATION LOADING
    # ========================================================================
    
    def _load_json(self, filename: str) -> Dict:
        """Load JSON configuration file"""
        json_path = self.template_dir / filename
        if not json_path.exists():
            return {}
        with open(json_path, 'r', encoding='utf-8') as f:
            return json.load(f)
    
    def _get_default_theme(self) -> Dict:
        """Get default theme configuration"""
        return {
            "colors": {
                "primary": "#01415C",
                "secondary": "#84BA93",
                "text_primary": "#0D2026",
                "text_inverse": "#FFFCEC"
            },
            "typography": {
                "font_families": {"primary": "Calibri Light", "body": "Calibri"},
                "font_sizes": {"title": 44, "heading_1": 32, "body": 18}
            }
        }
    
    # ========================================================================
    # LANGUAGE CONFIGURATION
    # ========================================================================
    
    def _normalize_language_code(self, language: str) -> str:
        """Normalize language to 2-letter code"""
        if not language:
            return 'en'
        lang_lower = language.lower().strip()
        language_map = {'arabic': 'ar', 'ar': 'ar', 'english': 'en', 'en': 'en'}
        return language_map.get(lang_lower, 'en')
    
    def _configure_language(self, presentation_data: PresentationData) -> None:
        """Configure language settings"""
        lang_settings = self.config.get('language_settings', {})
        
        if presentation_data and hasattr(presentation_data, 'language') and presentation_data.language:
            self.target_language = self._normalize_language_code(presentation_data.language)
        
        self.lang_config = lang_settings.get(self.target_language, {
            'rtl': False,
            'default_font': 'Calibri',
            'heading_font': 'Calibri Light',
            'alignment': 'left'
        })
        
        logger.info(f"Language: {self.target_language}, RTL={self.lang_config.get('rtl', False)}")
    
    # ========================================================================
    # HELPER METHODS
    # ========================================================================
    
    def _get_alignment(self, force_center: bool = False) -> PP_ALIGN:
        """Get text alignment based on language"""
        if force_center:
            return PP_ALIGN.CENTER
        alignment = self.lang_config.get('alignment', 'left')
        if alignment == 'right':
            return PP_ALIGN.RIGHT
        elif alignment == 'center':
            return PP_ALIGN.CENTER
        return PP_ALIGN.LEFT
    
    # ========================================================================
    # ICON SELECTION METHODS
    # ========================================================================
    
    def _get_available_icons(self, icon_type: str = 'title') -> List[str]:
        """Get list of available icons from template directory"""
        icons_dir = self.template_dir / "Icons"
        if not icons_dir.exists():
            return []
        
        prefix = f"icon_{icon_type}_"
        icons = []
        for f in icons_dir.iterdir():
            if f.is_file() and f.name.startswith(prefix):
                icons.append(f"Icons/{f.name}")
        
        return sorted(icons)
    
    def _normalize_heading_for_icon_match(self, text: str) -> str:
        """Normalize heading for icon keyword match: lowercase, trim, collapse spaces."""
        if not text:
            return ""
        return re.sub(r"\s+", " ", text.strip().lower())

    def _pick_icon_path_from_rule(self, rule: Dict, icon_type: str) -> Optional[str]:
        """From a rule's icons dict, pick the best path for title vs section. Returns path if file exists."""
        icons = rule.get("_icons") or rule.get("icons") or {}
        candidates = []
        if icon_type == "section":
            if icons.get("section"):
                candidates.append(icons["section"])
            if icons.get("title"):
                candidates.append(icons["title"])
        else:
            if icons.get("title"):
                candidates.append(icons["title"])
        alt = icons.get("alt") or []
        if isinstance(alt, list):
            candidates.extend(alt)
        else:
            if alt:
                candidates.append(alt)
        for path in candidates:
            if path and (self.template_dir / path).exists():
                return path
        return None

    def _select_icon_for_content(self, title: str, icon_type: str = 'title') -> Optional[str]:
        """
        Select an appropriate icon based on slide title content.
        Uses template icon_keywords.json: if "rules" format, matches by keywords (substring) and
        priority (higher wins), then picks icons.title or icons.section; falls back to legacy
        category_to_icon or config default or cycling when no match.
        """
        icons = self.available_title_icons if icon_type == 'title' else self.available_section_icons

        def _default_icon() -> Optional[str]:
            cand = self.icons_config.get('default_section' if icon_type == 'section' else 'default_title')
            if cand and (self.template_dir / cand).exists():
                return cand
            if icons:
                return icons[0]
            other = self.available_title_icons or self.available_section_icons
            return other[0] if other else None

        if not icons:
            return _default_icon()

        title_norm = self._normalize_heading_for_icon_match(title or "")

        # New format: rules with priority + keywords + icons.title / icons.section / icons.alt
        if self.icon_keyword_rules:
            for rule in self.icon_keyword_rules:
                keywords = rule.get("_keywords") or []
                for kw in keywords:
                    if kw and kw in title_norm:
                        path = self._pick_icon_path_from_rule(rule, icon_type)
                        if path:
                            return path
                        break
            # No rule matched; fall back to default or cycle
            if not icons:
                return _default_icon()
            self.icon_index += 1
            return icons[self.icon_index % len(icons)]

        # Legacy format: category_to_icon + ICON_KEYWORDS
        best_match = None
        best_score = 0
        for category, keywords in ICON_KEYWORDS.items():
            for keyword in keywords:
                if keyword in title_norm:
                    score = len(keyword)
                    if score > best_score:
                        best_score = score
                        best_match = category
        if best_match and best_match in self.category_to_icon:
            mapped = self.category_to_icon[best_match]
            if (self.template_dir / mapped).exists():
                return mapped
        if best_match and icons:
            return icons[hash(best_match) % len(icons)]
        self.icon_index += 1
        return icons[self.icon_index % len(icons)] if icons else _default_icon()

    def _ensure_header_icon(self, icon_path: Optional[str], icon_type: str = 'title') -> Optional[str]:
        """Ensure we have an icon for a slide header; every slide with a header must show an icon."""
        if icon_path:
            return icon_path
        icons = self.available_title_icons if icon_type == 'title' else self.available_section_icons
        default = self.icons_config.get('default_section' if icon_type == 'section' else 'default_title')
        if default and (self.template_dir / default).exists():
            return default
        if icons:
            return icons[0]
        return (self.available_title_icons or self.available_section_icons)[0] if (self.available_title_icons or self.available_section_icons) else None
    
    def _get_font(self, font_type: str = 'body') -> str:
        """Get font based on type and language"""
        if font_type in ['title', 'heading']:
            return self.lang_config.get('heading_font', 'Calibri Light')
        return self.lang_config.get('default_font', 'Calibri')
    
    def _get_slide_tone(self, content_type: Optional[str]) -> str:
        """Return 'light' or 'dark' for the slide background. Used to pick contrasting text/icon colors."""
        if not content_type:
            return 'light'
        slide_color = self.config.get('slide_color') or self.constraints.get('slide_color') or {}
        tone = slide_color.get(content_type, slide_color.get('content', 'light'))
        return (tone or 'light').lower()

    def _get_text_color_for_slide(self, content_type: Optional[str]) -> str:
        """When slide_color is light → use dark font; when dark → use light font. Returns hex (no #)."""
        text_colors = self.colors_config.get('text', {})
        dark_hex = (text_colors.get('dark') or '0D2026').lstrip('#')
        light_hex = (text_colors.get('light') or 'FFFCEC').lstrip('#')
        tone = self._get_slide_tone(content_type)
        return dark_hex if tone == 'light' else light_hex

    def _get_separator_color_for_slide(self, content_type: Optional[str]) -> str:
        """Separator/line color adapted to slide: dark line on light slides, light line on dark slides."""
        return self._get_text_color_for_slide(content_type)

    def _get_font_config(self, slide_type: str, element: str, content_type: Optional[str] = None) -> Dict:
        """Get font configuration. If content_type is set, text color is derived from slide_color (light slide → dark font, dark slide → light font)."""
        font_key = 'name_ar' if self.target_language == 'ar' else 'name_en'
        slide_fonts = self.fonts_config.get(slide_type, {})
        element_fonts = slide_fonts.get(element, {})
        default_color = self.colors_config.get('text', {}).get('dark', '0D2026')
        if isinstance(default_color, str):
            default_color = default_color.lstrip('#')
        if content_type:
            color = self._get_text_color_for_slide(content_type)
        else:
            color = element_fonts.get('color', default_color)
            if isinstance(color, str):
                color = color.lstrip('#')
        return {
            'name': element_fonts.get(font_key, element_fonts.get('name', self._get_font('body'))),
            'size': element_fonts.get('size', 18),
            'bold': element_fonts.get('bold', False),
            'color': color
        }

    def _get_color(self, color_path: str, default: str = '0D2026') -> str:
        """Get color from config by path (e.g., 'text.dark', 'elements.separator_line')"""
        parts = color_path.split('.')
        value = self.colors_config
        
        for part in parts:
            if isinstance(value, dict) and part in value:
                value = value[part]
            else:
                return default
        
        return value if isinstance(value, str) else default
    
    def _get_tinted_icon_path(self, icon_path: str, tint_hex: str) -> Optional[str]:
        """
        Treat PNG as alpha mask: replace RGB with tint color, preserve alpha.
        Dark slide → light tint (#FFFCEC); light slide → dark tint (#0D2026).
        Returns path to cached tinted PNG, or None if tinting fails (caller can use original).
        """
        if not icon_path or not tint_hex:
            return None
        full_path = self.template_dir / icon_path
        if not full_path.exists():
            return None
        hex_clean = tint_hex.lstrip("#").upper()
        if len(hex_clean) != 6:
            return None
        try:
            r, g, b = int(hex_clean[0:2], 16), int(hex_clean[2:4], 16), int(hex_clean[4:6], 16)
        except ValueError:
            return None
        cache_key = hashlib.sha256(f"{icon_path}:{hex_clean}".encode()).hexdigest()[:16]
        cache_dir = self.template_dir / ".icon_tint_cache"
        cache_dir.mkdir(parents=True, exist_ok=True)
        out_path = cache_dir / f"{cache_key}.png"
        if out_path.exists():
            return str(out_path)
        try:
            from PIL import Image
            img = Image.open(full_path).convert("RGBA")
            w, h = img.size
            # Replace RGB with tint, keep original alpha (alpha-mask behavior)
            r_band = Image.new("L", (w, h), r)
            g_band = Image.new("L", (w, h), g)
            b_band = Image.new("L", (w, h), b)
            _, _, _, a_band = img.split()
            tinted = Image.merge("RGBA", (r_band, g_band, b_band, a_band))
            tinted.save(str(out_path), "PNG")
            return str(out_path)
        except Exception as e:
            logger.debug(f"Icon tint failed {icon_path}: {e}")
            return None

    def _add_icon(
        self,
        slide,
        icon_path: str,
        pos: Dict,
        content_type: Optional[str] = None,
    ) -> None:
        """
        Add an icon to the slide. If content_type is set, the icon is tinted to the slide's
        foreground color (light slide → dark icon, dark slide → light icon) so it never
        conflicts with the theme. Every slide header should call this with an icon.
        """
        if not icon_path:
            return
        full_path = self.template_dir / icon_path
        if not full_path.exists():
            logger.debug(f"Icon not found: {full_path}")
            return
        path_to_add = str(full_path)
        if content_type:
            tint_hex = self._get_text_color_for_slide(content_type)
            tinted = self._get_tinted_icon_path(icon_path, tint_hex)
            if tinted:
                path_to_add = tinted
        try:
            slide.shapes.add_picture(
                path_to_add,
                Inches(pos.get('x', 0)),
                Inches(pos.get('y', 0)),
                Inches(pos.get('width', 0.5)),
                Inches(pos.get('height', 0.5))
            )
        except Exception as e:
            logger.debug(f"Icon error: {e}")
    
    def _hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex color to RGB tuple"""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    
    def _scrub_title(self, title: str) -> str:
        """Clean title text"""
        if not title:
            return ""
        t = title.strip()
        t = re.sub(r"\(\s*continued\s*\)", "", t, flags=re.IGNORECASE)
        return re.sub(r"\s{2,}", " ", t).strip()
    
    def _set_rtl_paragraph(self, paragraph) -> None:
        """Set RTL on a paragraph"""
        if not self.lang_config.get('rtl', False):
            return
        try:
            pPr = paragraph._element.get_or_add_pPr()
            pPr.set('rtl', '1')
        except:
            pass
    
    # ========================================================================
    # BACKGROUND AND PAGE NUMBERS
    # ========================================================================
    
    def _add_background(self, slide, content_type: str) -> None:
        """Add background image to slide"""
        bg_images = self.config.get('background_images', {})
        bg_path_str = bg_images.get(content_type)
        
        if not bg_path_str:
            # Try fallback
            bg_path_str = bg_images.get('content')
        
        if not bg_path_str:
            return
        
        bg_path = self.template_dir / bg_path_str
        if not bg_path.exists():
            logger.debug(f"Background not found: {bg_path}")
            return
        
        try:
            picture = slide.shapes.add_picture(
                str(bg_path),
                Inches(0), Inches(0),
                width=self.prs.slide_width,
                height=self.prs.slide_height
            )
            # Send to back
            slide.shapes._spTree.remove(picture._element)
            slide.shapes._spTree.insert(2, picture._element)
        except Exception as e:
            logger.warning(f"Background error: {e}")
    
    def _add_page_number(self, slide, page_num: int, content_type: Optional[str] = None) -> None:
        """Add page number diamond in the same position as the Agenda slide's embedded rhombus (bottom-right), so all slides match. Text color adapts to slide_color. Double-digit numbers use a wider shape so text stays horizontal."""
        try:
            page_config = self.config.get('page_numbering', {})
            if not page_config.get('enabled', True):
                return

            # Use the same position as the Agenda slide's embedded rhombus (bottom-right) for all slides
            position = page_config.get('position_unified') or page_config.get('position_ar', page_config.get('position_en', {}))
            offset_x = float(position.get('offset_x', 12.73))
            offset_y = float(position.get('offset_y', 6.8))

            shape_config = page_config.get('shape', {})
            base_w = float(shape_config.get('width', 0.4))
            base_h = float(shape_config.get('height', 0.4))
            # Slightly larger shape for double-digit page numbers so digits stay horizontal (not stacked)
            two_digits = page_num >= 10
            shape_w = base_w * 1.25 if two_digits else base_w
            shape_h = base_h

            diamond = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(offset_x), Inches(offset_y),
                Inches(shape_w), Inches(shape_h)
            )

            bg_color = self._get_color('elements.page_number_bg', shape_config.get('fill_color', 'C6C3BE'))
            fill_color = self._hex_to_rgb(bg_color)
            diamond.fill.solid()
            diamond.fill.fore_color.rgb = RGBColor(*fill_color)
            diamond.line.fill.background()

            tf = diamond.text_frame
            tf.clear()
            tf.word_wrap = False
            p = tf.paragraphs[0]
            p.text = str(page_num)
            p.alignment = PP_ALIGN.CENTER

            pn_font = self.fonts_config.get('page_number', {})
            font_config = page_config.get('font', {})

            p.font.size = Pt(pn_font.get('size', font_config.get('size', 14)))
            p.font.name = pn_font.get('name', font_config.get('name', 'Cairo'))
            p.font.bold = pn_font.get('bold', font_config.get('bold', False))

            if content_type:
                text_color_hex = self._get_text_color_for_slide(content_type)
            else:
                text_color_hex = self._get_color('elements.page_number_text', pn_font.get('color', 'FFFCEC'))
            text_color = self._hex_to_rgb(text_color_hex)
            p.font.color.rgb = RGBColor(*text_color)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        except Exception as e:
            logger.warning(f"Page number error: {e}")
    
    def _add_separator_line(self, slide, pos: Dict, content_type: Optional[str] = None) -> None:
        """Add a separator line. Color adapts to slide_color (light slide → dark line, dark slide → light line)."""
        try:
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(pos.get('x', 1.7)),
                Inches(pos.get('y', 1.55)),
                Inches(pos.get('width', 9.6)),
                Inches(pos.get('height', 0.02))
            )
            line.fill.solid()
            if content_type:
                separator_color = self._get_separator_color_for_slide(content_type)
            else:
                separator_color = self.colors_config.get('elements', {}).get('separator_line', '01415C')
            if isinstance(separator_color, str):
                separator_color = separator_color.lstrip('#')
            line.fill.fore_color.rgb = RGBColor(*self._hex_to_rgb(separator_color))
            line.line.fill.background()
        except Exception as e:
            logger.debug(f"Separator error: {e}")
    
    # ========================================================================
    # TEXT BOX CREATION
    # ========================================================================
    
    def _add_text_box(self, slide, text: str, pos: Dict, 
                      font_name: str = None, font_size: int = None,
                      bold: bool = False, color: str = None,
                      alignment: PP_ALIGN = None) -> Any:
        """Add a text box at specified position"""
        if not text:
            return None
        
        try:
            shape = slide.shapes.add_textbox(
                Inches(pos.get('x', 0)),
                Inches(pos.get('y', 0)),
                Inches(pos.get('width', 10)),
                Inches(pos.get('height', 1))
            )
            
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            
            # Apply formatting
            p.font.name = font_name or self._get_font('body')
            if font_size:
                p.font.size = Pt(font_size)
            p.font.bold = bold
            
            if color:
                p.font.color.rgb = RGBColor(*self._hex_to_rgb(color))
            
            if alignment:
                p.alignment = alignment
            else:
                p.alignment = self._get_alignment()
            
            self._set_rtl_paragraph(p)
            
            return shape
        except Exception as e:
            logger.warning(f"Text box error: {e}")
            return None
    
    def _add_bullets_textbox(self, slide, bullets: List[BulletPoint], pos: Dict, content_type: Optional[str] = None) -> Any:
        """Add a text box with bullet points. Colors adapt to slide_color when content_type is set."""
        if not bullets:
            return None
        
        try:
            shape = slide.shapes.add_textbox(
                Inches(pos.get('x', 1)),
                Inches(pos.get('y', 1.8)),
                Inches(pos.get('width', 11)),
                Inches(pos.get('height', 5))
            )
            
            tf = shape.text_frame
            tf.word_wrap = True
            
            bullet_font = self._get_font_config('content', 'bullet', content_type=content_type)
            sub_bullet_font = self._get_font_config('content', 'sub_bullet', content_type=content_type)
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                # Clean bullet text
                text = (bullet.text or "").replace("●", "").replace("**", "").strip()
                p.text = f"• {text}"
                p.font.name = bullet_font['name']
                p.font.size = Pt(bullet_font['size'])
                p.font.bold = bullet_font['bold']
                if bullet_font['color']:
                    p.font.color.rgb = RGBColor(*self._hex_to_rgb(bullet_font['color']))
                p.level = 0
                p.alignment = self._get_alignment()
                self._set_rtl_paragraph(p)
                
                # Sub-bullets
                if hasattr(bullet, 'sub_bullets') and bullet.sub_bullets:
                    for sub in bullet.sub_bullets:
                        sub_p = tf.add_paragraph()
                        sub_text = sub.text if hasattr(sub, 'text') else str(sub)
                        sub_p.text = f"   ○ {sub_text.strip()}"
                        sub_p.font.name = sub_bullet_font['name']
                        sub_p.font.size = Pt(sub_bullet_font['size'])
                        sub_p.font.bold = sub_bullet_font['bold']
                        if sub_bullet_font['color']:
                            sub_p.font.color.rgb = RGBColor(*self._hex_to_rgb(sub_bullet_font['color']))
                        sub_p.level = 1
                        sub_p.alignment = self._get_alignment()
                        self._set_rtl_paragraph(sub_p)
            
            return shape
        except Exception as e:
            logger.warning(f"Bullets error: {e}")
            return None
    
    # ========================================================================
    # SLIDE TYPE CREATION
    # ========================================================================
    
    def _get_blank_layout(self):
        """Get the blank layout"""
        layout_idx = self.config.get('layout_mapping', {}).get('blank', 6)
        try:
            return self.prs.slide_layouts[layout_idx]
        except IndexError:
            return self.prs.slide_layouts[0]
    
    def _create_title_slide(self, presentation_data: PresentationData) -> None:
        """Create title slide matching sample layout"""
        slide = self.prs.slides.add_slide(self._get_blank_layout())
        
        # Add background
        self._add_background(slide, 'title_slide')
        content_type = 'title_slide'
        # Get positions
        positions = self.element_positions.get('title_slide', {})
        title_pos = positions.get('title', {'x': 1.5, 'y': 2.8, 'width': 10.33, 'height': 1.5})
        subtitle_pos = positions.get('subtitle', {'x': 2.0, 'y': 4.5, 'width': 9.33, 'height': 1.0})
        # Get font config (colors adapted to slide_color: dark slide → light text)
        title_font = self._get_font_config('title_slide', 'title', content_type=content_type)
        subtitle_font = self._get_font_config('title_slide', 'subtitle', content_type=content_type)
        
        # Add title
        title = presentation_data.title or "Untitled Presentation"
        self._add_text_box(
            slide, self._scrub_title(title), title_pos,
            font_name=title_font['name'],
            font_size=title_font['size'],
            bold=title_font['bold'],
            color=title_font['color'],
            alignment=PP_ALIGN.CENTER
        )
        
        # Add subtitle
        subtitle_parts = []
        if presentation_data.subtitle:
            subtitle_parts.append(presentation_data.subtitle)
        if presentation_data.author:
            subtitle_parts.append(presentation_data.author)
        
        if subtitle_parts:
            self._add_text_box(
                slide, "\n".join(subtitle_parts), subtitle_pos,
                font_name=subtitle_font['name'],
                font_size=subtitle_font['size'],
                bold=subtitle_font['bold'],
                color=subtitle_font['color'],
                alignment=PP_ALIGN.CENTER
            )
        
        logger.info("Title slide created")
    
    def _create_section_slide(self, slide_data: SlideContent, page_num: int = None) -> None:
        """Create section header slide matching sample layout"""
        slide = self.prs.slides.add_slide(self._get_blank_layout())
        
        # Add background
        self._add_background(slide, 'section')
        content_type = 'section'
        # Get positions
        positions = self.element_positions.get('section_header', {})
        icon_pos = positions.get('icon', {'x': 6.07, 'y': 2.2, 'width': 1.2, 'height': 1.2})
        title_pos = positions.get('title', {'x': 1.0, 'y': 3.8, 'width': 11.33, 'height': 1.8})
        # Get font config (colors adapted to slide_color: dark section → light text)
        title_font = self._get_font_config('section_header', 'title', content_type=content_type)
        # Every section slide has an icon; tint to slide foreground (dark slide → light icon)
        section_icon = self._ensure_header_icon(
            self._select_icon_for_content(slide_data.title, icon_type='section'),
            icon_type='section'
        )
        if section_icon:
            self._add_icon(slide, section_icon, icon_pos, content_type=content_type)
        # Add title
        self._add_text_box(
            slide, self._scrub_title(slide_data.title), title_pos,
            font_name=title_font['name'],
            font_size=title_font['size'],
            bold=title_font['bold'],
            color=title_font['color'],
            alignment=PP_ALIGN.CENTER
        )
        if page_num:
            self._add_page_number(slide, page_num, content_type=content_type)
    
    def _create_agenda_slide(self, slide_data: SlideContent, page_num: int = None) -> None:
        """
        Create agenda slide: left (beige) = topics with icons, right (dark) = AGENDA label centered.
        Uses bg_blank_3246436f.jpg which has beige left and dark right split.
        """
        slide = self.prs.slides.add_slide(self._get_blank_layout())
        self._add_background(slide, 'agenda')
        positions = self.element_positions.get('agenda', {})
        agenda_label_pos = positions.get('agenda_label', {'x': 7.20, 'y': 2.80, 'width': 4.90, 'height': 1.80})
        items_pos = positions.get('items', {'x': 0.70, 'y': 1.35, 'width': 5.50, 'height': 5.00})
        icon_size = positions.get('icon_size', {'width': 0.45, 'height': 0.45})
        item_height = float(positions.get('item_height', 0.85))
        # Right side (dark): AGENDA label → light text. Left side (light): items → dark text.
        agenda_label_font = self._get_font_config('agenda', 'agenda_label', content_type='agenda')
        agenda_text = self.config.get('localization', {}).get(self.target_language, {}).get('agenda_title', 'Agenda')
        if self.target_language == 'ar':
            agenda_text = self.config.get('localization', {}).get('ar', {}).get('agenda_title', agenda_text)
        
        self._add_text_box(
            slide, agenda_text, agenda_label_pos,
            font_name=agenda_label_font.get('name', agenda_label_font.get('name_en', 'Cairo')),
            font_size=agenda_label_font.get('size', 36),
            bold=agenda_label_font.get('bold', True),
            color=agenda_label_font.get('color', 'FFFCEC'),
            alignment=PP_ALIGN.CENTER
        )
        
        # Left side (beige): Agenda items with icons. Keep content above page-number zone (y >= 6.25).
        items_font = self._get_font_config('agenda', 'items', content_type='agenda_items')
        bullets = slide_data.bullets or []
        items_top = float(items_pos.get('y', 1.35))
        items_width = float(items_pos.get('width', 5.50))
        items_left = float(items_pos.get('x', 0.70))
        agenda_bottom_safe = 6.25  # Page number zone starts ~6.8; leave margin
        max_agenda_items = 6
        num_items = min(len(bullets), max_agenda_items)
        if num_items > 0:
            item_height = min(item_height, (agenda_bottom_safe - items_top - 0.2) / num_items)
        for i, bullet in enumerate(bullets):
            if i >= max_agenda_items:
                break
            text = getattr(bullet, 'text', bullet) if hasattr(bullet, 'text') else str(bullet)
            if not text:
                continue
            y = items_top + i * item_height
            icon_path = self._ensure_header_icon(self._select_icon_for_content(text, icon_type='title'), icon_type='title')
            if icon_path:
                self._add_icon(slide, icon_path, {
                    'x': items_left, 'y': y + 0.1,
                    'width': icon_size.get('width', 0.45),
                    'height': icon_size.get('height', 0.45)
                }, content_type='agenda_items')

            text_box_pos = {
                'x': items_left + 0.6,
                'y': y,
                'width': items_width - 0.6,
                'height': item_height
            }
            self._add_text_box(
                slide, text, text_box_pos,
                font_name=items_font.get('name', items_font.get('name_en', 'Open Sans')),
                font_size=items_font.get('size', 20),
                bold=items_font.get('bold', False),
                color=items_font.get('color', '0D2026'),
                alignment=self._get_alignment()
            )
        
        if page_num:
            self._add_page_number(slide, page_num, content_type='agenda')
    
    def _create_content_slide(self, slide_data: SlideContent, page_num: int = None) -> None:
        """Create content slide matching sample layout"""
        slide = self.prs.slides.add_slide(self._get_blank_layout())
        
        # Determine content type for background
        if slide_data.table_data:
            content_type = 'table'
        elif slide_data.chart_data:
            content_type = 'chart'
        else:
            content_type = 'content'
        
        # Add background
        self._add_background(slide, content_type)
        # Get positions
        positions = self.element_positions.get('content', {})
        icon_pos = positions.get('icon', {'x': 1.0, 'y': 0.65, 'width': 0.5, 'height': 0.5})
        title_pos = positions.get('title', {'x': 1.7, 'y': 0.6, 'width': 9.6, 'height': 0.8})
        separator_pos = positions.get('separator', {'x': 1.7, 'y': 1.55, 'width': 9.6, 'height': 0.02})
        body_pos = positions.get('body', {'x': 1.0, 'y': 1.8, 'width': 11.33, 'height': 5.2})
        # Get font configs (colors adapted to slide_color for this content_type)
        title_font = self._get_font_config('content', 'title', content_type=content_type)
        body_font = self._get_font_config('content', 'body', content_type=content_type)
        # Every content slide has a header icon; tint to slide foreground (light/dark by content_type)
        title_icon = self._ensure_header_icon(
            self._select_icon_for_content(slide_data.title, icon_type='title'),
            icon_type='title'
        )
        if title_icon:
            self._add_icon(slide, title_icon, icon_pos, content_type=content_type)
        # Add title
        self._add_text_box(
            slide, self._scrub_title(slide_data.title), title_pos,
            font_name=title_font['name'],
            font_size=title_font['size'],
            bold=title_font['bold'],
            color=title_font['color'],
            alignment=self._get_alignment()
        )
        # Add separator line (color adapted to slide)
        self._add_separator_line(slide, separator_pos, content_type=content_type)
        
        # Add content based on type
        content_added = False
        
        if slide_data.table_data and slide_data.table_data.rows:
            self._add_table(slide, slide_data.table_data, body_pos)
            content_added = True
        elif slide_data.chart_data and self.chart_service:
            self._add_chart(slide, slide_data, body_pos)
            content_added = True
        elif slide_data.bullets:
            self._add_bullets_textbox(slide, slide_data.bullets, body_pos, content_type=content_type)
            content_added = True
        elif slide_data.paragraph:
            self._add_text_box(
                slide, slide_data.paragraph, body_pos,
                font_name=body_font['name'],
                font_size=body_font['size'],
                bold=body_font['bold'],
                color=body_font['color']
            )
            content_added = True
        elif slide_data.content:
            self._add_text_box(
                slide, slide_data.content, body_pos,
                font_name=body_font['name'],
                font_size=body_font['size'],
                bold=body_font['bold'],
                color=body_font['color']
            )
            content_added = True
        
        # Log warning if no content was added
        if not content_added:
            logger.warning(f"Slide '{slide_data.title}' has no body content (no bullets, paragraph, table, or chart data)")
        
        if page_num:
            self._add_page_number(slide, page_num, content_type=content_type)
    
    def _add_table(self, slide, table_data: TableData, pos: Dict) -> None:
        """Add table at specified position"""
        if not table_data or not table_data.rows:
            return
        
        rows = len(table_data.rows) + (1 if table_data.headers else 0)
        cols = len(table_data.headers) if table_data.headers else len(table_data.rows[0])
        
        if rows == 0 or cols == 0:
            return
        
        try:
            table_shape = slide.shapes.add_table(
                rows, cols,
                Inches(pos.get('x', 1.0)),
                Inches(pos.get('y', 1.8)),
                Inches(pos.get('width', 11.0)),
                Inches(pos.get('height', 5.0))
            )
            table = table_shape.table
            
            # Fill headers
            row_idx = 0
            if table_data.headers:
                for col_idx, header in enumerate(table_data.headers):
                    cell = table.cell(0, col_idx)
                    cell.text = str(header)
                    for p in cell.text_frame.paragraphs:
                        p.font.bold = True
                        p.font.name = self._get_font('heading')
                        p.font.size = Pt(14)
                row_idx = 1
            
            # Fill data
            for data_row in table_data.rows:
                for col_idx, value in enumerate(data_row):
                    if col_idx < cols:
                        cell = table.cell(row_idx, col_idx)
                        cell.text = str(value)
                        for p in cell.text_frame.paragraphs:
                            p.font.name = self._get_font('body')
                            p.font.size = Pt(12)
                row_idx += 1
                
        except Exception as e:
            logger.warning(f"Table error: {e}")
    
    def _add_chart(self, slide, slide_data: SlideContent, pos: Dict) -> None:
        """Add chart at specified position"""
        if not self.chart_service or not slide_data.chart_data:
            return
        
        try:
            self.chart_service.create_chart(
                slide=slide,
                chart_data=slide_data.chart_data,
                left=float(pos.get('x', 1.5)),
                top=float(pos.get('y', 1.8)),
                width=float(pos.get('width', 10.0)),
                height=float(pos.get('height', 5.0))
            )
        except Exception as e:
            logger.warning(f"Chart error: {e}")
    
    # ========================================================================
    # MAIN GENERATION
    # ========================================================================
    
    def _determine_content_type(self, slide_data: SlideContent) -> str:
        """Determine content type from slide data"""
        layout_hint = getattr(slide_data, 'layout_hint', None) or ''
        if layout_hint and 'agenda' in layout_hint.lower():
            return 'agenda'
        
        if slide_data.layout_type:
            lt = slide_data.layout_type.lower()
            if lt in ['section', 'section_header']:
                return 'section'
            if lt == 'agenda':
                return 'agenda'
            return lt
        
        if slide_data.table_data:
            return "table"
        if slide_data.chart_data:
            return "chart"
        if slide_data.bullets:
            return "bullets"
        if slide_data.paragraph:
            return "paragraph"
        
        return "content"
    
    def generate(self, presentation_data: PresentationData) -> str:
        """Generate PowerPoint presentation."""
        logger.info("=" * 60)
        logger.info("Starting presentation generation...")
        
        # Configure language
        self._configure_language(presentation_data)
        
        logger.info(f"  Template: {self.template_id}")
        logger.info(f"  Language: {self.target_language}")
        
        # Validate slides
        presentation_data.slides = validate_presentation(presentation_data.slides)
        
        # Load template PPTX to get layouts
        template_pptx = self.template_dir / "template.pptx"
        if template_pptx.exists():
            logger.info(f"  Loading template: {template_pptx.name}")
            self.prs = Presentation(str(template_pptx))
            
            # Clear existing slides - we only want the layouts
            existing_count = len(self.prs.slides)
            if existing_count > 0:
                logger.info(f"  Clearing {existing_count} template slides...")
                for i in range(existing_count - 1, -1, -1):
                    rId = self.prs.slides._sldIdLst[i].rId
                    self.prs.part.drop_rel(rId)
                    del self.prs.slides._sldIdLst[i]
        else:
            logger.info("  Creating blank presentation")
            self.prs = Presentation()
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
        
        # Clean titles
        for slide_data in presentation_data.slides:
            if slide_data.title:
                slide_data.title = self._scrub_title(slide_data.title)
        
        # Create title slide
        try:
            self._create_title_slide(presentation_data)
        except Exception as e:
            logger.error(f"Title slide error: {e}")
            import traceback
            traceback.print_exc()
        
        # Create content slides
        page_num = 2  # Start from 2 (title is 1, but usually not numbered)
        for idx, slide_data in enumerate(presentation_data.slides):
            content_type = self._determine_content_type(slide_data)
            logger.info(f"Creating slide {idx + 2}: {content_type} - {slide_data.title[:30] if slide_data.title else 'No title'}...")
            
            try:
                if content_type in ['section', 'section_header']:
                    self._create_section_slide(slide_data, page_num=page_num)
                elif content_type == 'agenda':
                    self._create_agenda_slide(slide_data, page_num=page_num)
                else:
                    self._create_content_slide(slide_data, page_num=page_num)
                
                page_num += 1
                    
            except Exception as e:
                logger.error(f"Slide error: {e}")
                import traceback
                traceback.print_exc()
        
        # Save presentation
        output_path = self._get_output_path(presentation_data.title)
        self.prs.save(output_path)
        
        logger.info(f"Generated: {output_path}")
        logger.info(f"  Total slides: {len(self.prs.slides)}")
        logger.info("=" * 60)
        
        return output_path

    # ========================================================================
    # TITLE SLIDE
    # ========================================================================

    def _create_title_slide_dynamic(self, presentation_data: PresentationData) -> None:
        """Create title slide with language-specific positioning"""
        layout_config = self.layouts.get('title_slide', {})
        slide_layout = self._get_slide_layout(layout_config)
        slide = self.prs.slides.add_slide(slide_layout)
        self._clear_default_placeholders(slide)

        bg_config = layout_config.get('background', {})
        if isinstance(bg_config, str):
            # Look up background from config.json backgrounds section
            bg_config = self.config.get('backgrounds', {}).get(bg_config, {})
        if isinstance(bg_config, dict) and bg_config.get('type') == 'image':
            self._add_background(slide, bg_config)

        elements = layout_config.get('elements', [])
        
        # Logo
        for element in elements:
            if element.get('type') == 'image' and element.get('id') == 'logo':
                self._add_logo(slide, element)

        # Title
        title = presentation_data.title or "Untitled Presentation"
        max_title_length = self.constraints['text_constraints']['max_title_length']
        if len(title) > max_title_length:
            title = title[:max_title_length - 3] + "..."
        
        title_element = next((e for e in elements if e.get('id') == 'title'), None)
        
        if title_element:
            pos = self.get_position(title_element, 'position')
            size = title_element.get('size', {})
            style = title_element.get('style', {})

            base_font_size = self.get_font_size('title')
            if len(title) > 60:
                font_size = 36
            elif len(title) > 40:
                font_size = 40
            else:
                font_size = base_font_size

            char_per_line = 50 if font_size >= 40 else 60
            lines = max(1, len(title) / char_per_line)
            required_height = lines * 0.6
            actual_height = max(size.get('height', 1.5), required_height)
            
            textbox = slide.shapes.add_textbox(
                Inches(pos.get('left', 1.5)),
                Inches(pos.get('top', 2.8)),
                Inches(size.get('width', 10.33)),
                Inches(actual_height)
            )

            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(font_size)
            p.font.name = self.get_font('heading')
            p.font.bold = True
            
            text_color = self.get_style_value(style, 'color', '#FFFCEC')
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = self.get_line_spacing('title')

        # Subtitle
        subtitle_parts = []
        if presentation_data.subtitle and presentation_data.subtitle != "None":
            subtitle_parts.append(presentation_data.subtitle)
        if presentation_data.author and presentation_data.author != "None":
            subtitle_parts.append(presentation_data.author)
        
        subtitle = "\n".join(subtitle_parts).strip()
        
        subtitle_element = next((e for e in elements if e.get('id') == 'subtitle'), None)
        
        if subtitle_element and subtitle:
            pos = self.get_position(subtitle_element, 'position')
            size = subtitle_element.get('size', {})
            style = subtitle_element.get('style', {})

            subtitle_top = pos.get('top', 4.5)
            if lines > 1.5:
                subtitle_top += 0.3

            textbox = slide.shapes.add_textbox(
                Inches(pos.get('left', 2.0)),
                Inches(subtitle_top),
                Inches(size.get('width', 9.33)),
                Inches(size.get('height', 1.0))
            )

            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()

            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(self.get_style_value(style, 'font_size', 24))
            p.font.name = self.get_font('body')
            
            text_color = self.get_style_value(style, 'color', '#FFFCEC')
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))
            p.alignment = PP_ALIGN.CENTER

    # ========================================================================
    # AGENDA SLIDE - ENHANCED WITH RTL SUPPORT
    # ========================================================================

    def _create_agenda_slide_enhanced(self, slide, layout_config: Dict, data) -> None:
        """Enhanced agenda with RTL/LTR support and proper localization"""
        # Background
        bg_config = layout_config.get('background', {})
        if isinstance(bg_config, str):
            # Look up background from config.json backgrounds section
            bg_config = self.config.get('backgrounds', {}).get(bg_config, {})
        if isinstance(bg_config, dict) and bg_config.get('type') == 'image':
            self._add_background(slide, bg_config)
        
        # *** FIX 1: Use localized agenda title ***
        localization = self.config.get('localization', {})
        lang_strings = localization.get(self.target_language, {})
        agenda_title_text = lang_strings.get('agenda_title', 'Agenda')
        
        # Add "Agenda" text on dark side
        elements = layout_config.get('elements', [])
        agenda_title_elem = next((e for e in elements if e.get('id') == 'agenda_title'), None)
        
        if agenda_title_elem:
            pos = self.get_position(agenda_title_elem, 'position')
            size = agenda_title_elem.get('size', {})
            style = agenda_title_elem.get('style', {})
            
            agenda_textbox = slide.shapes.add_textbox(
                Inches(pos.get('left', 7.5)),
                Inches(pos.get('top', 3.0)),
                Inches(size.get('width', 4.5)),
                Inches(size.get('height', 1.5))
            )
            
            tf = agenda_textbox.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # *** FIX: Set RTL for agenda title ***
            is_rtl = self.lang_config.get('rtl', False)
            self._set_text_frame_rtl(tf, is_rtl)
            
            p = tf.paragraphs[0]
            p.text = agenda_title_text
            p.font.size = Pt(self.get_font_size('agenda_title'))
            p.font.name = self.get_font('heading')
            p.font.bold = True
            
            text_color = self.get_style_value(style, 'color', '#FFFCEC')
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))
            p.alignment = PP_ALIGN.CENTER
        
        # Logo
        logo_elem = next((e for e in elements if e.get('id') == 'logo'), None)
        if logo_elem:
            self._add_logo(slide, logo_elem)
        
        # Content bullets with icons
        content_elem = next((e for e in elements if e.get('id') == 'content'), None)
        if not content_elem:
            return
        
        bullets = self._extract_bullet_items(data, 'content')
        if not bullets:
            bullets = self._extract_bullet_items(data, 'bullets')
        
        if not bullets:
            logger.warning("⚠️  Agenda slide has no content")
            return
        
        # Get agenda constraints
        agenda_config = self.constraints.get('agenda', {})
        max_items = agenda_config.get('max_items', 5)
        bullets = bullets[:max_items]
        
        # Get positioning
        pos = self.get_position(content_elem, 'position')
        size = content_elem.get('size', {})
        style = content_elem.get('style', {})
        bullet_style = content_elem.get('bullet_style', {})
        
        # Calculate vertical centering
        num_items = len(bullets)
        item_spacing = agenda_config.get('item_spacing', 0.85)
        total_height = num_items * item_spacing
        start_y = (7.5 - total_height) / 2
        
        # Text styling
        text_color = self.get_style_value(style, 'color', '#0D2026')
        text_rgb = self.hex_to_rgb(text_color)
        font_name = self.get_font('body')
        font_size = agenda_config.get('item_font_size', 20)
        line_spacing = agenda_config.get('line_spacing', 1.8)
        
        # Icon configuration
        icon_size = bullet_style.get('icon_size', 0.45)
        is_rtl = self.lang_config.get('rtl', False)
        
        # Get icon positioning based on language
        if is_rtl:
            icon_offset = bullet_style.get('icon_offset_ar', 0.6)
            text_left = pos.get('left', 6.83)
            icon_align = 'right'
        else:
            icon_offset = bullet_style.get('icon_offset_en', 0.3)
            text_left = pos.get('left', 0.7) + icon_size + icon_offset
            icon_align = 'left'
        
        # Render each agenda item
        for idx, bullet in enumerate(bullets):
            item_y = start_y + (idx * item_spacing)
            
            # Add icon
            if self.icon_service and agenda_config.get('use_icons', True):
                icon_name = self.icon_service.auto_select_icon(bullet.text or "", "")
                try:
                    icon_data = self.icon_service.render_to_png(
                        icon_name,
                        size=int(icon_size * 96),
                        color=text_color
                    )
                    if icon_data:
                        if is_rtl:
                            # *** FIX: Icon on the RIGHT for RTL ***
                            icon_left = Inches(text_left + size.get('width', 5.8) - icon_size - 0.2)
                        else:
                            # Icon on the left for LTR
                            icon_left = Inches(pos.get('left', 0.7))
                        
                        icon_top = Inches(item_y + 0.15)
                        
                        slide.shapes.add_picture(
                            icon_data,
                            icon_left,
                            icon_top,
                            width=Inches(icon_size),
                            height=Inches(icon_size)
                        )
                except Exception as e:
                    logger.debug(f"Icon skip: {e}")
            
            # Add text
            text_width = size.get('width', 5.8) - icon_size - icon_offset
            text_box = slide.shapes.add_textbox(
                Inches(text_left),
                Inches(item_y),
                Inches(text_width),
                Inches(0.7)
            )
            
            tf = text_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Inches(0.05)
            
            # *** FIX: Set RTL for agenda items ***
            self._set_text_frame_rtl(tf, is_rtl)
            
            p = tf.paragraphs[0]
            p.text = (bullet.text or "").strip()
            p.font.size = Pt(font_size)
            p.font.name = font_name
            p.font.color.rgb = RGBColor(*text_rgb)
            p.font.bold = False
            
            # *** FIX: Use constraint-based alignment ***
            alignment_config = self.constraints.get('alignment', {})
            lang_alignment = alignment_config.get(self.target_language, {})
            agenda_alignment = lang_alignment.get('agenda', 'left')
            p.alignment = self._alignment_enum(agenda_alignment)
            p.line_spacing = line_spacing
        
        logger.info(f"✅ Agenda created: {len(bullets)} items, RTL={is_rtl}")

    # ========================================================================
    # SLIDE CREATION FROM JSON
    # ========================================================================

    def _create_slide_from_json(self, content_type: str, data, page_num: int = None) -> None:
        """Create slide with RTL/LTR support"""
        try:
            # Check for agenda
            layout_hint = self._get_layout_hint(data)
            if layout_hint and 'agenda' in layout_hint.lower():
                content_type = 'agenda'
            
            layout_config = self.get_layout_for_content(content_type, data)
            slide_layout = self._get_slide_layout(layout_config)
            slide = self.prs.slides.add_slide(slide_layout)
            self._clear_default_placeholders(slide)
            
            # Background - handle both string references and dict configs
            bg_config = layout_config.get('background', {})
            if isinstance(bg_config, str):
                # Look up background from config.json backgrounds section
                bg_config = self.config.get('backgrounds', {}).get(bg_config, {})
            if isinstance(bg_config, dict) and bg_config.get('type') == 'image':
                self._add_background(slide, bg_config)
            
            # Special: Agenda
            if content_type == 'agenda' or (layout_hint and 'agenda' in layout_hint.lower()):
                self._create_agenda_slide_enhanced(slide, layout_config, data)
                if page_num:
                    self.add_page_number(slide, page_num)
                return
            
            elements = layout_config.get('elements', [])
            
            # Special: Section headers
            if content_type == 'section':
                title_text = self._extract_text_value(data, 'title')
                
                max_title = self.constraints['text_constraints']['max_title_length']
                if len(title_text) > max_title:
                    title_text = title_text[:max_title - 3] + "..."
                    data.title = title_text
                
                layout_config = self.get_layout_for_content(content_type, data)
                
                bg_config = layout_config.get('background', {})
                if isinstance(bg_config, str):
                    # Look up background from config.json backgrounds section
                    bg_config = self.config.get('backgrounds', {}).get(bg_config, {})
                if isinstance(bg_config, dict) and bg_config.get('type') == 'image':
                    self._add_background(slide, bg_config)
                    logger.info(f"✓ Section background added: {bg_config.get('path')}")
                
                # Add centered icon for section
                if self.icon_service and title_text:
                    is_thank_you = any(
                        word in title_text.lower() 
                        for word in ['thank', 'thanks', 'شكر', 'شكراً', 'شكرا', 'thank you', 'conclusion']
                    )
                    
                    if is_thank_you:
                        icon_name = 'hand-waving'
                    else:
                        icon_name = self.icon_service.auto_select_icon(title_text, "")
                    
                    # *** FIX: Pass slide_data ***
                    self._add_centered_section_icon(slide, icon_name, title_text, layout_config, slide_data=data)
                
                # Add title text
                for element in elements:
                    if element.get('type') == 'text' and element.get('id') == 'title':
                        title_element = element
                        self._add_text_master(slide, element, data, content_type)
                
                logger.info(f"✓ Section slide complete: {title_text}")
                return
            
            # Detect content type from data
            has_chart = bool(getattr(data, 'chart_data', None) or getattr(data, 'chart', None))
            has_table = bool(getattr(data, 'table_data', None) or getattr(data, 'table', None))
            has_paragraph = bool(getattr(data, 'content', None) and len(str(getattr(data, 'content', '')).strip()) > 0)
            has_bullets = bool(getattr(data, 'bullets', None) and len(getattr(data, 'bullets', [])) > 0)
            
            if has_chart:
                logger.info(f"   → Chart slide detected")
            if has_table:
                logger.info(f"   → Table slide detected")
            if has_paragraph:
                logger.info(f"   → Paragraph slide detected")
            if has_bullets:
                logger.info(f"   → Bullets slide detected")
            
            # Process elements
            title_element = None
            content_rendered = False
            
            for element in elements:
                try:
                    element_type = element.get('type')
                    element_id = element.get('id', 'unknown')
                    
                    # Logo
                    if element_type == 'image' and element_id == 'logo':
                        self._add_logo(slide, element)
                    
                    # Title with icon
                    elif element_type == 'text' and element_id == 'title':
                        title_element = element
                        title_text = self._extract_text_value(data, 'title')
                        
                        max_title = self.constraints['text_constraints']['max_title_length']
                        if len(title_text) > max_title:
                            title_text = title_text[:max_title - 3] + "..."
                            data.title = title_text
                        
                        # ✅ FIX: Add icon to title for ALL slides including table/chart
                        if self.icon_service and title_text:
                            icon_name = self.icon_service.auto_select_icon(title_text, "")
                            self._add_icon_to_title(slide, icon_name, element, slide_data=data)
                        
                        self._add_text_master(slide, element, data, content_type)
                    
                    # Title underline
                    elif element_type == 'line' and element_id == 'title_line':
                        if title_element:
                            self._add_title_underline(slide, title_element)
                    
                    # Chart
                    elif element_type == 'chart' and has_chart and not content_rendered:
                        logger.info(f"   → Rendering CHART")
                        self._add_chart_master(slide, element, data)
                        content_rendered = True
                    
                    # Table
                    elif element_type == 'table' and has_table and not content_rendered:
                        logger.info(f"   → Rendering TABLE")
                        self._add_table_master(slide, element, data)
                        content_rendered = True
                    
                    # Paragraph
                    elif element_type == 'text_paragraph' and has_paragraph and not content_rendered:
                        logger.info(f"   → Rendering PARAGRAPH")
                        self._add_paragraph_text(slide, element, data)
                        content_rendered = True
                    
                    # Bullets
                    elif (element_type in ['bullets', 'text_bullets'] or element_id == 'content') and has_bullets and not content_rendered:
                        logger.info(f"   → Rendering BULLETS")
                        self._add_bullets_master(slide, element, data)
                        content_rendered = True
                    
                    # Four-box layout
                    elif element_type == 'boxes' and element_id == 'content_boxes':
                        logger.info(f"   → Rendering FOUR-BOX")
                        self._add_content_box_with_icon_enhanced(slide, element, data)
                        content_rendered = True
                    
                    # Standalone icons
                    elif element_type == 'icon' and element_id != 'icon':
                        self._add_icon_master(slide, element, data)
                    
                    # Content boxes
                    elif element_type == 'content_box':
                        self._add_content_box_master(slide, element, data)
                
                except Exception as e:
                    logger.warning(f"⚠️  Element '{element_id}' error: {e}")
            
            # Add page number
            if page_num and self.config.get('page_numbering', {}).get('enabled', True):
                skip_title = self.config['page_numbering'].get('skip_title_slide', True)
                skip_sections = self.config['page_numbering'].get('skip_section_headers', False)
                
                should_add = True
                
                if content_type == "section" and skip_sections:
                    should_add = False
                
                if should_add:
                    self.add_page_number(slide, page_num)
        
        except Exception as e:
            logger.error(f"❌ Slide creation error: {e}")
            raise
        
    # ========================================================================
    # ICON INTEGRATION - RTL/LTR AWARE
    # ========================================================================

    def _add_icon_to_title(self, slide, icon_name: str, title_element: Dict, slide_data=None) -> None:
        """Add icon next to title with language-specific positioning"""
        if not self.icon_service:
            return
        
        try:
            # *** FIX: Check if slide_data has icon_name field ***
            if slide_data and hasattr(slide_data, 'icon_name') and slide_data.icon_name:
                icon_name_from_data = slide_data.icon_name
                logger.info(f"🎯 Using icon_name from slide_data: {icon_name_from_data}")
                
                # Try to match the icon_name
                matched_icon = self.icon_service.fuzzy_match_icon_name(icon_name_from_data)
                if matched_icon:
                    icon_name = matched_icon
                    logger.info(f"✅ Matched icon: {icon_name_from_data} → {matched_icon}")
                else:
                    logger.warning(f"⚠️ Could not match icon_name: {icon_name_from_data}, using auto-selected: {icon_name}")
            
            pos = self.get_position(title_element, 'position')
            style = title_element.get('style', {})
            
            text_color = self.get_style_value(style, 'color', '#01415C')
            
            # Get icon positioning from constraints
            positioning = self.constraints.get('positioning', {})
            lang_pos = positioning.get(self.target_language, {})
            
            is_rtl = self.lang_config.get('rtl', False)
            
            if is_rtl:
                # Icon on the right for RTL
                icon_left = lang_pos.get('icon_offset', 11.83)
            else:
                # Icon on the left for LTR
                icon_left = lang_pos.get('icon_offset', 1.0)
            
            icon_top = pos.get('top', 0.6) + 0.05
            
            icon_config = self.constraints.get('icons', {})
            icon_size = icon_config.get('size', 0.5)
            
            icon_data = self.icon_service.render_to_png(
                icon_name,
                size=int(icon_size * 96),
                color=text_color
            )
            
            if icon_data:
                slide.shapes.add_picture(
                    icon_data,
                    Inches(icon_left),
                    Inches(icon_top),
                    width=Inches(icon_size),
                    height=Inches(icon_size)
                )
                logger.debug(f"✅ Icon added: {icon_name} ({'RTL' if is_rtl else 'LTR'})")
        except Exception as e:
            logger.warning(f"⚠️ Icon render failed: {e}")


    def _add_centered_section_icon(self, slide, icon_name: str, title_text: str, layout_config: Dict, slide_data=None) -> None:
        """Add icon centered above title for section headers"""
        if not self.icon_service:
            return
        
        try:
            # *** FIX: Check if slide_data has icon_name field ***
            if slide_data and hasattr(slide_data, 'icon_name') and slide_data.icon_name:
                icon_name_from_data = slide_data.icon_name
                logger.info(f"🎯 Section using icon_name from slide_data: {icon_name_from_data}")
                
                # Try to match the icon_name
                matched_icon = self.icon_service.fuzzy_match_icon_name(icon_name_from_data)
                if matched_icon:
                    icon_name = matched_icon
                    logger.info(f"✅ Section matched icon: {icon_name_from_data} → {matched_icon}")
                else:
                    logger.warning(f"⚠️ Section could not match icon_name: {icon_name_from_data}, using auto-selected: {icon_name}")
            
            elements = layout_config.get('elements', [])
            title_element = next((e for e in elements if e.get('id') == 'title'), None)
            
            if not title_element:
                return
            
            title_style = title_element.get('style', {})
            text_color = self.get_style_value(title_style, 'color', '#FFFCEC')
            
            slide_width = self.constraints['layout']['slide_width']
            
            icon_config = self.constraints.get('icons', {})
            icon_size = icon_config.get('section_icon_size', 1.2)
            
            icon_left = (slide_width - icon_size) / 2
            icon_top = 2.2
            
            icon_data = self.icon_service.render_to_png(
                icon_name,
                size=int(icon_size * 96),
                color=text_color
            )
            
            if icon_data:
                slide.shapes.add_picture(
                    icon_data,
                    Inches(icon_left),
                    Inches(icon_top),
                    width=Inches(icon_size),
                    height=Inches(icon_size)
                )
                logger.info(f"✅ Section icon: {icon_name}")
        
        except Exception as e:
            logger.warning(f"⚠️ Section icon failed: {e}")

    def _add_title_underline(self, slide, title_element: Dict) -> None:
        """Add horizontal line below title"""
        try:
            # Position the underline dynamically based on the title's position/size
            pos = self.get_position(title_element, 'position')
            size = title_element.get('size', {})

            left = float(pos.get('left', 1.0))
            width = float(size.get('width', 11.33))

            # Place the line just below the title box with a small margin
            title_top = float(pos.get('top', 1.0))
            title_height = float(size.get('height', 0.8))
            margin = 0.15
            top = title_top + title_height + margin

            line_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(0.02)
            )
            
            line_shape.fill.solid()
            line_color = self.get_color_rgb('accent_teal')
            line_shape.fill.fore_color.rgb = RGBColor(*line_color)
            line_shape.line.fill.background()
            
            logger.debug("✅ Title underline added")
        
        except Exception as e:
            logger.warning(f"⚠️  Title underline failed: {e}")

    # ========================================================================
    # TEXT ELEMENTS - RTL/LTR AWARE
    # ========================================================================

    def _add_text_master(self, slide, element: Dict, data, content_type: str) -> None:
        """Add text with proper RTL/LTR support including alignment"""
        pos = self.get_position(element, 'position')
        size = element.get('size', {})
        style = element.get('style', {})
        element_id = element.get('id', '')

        # Extract text
        if element_id == 'title':
            text = self._extract_text_value(data, 'title')
        elif element_id == 'content':
            text = self._extract_text_value(data, 'content')
        else:
            text = self._extract_text_value(data, element_id)

        if not text:
            return

        textbox = slide.shapes.add_textbox(
            Inches(pos.get('left', 0)),
            Inches(pos.get('top', 0)),
            Inches(size.get('width', 5)),
            Inches(size.get('height', 1))
        )

        tf = textbox.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP

        # *** FIX: Set RTL/LTR at text frame level ***
        is_rtl = self.lang_config.get('rtl', False)
        self._set_text_frame_rtl(tf, is_rtl)

        padding = float(style.get('padding', 0))
        margin_value = Inches(max(padding, 0))
        
        # *** FIX: Language-specific margins ***
        if is_rtl:
            tf.margin_left = Inches(0.1)
            tf.margin_right = margin_value
        else:
            tf.margin_left = margin_value
            tf.margin_right = Inches(0.1)
        
        tf.margin_top = tf.margin_bottom = margin_value

        line_spacing = self.get_style_value(style, 'line_spacing', 1.45)
        
        # *** FIX: Proper alignment logic for sections and titles ***
        if content_type in ['title']:
            # Title slide - always center
            alignment = PP_ALIGN.CENTER
        elif content_type == 'section':
            # Section headers - always center
            alignment = PP_ALIGN.CENTER
        else:
            # Content slides - follow constraints
            alignment = self.get_text_alignment(style)
        
        # Get text color
        text_color = self.get_style_value(style, 'color')
        if not text_color:
            if content_type in ['title', 'section']:
                text_color = '#FFFCEC'
            else:
                text_color = '#0D2026'
        
        # Get font
        font_name = self.get_style_value(style, 'font')
        if not font_name:
            if element_id == 'title':
                font_name = self.get_font('heading')
            else:
                font_name = self.get_font('body')

        # Split into paragraphs
        paragraphs = [line.strip() for line in str(text).splitlines() if line.strip()]
        if not paragraphs:
            paragraphs = [str(text).strip()]

        for idx, line in enumerate(paragraphs):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = line
            
            font_size = self.get_style_value(style, 'font_size', 18)
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = font_name
            
            if self.get_style_value(style, 'bold'):
                paragraph.font.bold = True

            rgb = self.hex_to_rgb(text_color)
            paragraph.font.color.rgb = RGBColor(*rgb)

            paragraph.line_spacing = line_spacing
            paragraph.space_after = Pt(self.get_style_value(style, 'space_after', 6))
            
            # *** FIX: Apply alignment ***
            paragraph.alignment = alignment
            
            logger.debug(f"Paragraph: '{line[:30]}...' align={alignment} rtl={is_rtl}")


    def _set_text_frame_rtl(self, text_frame, is_rtl: bool) -> None:
        """Set text frame direction to RTL or LTR via XML manipulation"""
        try:
            from pptx.oxml import parse_xml
            
            tf_element = text_frame._element
            
            # Find or create bodyPr (body properties)
            bodyPr = tf_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
            
            if bodyPr is not None:
                # Set RTL attribute
                if is_rtl:
                    bodyPr.set('rtlCol', '1')
                    bodyPr.set('anchor', 'ctr')
                else:
                    # Remove RTL for LTR
                    if 'rtlCol' in bodyPr.attrib:
                        del bodyPr.attrib['rtlCol']
                    bodyPr.set('anchor', 'ctr')
            
            logger.debug(f"Text frame direction set to {'RTL' if is_rtl else 'LTR'}")
        except Exception as e:
            logger.debug(f"Could not set text frame direction: {e}")

    def _add_paragraph_text(self, slide, element: Dict, data) -> None:
        """Render paragraph content with proper RTL/LTR support"""
        from apps.app.utils.text_formatter import should_convert_to_bullets, break_long_paragraph_to_bullets
        
        pos = self.get_position(element, 'position')
        size = element.get('size', {})
        style = element.get('style', {})

        # Get text
        text = getattr(data, 'content', None)
        
        if not text or len(str(text).strip()) == 0:
            if hasattr(data, 'bullets') and data.bullets:
                text = "\n\n".join(b.text for b in data.bullets if getattr(b, 'text', None))

        text = (text or "").strip()
        if not text:
            logger.warning(f"⚠️  No paragraph text found")
            return
        
        # ✅ NEW: Check if text should be converted to bullets for better readability
        if should_convert_to_bullets(text):
            logger.info(f"   📝 Converting long paragraph to bullets ({len(text)} chars)")
            bullets = break_long_paragraph_to_bullets(text)
            if bullets and len(bullets) > 1:
                # Set bullets and render as bullet slide instead
                data.bullets = bullets
                self._add_bullets_master(slide, element, data)
                logger.info(f"   ✅ Converted to {len(bullets)} bullet points")
                return

        # Create textbox
        textbox = slide.shapes.add_textbox(
            Inches(pos.get('left', 1.5)),
            Inches(pos.get('top', 2.0)),
            Inches(size.get('width', 10.33)),
            Inches(size.get('height', 4.8))
        )

        tf = textbox.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP  # ✅ Keep TOP alignment
        
        # Set RTL direction for text frame
        is_rtl = self.lang_config.get('rtl', False)
        self._set_text_frame_rtl(tf, is_rtl)
        
        # Language-aware margins
        padding = float(style.get('padding', 0.2))
        if is_rtl:
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(padding)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)
        else:
            tf.margin_left = Inches(padding)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.1)
            tf.margin_bottom = Inches(0.1)

        # Split into paragraphs (prioritize double line breaks)
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        if not paragraphs:
            # Fallback to single line breaks
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        if not paragraphs:
            paragraphs = [text]

        # Get styling
        font_name = self.get_font('body')  # or self.get_font_body() depending on your method
        text_color = self.get_style_value(style, 'color', '#0D2026')
        font_size = self.get_font_size('body')
        line_spacing = self.get_line_spacing('body')
        alignment = self.get_text_alignment(style)

        # Add paragraphs
        for idx, chunk in enumerate(paragraphs):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = chunk
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = font_name
            paragraph.line_spacing = line_spacing
            paragraph.space_after = Pt(12)
            paragraph.alignment = alignment

            # Set color
            rgb = self.hex_to_rgb(text_color)
            paragraph.font.color.rgb = RGBColor(*rgb)
            
            # ✅ FIX: Set paragraph-level RTL if needed
            if is_rtl:
                pPr = paragraph._element.get_or_add_pPr()
                pPr.set('rtl', '1')
        
        logger.info(f"   ✅ Paragraph: {len(paragraphs)} blocks, RTL={is_rtl}")

    def _add_bullets_master(self, slide, element: Dict, data) -> None:
        """Add bullets with proper RTL/LTR alignment using native PowerPoint bullets"""
        element_id = element.get('id')
        bullets = self._extract_bullet_items(data, element_id)
        if not bullets:
            return
        
        # Get constraints
        bullet_config = self.constraints.get('bullets', {})
        max_bullets = bullet_config.get('max_bullets_per_slide', 6)
        bullets = bullets[:max_bullets]
        
        pos = self.get_position(element, 'position')
        size = element.get('size', {})
        style = element.get('style', {})
        bullet_style_cfg = element.get('bullet_style', {})
        
        # Create textbox – use same positioning conventions as paragraph text
        # so bullets align visually with paragraphs and make good use of space
        textbox = slide.shapes.add_textbox(
            Inches(pos.get('left', 1.5)),
            Inches(pos.get('top', 2.0)),
            Inches(size.get('width', 10.33)),
            Inches(size.get('height', 4.8))
        )
        
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP  # ✅ Keep TOP alignment
        
        # Set RTL at text frame level
        is_rtl = self.lang_config.get('rtl', False)
        self._set_text_frame_rtl(tf, is_rtl)
        
        # Language-aware margins (match paragraph padding for consistent start)
        padding = float(self.get_style_value(style, 'padding', 0.2))
        if is_rtl:
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(padding)
        else:
            tf.margin_left = Inches(padding)
            tf.margin_right = Inches(0.05)
        
        # Get styling
        font_size = self.get_font_size('body')
        line_spacing = bullet_config.get('line_spacing', 1.75)
        font_name = self.get_font('body')
        text_color = self.get_style_value(style, 'color', '#0D2026')
        text_rgb = self.hex_to_rgb(text_color)
        
        # Get alignment from constraints
        alignment_config = self.constraints.get('alignment', {})
        lang_alignment = alignment_config.get(self.target_language, {})
        bullets_alignment_str = lang_alignment.get('bullets', 'right' if is_rtl else 'left')
        text_alignment = self._alignment_enum(bullets_alignment_str)
        
        # Get bullet symbols
        bullet_symbols = bullet_config.get('bullet_symbols', {})
        bullet_char = bullet_symbols.get('level_1', '●')
        bullet_color = self.get_style_value(bullet_style_cfg, 'bullet_color', '#01415C')
        bullet_rgb = self.hex_to_rgb(bullet_color)
        
        logger.info(f"   Bullets: lang={self.target_language}, align={bullets_alignment_str}, RTL={is_rtl}")
        
        # Add bullets
        for idx, bullet in enumerate(bullets):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            
            # Clean bullet text (remove manual symbols)
            bullet_text = (bullet.text or "").replace("●", "").replace("**", "").strip()
            
            # ✅ FIX: Set text WITHOUT manual bullet symbol
            paragraph.text = bullet_text
            paragraph.level = 0  # Level 0 = main bullet
            
            # Font styling
            paragraph.font.bold = self.get_style_value(style, 'bold', False)
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = font_name
            paragraph.font.color.rgb = RGBColor(*text_rgb)
            
            # Alignment
            paragraph.alignment = text_alignment
            paragraph.line_spacing = line_spacing
            
            # Spacing
            spacing_before = bullet_config.get('spacing_before', 6)
            spacing_after = bullet_config.get('spacing_after', 6)
            paragraph.space_before = Pt(spacing_before)
            paragraph.space_after = Pt(spacing_after)
            
            # ✅ FIX: Set paragraph RTL property
            if is_rtl:
                pPr = paragraph._element.get_or_add_pPr()
                pPr.set('rtl', '1')
            
            # ✅ FIX: Apply native PowerPoint bullet formatting
            pPr = paragraph._element.get_or_add_pPr()
            
            # Add bullet formatting
            buFont = parse_xml(f'<a:buFont xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font_name}"/>')
            buChar = parse_xml(f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{bullet_char}"/>')
            
            # Add bullet color
            buClr = parse_xml(f'''<a:buClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                <a:srgbClr val="{bullet_color.replace("#", "")}"/>
            </a:buClr>''')
            
            pPr.append(buFont)
            pPr.append(buChar)
            pPr.append(buClr)
            
            # Handle sub-bullets
            if getattr(bullet, 'sub_bullets', None):
                sub_bullet_char = bullet_symbols.get('level_2', '○')
                max_sub = bullet_config.get('max_sub_bullets_per_bullet', 3)
                
                for sub in bullet.sub_bullets[:max_sub]:
                    sp = tf.add_paragraph()
                    sub_text = str(sub).replace("○", "").replace("●", "").strip()
                    
                    sp.text = sub_text
                    sp.level = 1  # Level 1 = sub-bullet
                    
                    sp.font.size = Pt(max(font_size - 2, 14))
                    sp.font.name = font_name
                    sp.font.color.rgb = RGBColor(*text_rgb)
                    sp.alignment = text_alignment
                    sp.line_spacing = line_spacing
                    
                    # RTL for sub-bullet
                    if is_rtl:
                        spPr = sp._element.get_or_add_pPr()
                        spPr.set('rtl', '1')
                    
                    # Apply sub-bullet formatting
                    spPr = sp._element.get_or_add_pPr()
                    sub_buFont = parse_xml(f'<a:buFont xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{font_name}"/>')
                    sub_buChar = parse_xml(f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{sub_bullet_char}"/>')
                    sub_buClr = parse_xml(f'''<a:buClr xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                        <a:srgbClr val="{bullet_color.replace("#", "")}"/>
                    </a:buClr>''')
                    
                    spPr.append(sub_buFont)
                    spPr.append(sub_buChar)
                    spPr.append(sub_buClr)
        
        logger.info(f"   ✅ Bullets: {len(bullets)} items added")


    # ========================================================================
    # TABLE - RTL/LTR AWARE
    # ========================================================================

    def _add_table_master(self, slide, element: Dict, data) -> None:
        """Add table with RTL/LTR support and proper language context"""
        table_data_obj = getattr(data, 'table_data', None) or getattr(data, 'table', None)
        
        if not table_data_obj:
            logger.warning(f"⚠️  No table_data found")
            return
        
        pos = self.get_position(element, 'position')
        size = element.get('size', {})
        style = element.get('style', {})
        
        # *** FIX: Import and initialize TableService with language ***
        from apps.app.services.table_service import TableService
        
        try:
            # Initialize TableService with template and language
            table_service = TableService(
                template_id=self.template_id,
                language=self.target_language  # Pass the target language
            )
            
            # Add table
            table_service.add_table(
                slide=slide,
                table_data=table_data_obj,
                position=pos,
                size=size
            )
            
            logger.info(f"✅ Table rendered (RTL={self.lang_config.get('rtl', False)})")
            
        except Exception as e:
            logger.error(f"❌ Table error: {e}")
            logger.exception(e)

    # ========================================================================
    # CHART - RTL/LTR AWARE
    # ========================================================================

    def _validate_chart_data(self, chart_data_obj) -> Tuple[bool, str]:
        """Validate chart has required data before rendering"""
        if not chart_data_obj:
            return False, "Chart data is None"
        
        # Convert to dict
        if hasattr(chart_data_obj, 'dict'):
            chart_data = chart_data_obj.dict()
        elif isinstance(chart_data_obj, dict):
            chart_data = chart_data_obj
        else:
            try:
                chart_data = chart_data_obj.__dict__
            except:
                return False, "Cannot extract chart data"
        
        categories = chart_data.get('categories', [])
        series = chart_data.get('series', [])
        
        if not categories or len(categories) == 0:
            return False, "Categories array is empty"
        
        if not series or len(series) == 0:
            return False, "Series array is empty"
        
        # Check first series has values
        first_series = series[0]
        if isinstance(first_series, dict):
            values = first_series.get('values', [])
        else:
            values = getattr(first_series, 'values', [])
        
        if not values or len(values) == 0:
            return False, "Series values are empty"
        
        if len(values) != len(categories):
            return False, f"Data mismatch: {len(categories)} categories but {len(values)} values"
        
        return True, "Valid"

    def _add_chart_master(self, slide, element: Dict, data) -> None:
        """Add chart with RTL/LTR support"""
        if not self.chart_service:
            logger.warning("⚠️  ChartService not available")
            return
        
        chart_data_obj = getattr(data, 'chart_data', None) or getattr(data, 'chart', None)
        
        if not chart_data_obj:
            logger.warning(f"⚠️  No chart_data found")
            return
        
        # ✅ NEW: Validate chart data before rendering
        is_valid, error_msg = self._validate_chart_data(chart_data_obj)
        if not is_valid:
            logger.error(f"❌ Invalid chart data: {error_msg}")
            logger.error(f"   Chart object: {chart_data_obj}")
            
            # Create error message textbox instead of crashing
            pos = self.get_position(element, 'position')
            textbox = slide.shapes.add_textbox(
                Inches(pos.get('left', 2.0)),
                Inches(pos.get('top', 3.0)),
                Inches(9.0),
                Inches(2.0)
            )
            tf = textbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"⚠️ Chart Generation Error\n\n{error_msg}\n\nPlease check source data and regenerate."
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(220, 53, 69)  # Red
            p.alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            logger.warning(f"   Created error placeholder for chart")
            return
        
        pos = self.get_position(element, 'position')
        size = element.get('size', {})
        style = element.get('style', {})
        
        try:
            # Convert to dict
            if hasattr(chart_data_obj, 'dict'):
                chart_data = chart_data_obj.dict()
            elif isinstance(chart_data_obj, dict):
                chart_data = chart_data_obj
            else:
                chart_data = chart_data_obj.__dict__
            
            logger.info(f"   → Chart type: {chart_data.get('chart_type')}")
            
            # Get chart config from constraints
            chart_config = self.constraints.get('chart', {})
            
            # Apply styling from constraints
            chart_data['font_color'] = chart_config.get('font_color', '#FFFFFF')
            chart_data['data_label_color'] = chart_config.get('data_label_color', '#FFFFFF')
            chart_data['axis_color'] = chart_config.get('axis_color', '#FFFFFF')
            chart_data['axis_label_color'] = chart_config.get('axis_label_color', '#FFFFFF')
            chart_data['grid_color'] = chart_config.get('grid_color', '#5C6F7A')
            chart_data['legend_font_color'] = chart_config.get('legend_font_color', '#FFFFFF')
            chart_data['title_color'] = chart_config.get('title_color', '#FFFCEC')
            
            bg_color = chart_config.get('background_color', '#0D2026')
            bg_rgb = self.hex_to_rgb(bg_color)
            
            self.chart_service.add_native_chart(
                slide=slide,
                chart_data=chart_data,
                position={"left": pos.get('left', 1.5), "top": pos.get('top', 2.0)},
                size={"width": size.get('width', 10.33), "height": size.get('height', 5.0)},
                background_rgb=bg_rgb
            )
            
            logger.info(f"✅ Chart rendered")
        except Exception as e:
            logger.error(f"❌ Chart error: {e}")
            logger.exception(e)

    # ========================================================================
    # FOUR-BOX LAYOUT - RTL/LTR AWARE
    # ========================================================================

    def _add_content_box_with_icon_enhanced(self, slide, element: Dict, data) -> None:
        """Enhanced four-box layout with RTL/LTR support"""
        box_config = element
        bullets = self._extract_bullet_items(data, 'content')
        
        if not bullets or len(bullets) == 0:
            return
        
        # Get box configuration from constraints
        box_constraints = self.constraints.get('boxes', {})
        max_boxes = box_constraints.get('max_boxes', 4)
        bullets = bullets[:max_boxes]
        
        layout_type = box_config.get('layout', '2x2')
        position = box_config.get('position', {})
        size = box_config.get('size', {})
        box_style = box_config.get('box_style', {})
        
        # Box dimensions
        box_width = box_constraints.get('width', 5.4)
        box_height = box_constraints.get('height', 2.2)
        gap_h = box_constraints.get('gap_horizontal', 0.53)
        gap_v = box_constraints.get('gap_vertical', 0.4)
        
        # Colors
        colors = self.constraints.get('colors', {}).get('box_colors', ['#B1D8BE', '#F9D462', '#C6C3BE', '#E09059'])
        text_color = box_style.get('text_color', '#0D2026')
        font_name = self.get_font('body')
        font_size = box_constraints.get('font_size', 16)
        icon_size = box_constraints.get('icon_size', 0.6)
        
        # Base position
        base_left = position.get('left', 1.0)
        base_top = position.get('top', 2.0)
        
        # Get alignment
        is_rtl = self.lang_config.get('rtl', False)
        text_alignment = PP_ALIGN.CENTER  # Boxes are always centered
        
        # Render each box
        for idx, bullet in enumerate(bullets[:4]):
            # Calculate grid position (2x2)
            row = idx // 2
            col = idx % 2
            
            # Calculate box position
            box_left = Inches(base_left + col * (box_width + gap_h))
            box_top = Inches(base_top + row * (box_height + gap_v))
            box_w = Inches(box_width)
            box_h = Inches(box_height)
            
            # 1. Create colored box
            box_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                box_left, box_top, box_w, box_h
            )
            
            box_color = colors[idx % len(colors)]
            box_rgb = self.hex_to_rgb(box_color)
            
            box_shape.fill.solid()
            box_shape.fill.fore_color.rgb = RGBColor(*box_rgb)
            box_shape.line.fill.background()
            
            # 2. Add icon INSIDE box (top-center)
            if self.icon_service:
                icon_name_to_use = None
                
                if hasattr(bullet, 'icon_name') and bullet.icon_name:
                    icon_name_from_bullet = bullet.icon_name
                    logger.info(f"🎯 Box {idx+1} using icon_name: {icon_name_from_bullet}")
                    
                    matched_icon = self.icon_service.fuzzy_match_icon_name(icon_name_from_bullet)
                    if matched_icon:
                        icon_name_to_use = matched_icon
                        logger.info(f"✅ Box {idx+1} matched icon: {icon_name_from_bullet} → {matched_icon}")
                
                # Fallback to auto-selection
                if not icon_name_to_use:
                    icon_name_to_use = self.icon_service.auto_select_icon(bullet.text or "", "")
                
                try:
                    icon_data = self.icon_service.render_to_png(
                        icon_name_to_use,
                        size=int(icon_size * 96),
                        color=text_color
                    )
                    
                    if icon_data:
                        # Center icon horizontally within box
                        icon_left = box_left + (box_w - Inches(icon_size)) / 2
                        icon_top = box_top + Inches(box_constraints.get('icon_top_padding', 0.25))
                        
                        slide.shapes.add_picture(
                            icon_data,
                            icon_left,
                            icon_top,
                            Inches(icon_size),
                            Inches(icon_size)
                        )
                except Exception as e:
                    logger.debug(f"Icon render skipped for box {idx + 1}: {e}")
            
            # 3. Add text BELOW icon (inside box) - WITH TRUNCATION
            text_top = box_top + Inches(icon_size + box_constraints.get('text_top_offset', 0.35))
            text_height = box_h - Inches(icon_size + 0.5)
            
            padding = box_constraints.get('padding', 0.3)
            text_box = slide.shapes.add_textbox(
                box_left + Inches(padding),
                text_top,
                box_w - Inches(padding * 2),
                text_height
            )
            
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = tf.margin_right = Inches(0.1)
            
            # ✅ FIX: Truncate text to prevent overflow (max 100 chars for four-box)
            box_text = (bullet.text or "").strip()
            max_box_length = 100  # Strict limit for four-box layouts
            if len(box_text) > max_box_length:
                # Truncate at word boundary
                truncated = box_text[:max_box_length]
                last_space = truncated.rfind(' ')
                if last_space > max_box_length * 0.75:
                    box_text = truncated[:last_space].strip() + "..."
                else:
                    box_text = truncated.strip() + "..."
                logger.warning(f"   Truncated four-box text: {len(bullet.text)} → {len(box_text)} chars")
            
            p = tf.paragraphs[0]
            p.text = box_text
            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))
            p.alignment = text_alignment
            p.line_spacing = box_constraints.get('line_spacing', 1.3)
            
            logger.debug(f"✅ Four-box item {idx + 1}")

    # ========================================================================
    # ICON MASTER - RTL/LTR AWARE
    # ========================================================================
    def _get_icon_name_from_data(self, data) -> Optional[str]:
        """Extract icon_name from slide data with various fallbacks"""
        # Try direct attribute
        if hasattr(data, 'icon_name') and data.icon_name:
            return data.icon_name
        
        # Try dict access
        if isinstance(data, dict) and 'icon_name' in data:
            return data['icon_name']
        
        return None

    def _add_icon_master(self, slide, element: Dict, data) -> None:
        """Add standalone icon with RTL/LTR positioning"""
        if not self.icon_service:
            return

        element_id = element.get("id", "")
        pos = self.get_position(element, 'position')
        size = element.get('size', {})
        style = element.get('style', {})

        # Determine content for icon selection
        content = ""
        
        if element_id == "icon":
            if hasattr(data, "title") and data.title:
                content = data.title
        elif element_id.endswith("_icon"):
            box_id = element_id.replace("_icon", "")
            if hasattr(data, "bullets") and data.bullets:
                try:
                    idx = int(box_id.replace("box", "")) - 1
                    if 0 <= idx < len(data.bullets):
                        content = data.bullets[idx].text
                except:
                    content = getattr(data, "title", "")

        if not content:
            content = getattr(data, "title", "") or ""

        if not content.strip():
            return

        # Select icon
        icon_name = self.icon_service.auto_select_icon(content, "")
        
        # Get icon config
        icon_config = self.constraints.get('icons', {})
        icon_color = self.get_style_value(style, 'color', '#0D2026')

        try:
            icon_data = self.icon_service.render_to_png(
                icon_name,
                size=int(size.get("width", 0.6) * 96),
                color=icon_color
            )

            if not icon_data:
                return

            slide.shapes.add_picture(
                icon_data,
                Inches(pos.get("left", 0)),
                Inches(pos.get("top", 0)),
                width=Inches(size.get("width", 0.6)),
                height=Inches(size.get("height", 0.6))
            )

            logger.debug(f"✅ Icon added: {icon_name}")

        except Exception as e:
            logger.warning(f"⚠️ Icon render failed: {e}")

    # ========================================================================
    # CONTENT BOX MASTER
    # ========================================================================

    def _add_content_box_master(self, slide, element: Dict, data) -> None:
        """Add colored content box with RTL/LTR support"""
        box_id = element.get('id', '')
        pos = self.get_position(element, 'position')
        size = element.get('size', {})
        style = element.get('style', {})
        
        content = ""
        if isinstance(data, dict):
            content = data.get(box_id, '')
        elif hasattr(data, 'bullets') and data.bullets:
            box_num = int(box_id.replace('box', '')) - 1
            if box_num < len(data.bullets):
                content = data.bullets[box_num].text
        
        if not content:
            return
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos.get('left', 0)),
            Inches(pos.get('top', 0)),
            Inches(size.get('width', 3)),
            Inches(size.get('height', 2))
        )
        
        bg_color = self.get_style_value(style, 'background_color', '#FFFFFF')
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))
        box.line.fill.background()
        
        padding = self.get_style_value(style, 'padding', 0.3)
        textbox = slide.shapes.add_textbox(
            Inches(pos.get('left') + padding),
            Inches(pos.get('top') + padding),
            Inches(size.get('width') - 2*padding),
            Inches(size.get('height') - 2*padding)
        )
        
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(self.get_style_value(style, 'font_size', 16))
        p.font.name = self.get_font('body')
        p.alignment = self.get_text_alignment(style)
        p.line_spacing = self.get_style_value(style, 'line_spacing', 1.45)
        p.space_after = Pt(self.get_style_value(style, 'space_after', 4))
        
        text_color = self.get_style_value(style, 'text_color', '#0D2026')
        p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))

    # ========================================================================
    # OUTPUT PATH
    # ========================================================================

    def _get_output_path(self, title: str) -> str:
        """Generate output file path"""
        output_dir = Path(settings.OUTPUT_DIR)
        output_dir.mkdir(parents=True, exist_ok=True)
        
        safe_title = re.sub(r'[^\w\s-]', '', title or 'presentation')[:50]
        safe_title = safe_title.strip().replace(' ', '_')
        
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        lang_code = self.target_language[:2]
        
        filename = f"{self.template_id}_{lang_code}_{safe_title}_{timestamp}.pptx"
        return str(output_dir / filename)
    
    def generate_to_bytes(self, presentation_data: PresentationData) -> bytes:
        """Generate presentation and return as bytes"""
        output_path = self.generate(presentation_data)
        with open(output_path, 'rb') as f:
            return f.read()
