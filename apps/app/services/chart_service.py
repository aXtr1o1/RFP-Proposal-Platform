from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from typing import Dict, Tuple, Optional, List
from pathlib import Path
import logging
import json

from ..config import settings

logger = logging.getLogger("chart_service")


def _chart_data_to_dict(chart_data) -> Dict:
    """Convert ChartData model or dict to dict for chart_service."""
    if chart_data is None:
        return {}
    if isinstance(chart_data, dict):
        return chart_data
    if hasattr(chart_data, "model_dump"):
        return chart_data.model_dump()
    if hasattr(chart_data, "dict"):
        return chart_data.dict()
    return {}

class ChartService:
    """Service for creating native PowerPoint charts - fully dynamic from constraints.json"""
    
    def __init__(self, template_id: str = "arweqah", language: str = "en"):
        self.template_id = template_id
        self.language = language
        
        # Load constraints dynamically
        self.constraints = self._load_constraints(template_id)
        
        # Get chart configuration from constraints
        self.chart_config = self.constraints.get('chart', {})
        self.colors_config = self.constraints.get('colors', {})
        self.typography_config = self.constraints.get('typography', {})
        
        # Get language-specific fonts
        fonts = self.typography_config.get('fonts', {})
        self.lang_fonts = fonts.get(language, fonts.get('en', {}))
        
        # Get chart-specific settings
        self.color_palette = self.chart_config.get('color_palette', [
            "#F9D462", "#B1D8BE", "#E09059", "#C6C3BE", "#FFFCEC"
        ])
        
        self.default_font = self.lang_fonts.get('body', 'Tajawal')
        self.font_size = self.chart_config.get('font_size', 12)
        self.title_font_size = self.chart_config.get('title_font_size', 16)
        
        # Colors from constraints
        self.font_color = self.chart_config.get('font_color', '#FFFFFF')
        self.axis_color = self.chart_config.get('axis_color', '#FFFFFF')
        self.axis_label_color = self.chart_config.get('axis_label_color', '#FFFFFF')
        self.grid_color = self.chart_config.get('grid_color', '#5C6F7A')
        self.legend_font_color = self.chart_config.get('legend_font_color', '#FFFFFF')
        self.data_label_color = self.chart_config.get('data_label_color', '#FFFFFF')
        self.background_color = self.chart_config.get('background_color', '#0D2026')
        
        logger.info(f"✅ ChartService initialized (template={template_id}, lang={language})")
    
    def _load_constraints(self, template_id: str) -> Dict:
        """Load constraints.json from template directory"""
        path = Path(settings.TEMPLATES_DIR) / template_id / "constraints.json"
        with open(path, "r", encoding="utf-8") as file:
            return json.load(file)
    
    def add_native_chart(
        self,
        slide,
        chart_data: Dict,
        position: Dict,
        size: Dict,
        background_rgb: Optional[Tuple[int, int, int]] = None
    ):
        """Add native PowerPoint chart - fully dynamic styling"""
        try:
            chart_type = chart_data.get('chart_type', 'column')
            title = chart_data.get('title', '')
            
            # Extract data dynamically
            categories, series_list = self._extract_chart_data(chart_data)
            
            if not categories or not series_list:
                logger.error(f"❌ Chart data extraction failed!")
                return None
            
            if len(categories) == 0 or len(series_list) == 0:
                logger.error(f"❌ Chart data empty!")
                return None
            
            logger.info(f"✓ Creating {chart_type} chart")
            logger.info(f"   Categories: {categories}")
            logger.info(f"   Series count: {len(series_list)}")
            
            # Map chart types
            chart_type_map = {
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'line': XL_CHART_TYPE.LINE_MARKERS,
                'pie': XL_CHART_TYPE.PIE,
                'area': XL_CHART_TYPE.AREA
            }
            
            xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Create chart data object
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = categories
            
            # Add all series
            for series_info in series_list:
                chart_data_obj.add_series(series_info['name'], series_info['values'])
            
            # Add chart to slide
            x = Inches(position['left'])
            y = Inches(position['top'])
            cx = Inches(size['width'])
            cy = Inches(size['height'])
            
            graphic_frame = slide.shapes.add_chart(
                xl_chart_type, x, y, cx, cy, chart_data_obj
            )
            
            chart = graphic_frame.chart

            # Get axis labels and unit
            x_axis_label = chart_data.get('x_axis_label', '')
            y_axis_label = chart_data.get('y_axis_label', '')
            unit = chart_data.get('unit', '')

            # Build color config from constraints
            color_config = {
                "font_color": self.font_color,
                "data_label_color": self.data_label_color,
                "axis_color": self.axis_color,
                "axis_label_color": self.axis_label_color,
                "grid_color": self.grid_color,
                "legend_font_color": self.legend_font_color,
                "title_color": self.font_color
            }
            
            # Apply styling
            self._apply_modern_chart_style(chart, chart_type, background_rgb, color_config)
            
            # Add data labels
            self._add_data_labels(chart, chart_type, unit, color_config)
            
            # Configure legend
            self._configure_legend(chart, chart_type, color_config)
            
            # Add axis labels
            self._add_axis_labels(chart, chart_type, x_axis_label, y_axis_label, unit, color_config)
            
            # Add title if provided
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(self.title_font_size)
                chart.chart_title.text_frame.paragraphs[0].font.bold = True
                chart.chart_title.text_frame.paragraphs[0].font.name = self.default_font
                title_color = self._get_rgb(color_config.get("font_color"))
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*title_color)
            
            logger.info(f"✅ Chart created: {title} ({chart_type})")
            return chart
            
        except Exception as e:
            logger.error(f"❌ Chart creation error: {e}")
            import traceback
            traceback.print_exc()
            return None

    def create_chart(
        self,
        slide,
        chart_data,
        left: float,
        top: float,
        width: float,
        height: float,
    ):
        """
        Create a chart on the slide (API used by pptx_generator).
        Accepts position/size in inches and ChartData model or dict.
        """
        data = _chart_data_to_dict(chart_data)
        if not data:
            logger.error("create_chart: no chart data")
            return None
        position = {"left": left, "top": top}
        size = {"width": width, "height": height}
        return self.add_native_chart(slide, data, position, size)
    
    def _extract_chart_data(self, chart_data: Dict) -> Tuple[List, List[Dict]]:
        """Dynamically extract chart data from multiple formats"""
        categories = []
        series_list = []
        
        # FORMAT 1: NEW - categories + series (PREFERRED)
        if 'series' in chart_data and chart_data.get('series'):
            raw_series = chart_data.get('series', [])
            categories = chart_data.get('categories', [])
            
            if isinstance(raw_series, list) and len(raw_series) > 0:
                for series_item in raw_series:
                    if isinstance(series_item, dict):
                        name = series_item.get('name', 'Series')
                        values = series_item.get('values', [])
                        if values:
                            series_list.append({'name': name, 'values': values})
                
                if categories and series_list:
                    logger.info(f"   ✓ NEW format: {len(categories)} categories, {len(series_list)} series")
                    return categories, series_list
        
        # FORMAT 2: OLD - labels + values (FALLBACK)
        if 'labels' in chart_data or 'values' in chart_data:
            categories = chart_data.get('labels', []) or chart_data.get('categories', [])
            values = chart_data.get('values', [])
            series_name = chart_data.get('series_name', 'Values')
            
            if categories and values:
                series_list = [{'name': series_name, 'values': values}]
                logger.info(f"   ✓ OLD format: {len(categories)} labels, 1 series")
                return categories, series_list
        
        # FORMAT 3: LAST RESORT
        if 'categories' in chart_data and chart_data.get('categories'):
            categories = chart_data.get('categories', [])
            for key in ['values', 'data', 'y']:
                if key in chart_data and chart_data[key]:
                    values = chart_data[key]
                    series_list = [{'name': 'Values', 'values': values}]
                    logger.info(f"   ✓ FALLBACK format: found data in '{key}'")
                    return categories, series_list
        
        logger.error(f"   ❌ Could not extract chart data")
        return [], []
        
    def _add_axis_labels(self, chart, chart_type: str, x_axis_label: str, y_axis_label: str, unit: str, color_config: Dict):
        """Add axis labels with dynamic colors from constraints"""
        try:
            if chart_type == 'pie':
                return
            
            axis_label_rgb = self._get_rgb(color_config.get("axis_label_color"))

            # Category axis (X-axis)
            if x_axis_label:
                try:
                    cat_axis = chart.category_axis
                    cat_axis.has_title = True
                    cat_axis.axis_title.text_frame.text = x_axis_label
                    
                    for paragraph in cat_axis.axis_title.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.bold = True
                        paragraph.font.name = self.default_font
                        paragraph.font.color.rgb = RGBColor(*axis_label_rgb)
                    
                    logger.info(f"  ✓ X-axis label: {x_axis_label}")
                except Exception as e:
                    logger.error(f"❌ X-axis label error: {e}")
            
            # Value axis (Y-axis)
            if y_axis_label or unit:
                try:
                    val_axis = chart.value_axis
                    val_axis.has_title = True
                    
                    if y_axis_label and unit:
                        axis_title = f"{y_axis_label} ({unit})"
                    elif y_axis_label:
                        axis_title = y_axis_label
                    else:
                        axis_title = unit
                    
                    val_axis.axis_title.text_frame.text = axis_title
                    
                    for paragraph in val_axis.axis_title.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.bold = True
                        paragraph.font.name = self.default_font
                        paragraph.font.color.rgb = RGBColor(*axis_label_rgb)
                    
                    logger.info(f"  ✓ Y-axis label: {axis_title}")
                except Exception as e:
                    logger.error(f"❌ Y-axis label error: {e}")
                    
        except Exception as e:
            logger.error(f"❌ Axis labels error: {e}")

    def _add_data_labels(self, chart, chart_type: str, unit: str, color_config: Dict):
        """Add data labels with dynamic colors from constraints"""
        try:
            for series_idx, series in enumerate(chart.series):
                series.has_data_labels = True
                data_labels = series.data_labels
                
                data_label_rgb = self._get_rgb(color_config.get("data_label_color"))

                if chart_type == 'pie':
                    data_labels.number_format = '0%'
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    data_labels.font.size = Pt(11)
                    data_labels.font.bold = True
                    data_labels.font.name = self.default_font
                    # Pie labels outside - use dark color for visibility
                    pie_label_rgb = self._get_rgb(self.colors_config.get('text_primary', '#0D2026'))
                    data_labels.font.color.rgb = RGBColor(*pie_label_rgb)
                    
                elif chart_type in ['column', 'bar']:
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    data_labels.font.size = Pt(11)
                    data_labels.font.bold = True
                    data_labels.font.name = self.default_font
                    data_labels.font.color.rgb = RGBColor(*data_label_rgb)
                    
                    if unit:
                        if unit == '%':
                            data_labels.number_format = '0"%"'
                        elif unit == '$':
                            data_labels.number_format = '"$"0'
                        else:
                            data_labels.number_format = f'0" {unit}"'
                    
                elif chart_type == 'line':
                    data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
                    data_labels.font.size = Pt(10)
                    data_labels.font.name = self.default_font
                    data_labels.font.color.rgb = RGBColor(*data_label_rgb)
                    
                    if unit:
                        if unit == '%':
                            data_labels.number_format = '0"%"'
                        elif unit == '$':
                            data_labels.number_format = '"$"0'
                        else:
                            data_labels.number_format = f'0" {unit}"'
                
                logger.debug(f"  ✓ Data labels added to series {series_idx + 1}")
                
        except Exception as e:
            logger.warning(f"⚠️  Data labels error: {e}")

    def _configure_legend(self, chart, chart_type: str, color_config: Dict):
        """Configure chart legend with dynamic colors from constraints"""
        try:
            legend_rgb = self._get_rgb(color_config.get("legend_font_color"))
            
            if chart_type == 'pie':
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
                chart.legend.font.size = Pt(11)
                chart.legend.font.name = self.default_font
                chart.legend.font.color.rgb = RGBColor(*legend_rgb)
                chart.legend.include_in_layout = False
                logger.debug("  ✓ Legend (right)")
                
            elif chart_type in ['column', 'bar', 'line', 'area']:
                if len(chart.series) > 1:
                    chart.has_legend = True
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.font.size = Pt(10)
                    chart.legend.font.name = self.default_font
                    chart.legend.font.color.rgb = RGBColor(*legend_rgb)
                    logger.debug("  ✓ Legend (bottom)")
                else:
                    chart.has_legend = False
                    logger.debug("  ✓ Legend hidden")
            
        except Exception as e:
            logger.warning(f"⚠️  Legend error: {e}")
    
    def _apply_modern_chart_style(self, chart, chart_type: str, background_rgb: Optional[Tuple], color_config: Dict):
        """Apply styling with dynamic colors from constraints"""
        try:
            # Set chart style
            try:
                chart.chart_style = 2
            except:
                pass
            
            # Get series colors from constraints
            base_colors = [self._get_rgb(color) for color in self.color_palette]
            
            for series_idx, series in enumerate(chart.series):
                # Cycle through colors for multiple series
                series_color = base_colors[series_idx % len(base_colors)]
                
                try:
                    if chart_type == 'line':
                        series.format.line.color.rgb = RGBColor(*series_color)
                        series.format.line.width = Pt(2.5)
                        series.smooth = True
                        
                        try:
                            series.marker.style = 8
                            series.marker.size = 7
                            series.marker.format.fill.solid()
                            series.marker.format.fill.fore_color.rgb = RGBColor(*series_color)
                            series.marker.format.line.color.rgb = RGBColor(255, 255, 255)
                            series.marker.format.line.width = Pt(1.5)
                        except:
                            pass
                    
                    elif chart_type in ['column', 'bar']:
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = RGBColor(*series_color)
                    
                    elif chart_type == 'area':
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = RGBColor(*series_color)
                        series.format.line.color.rgb = RGBColor(*series_color)
                except Exception as e:
                    logger.debug(f"Series styling skipped: {e}")
            
            # Axis styling
            if chart_type != 'pie':
                # Category axis (X-axis)
                try:
                    axis_rgb = self._get_rgb(color_config.get("axis_color"))
                    cat_axis = chart.category_axis
                    cat_axis.visible = True
                    
                    cat_axis.tick_labels.font.size = Pt(self.font_size)
                    cat_axis.tick_labels.font.bold = False
                    cat_axis.tick_labels.font.name = self.default_font
                    cat_axis.tick_labels.font.color.rgb = RGBColor(*axis_rgb)
                    
                    self._force_axis_text_xml(cat_axis, axis_rgb)
                    
                    logger.info(f"  ✓ X-axis styled")
                except Exception as e:
                    logger.error(f"❌ X-axis styling failed: {e}")
                
                # Value axis (Y-axis)
                try:
                    axis_rgb = self._get_rgb(color_config.get("axis_color"))
                    val_axis = chart.value_axis
                    val_axis.visible = True
                    
                    val_axis.tick_labels.font.size = Pt(self.font_size)
                    val_axis.tick_labels.font.bold = False
                    val_axis.tick_labels.font.name = self.default_font
                    val_axis.tick_labels.font.color.rgb = RGBColor(*axis_rgb)
                    
                    self._force_axis_text_xml(val_axis, axis_rgb)
                    
                    # Gridlines
                    try:
                        if val_axis.has_major_gridlines:
                            grid_rgb = self._get_rgb(color_config.get("grid_color"))
                            val_axis.major_gridlines.format.line.color.rgb = RGBColor(*grid_rgb)
                            val_axis.major_gridlines.format.line.width = Pt(0.75)
                    except:
                        pass
                    
                    logger.info(f"  ✓ Y-axis styled")
                except Exception as e:
                    logger.error(f"❌ Y-axis styling failed: {e}")
            
            logger.info("  ✅ Chart styling applied")
            
        except Exception as e:
            logger.error(f"❌ Chart styling error: {e}")

    def _force_axis_text_xml(self, axis, color_rgb: Tuple[int, int, int]):
        """Force axis text color using XML manipulation"""
        try:
            axis_element = axis._element
            
            # Convert RGB to hex
            color_hex = '%02x%02x%02x' % color_rgb
            
            white_color_xml = f'''
            <c:txPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" 
                     xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                <a:bodyPr/>
                <a:lstStyle/>
                <a:p>
                    <a:pPr>
                        <a:defRPr sz="1200" b="0">
                            <a:solidFill>
                                <a:srgbClr val="{color_hex.upper()}"/>
                            </a:solidFill>
                            <a:latin typeface="{self.default_font}"/>
                        </a:defRPr>
                    </a:pPr>
                    <a:endParaRPr lang="en-US"/>
                </a:p>
            </c:txPr>
            '''
            
            # Remove existing txPr if present
            existing_txPr = axis_element.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}txPr')
            if existing_txPr is not None:
                axis_element.remove(existing_txPr)
            
            # Add new XML
            txPr = parse_xml(white_color_xml)
            axis_element.append(txPr)
            
            logger.debug("  ✓ Forced axis text color via XML")
            return True
        except Exception as e:
            logger.debug(f"XML axis color failed: {e}")
            return False
    
    def _get_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        """Convert hex to RGB"""
        try:
            value = hex_color.lstrip("#")
            return tuple(int(value[i:i+2], 16) for i in (0, 2, 4))
        except Exception:
            return (255, 255, 255)  # Default to white