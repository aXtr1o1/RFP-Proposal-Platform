import os
import sys
import json
import logging
from typing import Dict, Any, Tuple, Optional, List, Iterator
from concurrent.futures import ThreadPoolExecutor
import requests
from openai import OpenAI
from dotenv import load_dotenv

load_dotenv(override=True)

def _emit_stdout(text: str) -> None:
    data = (text + "\n").encode("utf-8", errors="replace")
    buffer = getattr(sys.stdout, "buffer", None)
    if buffer is not None:
        try:
            buffer.write(data)
            buffer.flush()
            return
        except Exception:
            pass
    safe_text = (text + "\n").encode("ascii", errors="replace").decode("ascii")
    print(safe_text, end="", flush=True)

from apps.api.services.supabase_service import (
    save_generated_markdown,
)
from apps.wordgenAgent.app import prompt5 as P
from apps.wordgenAgent.app.document import generate_word_from_markdown

logger = logging.getLogger("wordgen_api")


def _lang_flag(language: str) -> str:
    lang = (language or "").strip().lower()
    if lang == "arabic":
        return (
            "LANGUAGE_MODE: ARABIC (Modern Standard Arabic).\n"
            "TOP PRIORITY: Output ALL fields (title, headings, content, points, table headers/rows) in Arabic only.\n"
        )
    return (
        "LANGUAGE_MODE: ENGLISH.\n"
        "TOP PRIORITY: Output ALL fields (title, headings, content, points, table headers/rows) in English only."
    )


def _sse_event_raw(event: str, data: str) -> bytes:
    if data is None:
        data = ""
    encoded_data = json.dumps(data, ensure_ascii=False)
    return f"event: {event}\ndata: {encoded_data}\n\n".encode("utf-8")


def _sse_event_json(event: str, obj: Dict[str, Any]) -> bytes:
    return f"event: {event}\ndata: {json.dumps(obj, ensure_ascii=False)}\n\n".encode("utf-8")


class WordGenAPI:
    def __init__(self) -> None:
        api_key = os.getenv("OPENAI_API_KEY")
        if not api_key:
            raise RuntimeError("OPENAI_API_KEY is required")
        self.client = OpenAI(api_key=api_key)
        logger.info("OpenAI client initialized")

    @staticmethod
    def _clean_url(url: str) -> str:
        cleaned_url = (url or "").split("?")[0]
        logger.debug(f"Cleaned URL: {cleaned_url}")
        return cleaned_url
    
    def _get_mime_type(self, filename: str) -> str:
        """Add this new method to detect MIME type from filename"""
        ext = filename.lower().split('.')[-1]
        mime_types = {
            'pdf': 'application/pdf',
            'doc': 'application/msword',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'ppt': 'application/vnd.ms-powerpoint',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'txt': 'text/plain',
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'png': 'image/png'
        }
        return mime_types.get(ext, 'application/octet-stream')
    
    def _convert_to_pdf(self, file_bytes: bytes, filename: str) -> Tuple[bytes, str]:
        """Convert document formats to PDF"""
        import io
        import subprocess
        import tempfile
        import os
        ext = filename.lower().split('.')[-1]
        
        # If already PDF, return as-is
        if ext == 'pdf':
            logger.info(f"{filename} is already PDF")
            return file_bytes, filename
        
        # For text-based formats, create simple PDF
        if ext in ['txt', 'md', 'markdown']:
            try:
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
                
                # Decode text
                text = file_bytes.decode('utf-8', errors='ignore')
                
                # Create PDF in memory
                pdf_buffer = io.BytesIO()
                doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
                styles = getSampleStyleSheet()
                story = []
                
                # Add text paragraphs
                for line in text.split('\n'):
                    if line.strip():
                        para = Paragraph(line, styles['Normal'])
                        story.append(para)
                        story.append(Spacer(1, 12))
                
                doc.build(story)
                pdf_bytes = pdf_buffer.getvalue()
                new_filename = filename.rsplit('.', 1)[0] + '.pdf'
                logger.info(f"Converted {filename} to {new_filename}")
                return pdf_bytes, new_filename
            except Exception as e:
                logger.error(f"Failed to convert {filename} to PDF: {e}")
                raise RuntimeError(f"Cannot convert text to PDF: {str(e)}")
        
        # For DOCX/PPTX, use LibreOffice (best quality)
        if ext in ['docx', 'doc', 'pptx', 'ppt']:
            try:
                # Create temporary files
                with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{ext}') as tmp_input:
                    tmp_input.write(file_bytes)
                    tmp_input_path = tmp_input.name
                
                tmp_output_dir = tempfile.mkdtemp()
                
                # Convert using LibreOffice
                cmd = [
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', tmp_output_dir,
                    tmp_input_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
                
                if result.returncode != 0:
                    raise RuntimeError(f"LibreOffice conversion failed: {result.stderr}")
                
                # Read converted PDF
                base_name = os.path.splitext(os.path.basename(tmp_input_path))[0]
                pdf_path = os.path.join(tmp_output_dir, f"{base_name}.pdf")
                
                if not os.path.exists(pdf_path):
                    raise RuntimeError("PDF output file not created")
                
                with open(pdf_path, 'rb') as f:
                    pdf_bytes = f.read()
                
                # Cleanup
                os.unlink(tmp_input_path)
                os.unlink(pdf_path)
                os.rmdir(tmp_output_dir)
                
                new_filename = filename.rsplit('.', 1)[0] + '.pdf'
                logger.info(f"Converted {filename} to {new_filename} using LibreOffice ({len(pdf_bytes)} bytes)")
                return pdf_bytes, new_filename
                
            except FileNotFoundError:
                logger.error("LibreOffice not installed. Falling back to text extraction.")
                # Fallback: extract text and create simple PDF
                return self._convert_office_to_pdf_fallback(file_bytes, filename, ext)
            except Exception as e:
                logger.error(f"LibreOffice conversion failed: {e}. Falling back to text extraction.")
                return self._convert_office_to_pdf_fallback(file_bytes, filename, ext)
        
        raise RuntimeError(f"Unsupported file format for PDF conversion: {ext}")


    def _convert_office_to_pdf_fallback(self, file_bytes: bytes, filename: str, ext: str) -> Tuple[bytes, str]:
        """Fallback: Extract text from Office files and create simple PDF"""
        import io
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        
        text_lines = []
        
        # Extract text from DOCX
        if ext == 'docx':
            from docx import Document
            doc = Document(io.BytesIO(file_bytes))
            for para in doc.paragraphs:
                if para.text.strip():
                    text_lines.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        text_lines.append(row_text)
        
        # Extract text from PPTX
        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            for slide_num, slide in enumerate(prs.slides, 1):
                text_lines.append(f"Slide {slide_num}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_lines.append(shape.text)
        
        # Create PDF
        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        for line in text_lines:
            para = Paragraph(line.replace('<', '&lt;').replace('>', '&gt;'), styles['Normal'])
            story.append(para)
            story.append(Spacer(1, 12))
        
        doc.build(story)
        pdf_bytes = pdf_buffer.getvalue()
        new_filename = filename.rsplit('.', 1)[0] + '.pdf'
        logger.info(f"Converted {filename} to {new_filename} (text extraction fallback)")
        return pdf_bytes, new_filename

    def _download_file(self, url: str) -> bytes:
        cleaned_url = self._clean_url(url)
        logger.info(f"Starting download of PDF: {cleaned_url}")
        response = requests.get(cleaned_url, timeout=60)
        response.raise_for_status()
        logger.info(f"Finished download of PDF: {cleaned_url} ({len(response.content)} bytes)")
        return response.content

    def _download_two_files(self, rfp_url: str, supporting_url: str) -> Tuple[bytes, bytes]:
        logger.info("Beginning parallel download of two PDFs")
        with ThreadPoolExecutor() as executor:
            future_rfp = executor.submit(self._download_file, rfp_url)
            future_sup = executor.submit(self._download_file, supporting_url)
            rfp_bytes = future_rfp.result()
            sup_bytes = future_sup.result()
        logger.info("Completed parallel download of two PDFs")
        return rfp_bytes, sup_bytes

    def _upload_file_bytes_to_openai(self, file_bytes: bytes, filename: str) -> str:
        logger.info(f"Processing {filename} for OpenAI upload")
        
        # Convert to PDF
        pdf_bytes, pdf_filename = self._convert_to_pdf(file_bytes, filename)
        
        mime_type = 'application/pdf'
        logger.info(f"Uploading {pdf_filename} (MIME: {mime_type}) to OpenAI")
        
        file_obj = self.client.files.create(
            file=(pdf_filename, pdf_bytes, mime_type), 
            purpose="user_data"
        )
        logger.info(f"Uploaded {pdf_filename} as file ID: {file_obj.id}")
        return file_obj.id
    
    def _extract_filename_from_url(self, url: str) -> str:
        """Extract filename from URL, default if not found"""
        from urllib.parse import urlparse, unquote
        path = urlparse(url).path
        filename = unquote(path.split('/')[-1])
        return filename if filename else "document"
    
    def _upload_file_urls_to_openai(self, rfp_url: str, supporting_url: str) -> Tuple[str, str]:
        rfp_bytes, sup_bytes = self._download_two_files(rfp_url, supporting_url)
        logger.info("Beginning parallel upload of two files to OpenAI")
        
        # Extract filenames from URLs
        rfp_filename = self._extract_filename_from_url(rfp_url)
        sup_filename = self._extract_filename_from_url(supporting_url)
        
        with ThreadPoolExecutor() as executor:
            future_rfp = executor.submit(self._upload_file_bytes_to_openai, rfp_bytes, rfp_filename)
            future_sup = executor.submit(self._upload_file_bytes_to_openai, sup_bytes, sup_filename)
            rfpf_id = future_rfp.result()
            supf_id = future_sup.result()
        
        logger.info(f"Completed parallel upload of files with IDs: rfp={rfpf_id}, supporting={supf_id}")
        return rfpf_id, supf_id

    def generate_complete_proposal(
        self,
        uuid: str,
        gen_id: str,
        rfp_url: str,
        supporting_url: str,
        user_config: str = "",
        doc_config: Optional[Dict[str, Any]] = None,
        language: str = "english",
        outline: Optional[str] = None,
    ) -> Iterator[bytes]:
        try:
            yield _sse_event_json("stage", {"stage": "starting"})
            yield _sse_event_json("stage", {"stage": "uploading_files"})
            rfp_id, sup_id = self._upload_file_urls_to_openai(rfp_url, supporting_url)

            rfp_label = "RFP/BRD: requirements, evaluation criteria, project details, and timelines"
            supporting_label = "Supporting: company profile, portfolio, capabilities, certifications, differentiators"
            system_prompts = P.system_prompts
            lang_block = _lang_flag(language)
            user_cfg_notes = user_config if isinstance(user_config, str) else ""

            additive_block = P.build_task_instructions_with_config(
                language=language,
                user_config_json=(user_config if isinstance(user_config, str) else "null"),
                rfp_label=rfp_label,
                supporting_label=supporting_label,
                user_config_notes=user_cfg_notes,
            )
            task_instructions = f"\nIMPORTANT: The proposal must follow this structure:\n{additive_block}"
            yield _sse_event_json("stage", {"stage": "prompting_model"})
            logger.info(doc_config)
            logger.info("Calling OpenAI Responses API…")
            response = self.client.responses.create(
                model=P.MODEL,
                max_output_tokens=18000,
                input=[{
                    "role": "user",
                    "content": [
                        {"type": "input_text", "text": lang_block},
                        {"type": "input_file", "file_id": rfp_id},
                        {"type": "input_file", "file_id": sup_id},
                        {"type": "input_text", "text": user_cfg_notes},
                        {"type": "input_text", "text": system_prompts},
                        {"type": "input_text", "text": task_instructions},
                    ],
                }],
                reasoning={"effort": "minimal"},
                stream=True,
            )

            buffer_chunks: List[str] = []
            line_buffer = ""

            for event in response:
                et = getattr(event, "type", "")
                if et == "response.output_text.delta":
                    delta = getattr(event, "delta", "")
                    if not delta:
                        continue
                    buffer_chunks.append(delta)
                    yield _sse_event_raw("chunk", delta)

                    line_buffer += delta
                    while "\n" in line_buffer:
                        line, line_buffer = line_buffer.split("\n", 1)
                        _emit_stdout(line)

                elif et == "response.error":
                    err_msg = getattr(event, "error", "stream error")
                    logger.error(f"OpenAI stream error: {err_msg}")
                    yield _sse_event_json("error", {"message": str(err_msg)})
                    return

                elif et == "response.completed":
                    break

            if line_buffer:
                _emit_stdout(line_buffer)

            full_markdown = "".join(buffer_chunks)
            _emit_stdout(full_markdown)
            yield _sse_event_json("stage", {"stage": "saving_markdown"})

            saved_ok = False
            try:
                saved_ok = save_generated_markdown(uuid, gen_id, full_markdown)
            except Exception as e:
                logger.exception("Failed to save generated markdown")
                yield _sse_event_json("error", {"message": f"save error: {str(e)}"})
            yield _sse_event_json("stage", {"stage": "building_word"})
            try:
                _ = generate_word_from_markdown(
                    uuid=uuid,
                    gen_id=gen_id,
                    markdown=full_markdown,
                    doc_config=doc_config,
                    language=(language or "english").lower(),
                )
            except Exception as e:
                logger.exception("Word generation/upload failed")
                yield _sse_event_json("error", {"message": f"word build error: {str(e)}"})

            yield _sse_event_json("done", {"status": "saved" if saved_ok else "not_saved"})

        except Exception as e:
            logger.exception("generate_complete_proposal failed")
            yield _sse_event_json("error", {"message": str(e)})


wordgen_api = WordGenAPI()
