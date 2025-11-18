from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from typing import Dict, Tuple, Optional, List
import logging

logger = logging.getLogger("chart_service")

class ChartService:
    """Service for creating native PowerPoint charts with dynamic data format support"""
    
    def __init__(self, template_id: str = "arweqah"):
        self.template_id = template_id
    
    def add_native_chart(
        self,
        slide,
        chart_data: Dict,
        position: Dict,
        size: Dict,
        background_rgb: Optional[Tuple[int, int, int]] = None
    ):
        """Add native PowerPoint chart - supports both old and new data formats"""
        try:
            chart_type = chart_data.get('chart_type', 'column')
            title = chart_data.get('title', '')
            
            # CRITICAL: Extract data dynamically
            categories, series_list = self._extract_chart_data(chart_data)
            
            if not categories or not series_list:
                logger.error(f"❌ Chart data extraction failed!")
                logger.error(f"   Raw data: {chart_data}")
                return None
            
            if len(categories) == 0 or len(series_list) == 0:
                logger.error(f"❌ Chart data empty! categories={len(categories)}, series={len(series_list)}")
                return None
            
            logger.info(f"✓ Creating {chart_type} chart")
            logger.info(f"   Categories: {categories}")
            logger.info(f"   Series count: {len(series_list)}")
            for idx, s in enumerate(series_list):
                logger.info(f"   Series {idx + 1}: {s['name']} = {s['values']}")
            
            # Map chart types
            chart_type_map = {
                'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'bar': XL_CHART_TYPE.BAR_CLUSTERED,
                'line': XL_CHART_TYPE.LINE_MARKERS,
                'pie': XL_CHART_TYPE.PIE,
                'area': XL_CHART_TYPE.AREA
            }
            
            xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            
            # Create chart data object
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = categories
            
            # Add all series
            for series_info in series_list:
                chart_data_obj.add_series(series_info['name'], series_info['values'])
            
            # Add chart to slide
            x = Inches(position['left'])
            y = Inches(position['top'])
            cx = Inches(size['width'])
            cy = Inches(size['height'])
            
            graphic_frame = slide.shapes.add_chart(
                xl_chart_type, x, y, cx, cy, chart_data_obj
            )
            
            chart = graphic_frame.chart

            # Get axis labels and unit
            x_axis_label = chart_data.get('x_axis_label', '')
            y_axis_label = chart_data.get('y_axis_label', '')
            unit = chart_data.get('unit', '')

            # WHITE COLORS FOR DARK BACKGROUNDS
            color_config = {
                "font_color": chart_data.get("font_color", "#FFFFFF"),
                "data_label_color": chart_data.get("data_label_color", "#FFFFFF"),
                "axis_color": chart_data.get("axis_color", "#FFFFFF"),
                "axis_label_color": chart_data.get("axis_label_color", "#FFFFFF"),
                "grid_color": chart_data.get("grid_color", "#5C6F7A"),
                "legend_font_color": chart_data.get("legend_font_color", "#FFFFFF"),
                "plot_background_color": chart_data.get("plot_background_color"),
                "series_color": chart_data.get("series_color", "#F9D462")
            }
            
            # Apply modern styling with WHITE axis text
            self._apply_modern_chart_style(chart, chart_type, background_rgb, color_config)
            
            # Add data labels to all series
            self._add_data_labels(chart, chart_type, unit, color_config)
            
            # Configure legend
            self._configure_legend(chart, chart_type, color_config)
            
            # Add axis labels (WHITE text)
            self._add_axis_labels(chart, chart_type, x_axis_label, y_axis_label, unit, color_config)
            
            # Add title if provided
            if title:
                chart.has_title = True
                chart.chart_title.text_frame.text = title
                chart.chart_title.text_frame.paragraphs[0].font.size = Pt(16)
                chart.chart_title.text_frame.paragraphs[0].font.bold = True
                title_color = self._get_rgb(color_config.get("font_color"), (255, 255, 255))
                chart.chart_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(*title_color)
            
            logger.info(f"✅ Chart created successfully: {title} ({chart_type}, {len(series_list)} series)")
            return chart
            
        except Exception as e:
            logger.error(f"❌ Chart creation error: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _extract_chart_data(self, chart_data: Dict) -> Tuple[List, List[Dict]]:
        """
        Dynamically extract chart data from multiple formats
        Returns: (categories, series_list)
        """
        categories = []
        series_list = []
        
        # FORMAT 1: NEW - categories + series (PREFERRED)
        if 'series' in chart_data and chart_data.get('series'):
            raw_series = chart_data.get('series', [])
            categories = chart_data.get('categories', [])
            
            if isinstance(raw_series, list) and len(raw_series) > 0:
                for series_item in raw_series:
                    if isinstance(series_item, dict):
                        name = series_item.get('name', 'Series')
                        values = series_item.get('values', [])
                        if values:
                            series_list.append({'name': name, 'values': values})
                
                if categories and series_list:
                    logger.info(f"   ✓ Using NEW format: {len(categories)} categories, {len(series_list)} series")
                    return categories, series_list
        
        # FORMAT 2: OLD - labels + values (FALLBACK)
        if 'labels' in chart_data or 'values' in chart_data:
            categories = chart_data.get('labels', []) or chart_data.get('categories', [])
            values = chart_data.get('values', [])
            series_name = chart_data.get('series_name', 'Values')
            
            if categories and values:
                series_list = [{'name': series_name, 'values': values}]
                logger.info(f"   ✓ Using OLD format: {len(categories)} labels, 1 series")
                return categories, series_list
        
        # FORMAT 3: LAST RESORT - categories only (try to extract from chart_data)
        if 'categories' in chart_data and chart_data.get('categories'):
            categories = chart_data.get('categories', [])
            # Try to find any numeric data
            for key in ['values', 'data', 'y']:
                if key in chart_data and chart_data[key]:
                    values = chart_data[key]
                    series_list = [{'name': 'Values', 'values': values}]
                    logger.info(f"   ✓ Using FALLBACK format: found data in '{key}'")
                    return categories, series_list
        
        logger.error(f"   ❌ Could not extract chart data from any format")
        return [], []
        
    def _add_axis_labels(self, chart, chart_type: str, x_axis_label: str, y_axis_label: str, unit: str, color_config: Dict):
        """Add axis labels with WHITE text for visibility on dark backgrounds"""
        try:
            if chart_type == 'pie':
                return
            
            # WHITE color for axis labels
            axis_label_rgb = self._get_rgb(color_config.get("axis_label_color"), (255, 255, 255))

            # Category axis (X-axis) label
            if x_axis_label:
                try:
                    cat_axis = chart.category_axis
                    cat_axis.has_title = True
                    cat_axis.axis_title.text_frame.text = x_axis_label
                    
                    # CRITICAL: Make title text WHITE and visible
                    for paragraph in cat_axis.axis_title.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.bold = True
                        paragraph.font.name = 'Calibri'
                        paragraph.font.color.rgb = RGBColor(*axis_label_rgb)
                    
                    logger.info(f"  ✓ X-axis label (WHITE): {x_axis_label}")
                except Exception as e:
                    logger.error(f"❌ X-axis label error: {e}")
                    import traceback
                    traceback.print_exc()
            
            # Value axis (Y-axis) label
            if y_axis_label or unit:
                try:
                    val_axis = chart.value_axis
                    val_axis.has_title = True
                    
                    # Combine Y-axis label with unit
                    if y_axis_label and unit:
                        axis_title = f"{y_axis_label} ({unit})"
                    elif y_axis_label:
                        axis_title = y_axis_label
                    else:
                        axis_title = unit
                    
                    val_axis.axis_title.text_frame.text = axis_title
                    
                    # CRITICAL: Make title text WHITE and visible
                    for paragraph in val_axis.axis_title.text_frame.paragraphs:
                        paragraph.font.size = Pt(14)
                        paragraph.font.bold = True
                        paragraph.font.name = 'Calibri'
                        paragraph.font.color.rgb = RGBColor(*axis_label_rgb)
                    
                    logger.info(f"  ✓ Y-axis label (WHITE): {axis_title}")
                except Exception as e:
                    logger.error(f"❌ Y-axis label error: {e}")
                    import traceback
                    traceback.print_exc()
                    
        except Exception as e:
            logger.error(f"❌ Axis labels error: {e}")
            import traceback
            traceback.print_exc()

    def _add_data_labels(self, chart, chart_type: str, unit: str, color_config: Dict):
        """Add data labels to chart series with WHITE text"""
        try:
            for series_idx, series in enumerate(chart.series):
                series.has_data_labels = True
                data_labels = series.data_labels
                
                data_label_rgb = self._get_rgb(color_config.get("data_label_color"), (255, 255, 255))

                if chart_type == 'pie':
                    data_labels.number_format = '0%'
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    data_labels.font.size = Pt(11)
                    data_labels.font.bold = True
                    # Pie labels outside - use dark color for visibility
                    pie_label_rgb = self._get_rgb(color_config.get("data_label_color"), (13, 32, 38))
                    data_labels.font.color.rgb = RGBColor(*pie_label_rgb)
                    
                elif chart_type in ['column', 'bar']:
                    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
                    data_labels.font.size = Pt(11)
                    data_labels.font.bold = True
                    data_labels.font.color.rgb = RGBColor(*data_label_rgb)
                    
                    if unit:
                        if unit == '%':
                            data_labels.number_format = '0"%"'
                        elif unit == '$':
                            data_labels.number_format = '"$"0'
                        else:
                            data_labels.number_format = f'0" {unit}"'
                    
                elif chart_type == 'line':
                    data_labels.position = XL_DATA_LABEL_POSITION.ABOVE
                    data_labels.font.size = Pt(10)
                    data_labels.font.color.rgb = RGBColor(*data_label_rgb)
                    
                    if unit:
                        if unit == '%':
                            data_labels.number_format = '0"%"'
                        elif unit == '$':
                            data_labels.number_format = '"$"0'
                        else:
                            data_labels.number_format = f'0" {unit}"'
                
                logger.debug(f"  ✓ Data labels added to series {series_idx + 1}")
                
        except Exception as e:
            logger.warning(f"⚠️  Data labels error: {e}")

    def _configure_legend(self, chart, chart_type: str, color_config: Dict):
        """Configure chart legend with WHITE text"""
        try:
            legend_rgb = self._get_rgb(color_config.get("legend_font_color"), (255, 255, 255))
            
            if chart_type == 'pie':
                chart.has_legend = True
                chart.legend.position = XL_LEGEND_POSITION.RIGHT
                chart.legend.font.size = Pt(11)
                chart.legend.font.color.rgb = RGBColor(*legend_rgb)
                chart.legend.include_in_layout = False
                logger.debug("  ✓ Legend (right)")
                
            elif chart_type in ['column', 'bar', 'line', 'area']:
                if len(chart.series) > 1:
                    chart.has_legend = True
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.font.size = Pt(10)
                    chart.legend.font.color.rgb = RGBColor(*legend_rgb)
                    logger.debug("  ✓ Legend (bottom)")
                else:
                    chart.has_legend = False
                    logger.debug("  ✓ Legend hidden")
            
        except Exception as e:
            logger.warning(f"⚠️  Legend error: {e}")
    
    def _apply_modern_chart_style(self, chart, chart_type: str, background_rgb: Optional[Tuple], color_config: Dict):
        """Apply styling with WHITE axis text for dark backgrounds"""
        try:
            # Set chart style
            try:
                chart.chart_style = 2
            except:
                pass
            
            # Series styling - use different colors for multiple series
            base_colors = [
                (249, 212, 98),   # Yellow #F9D462
                (177, 216, 190),  # Green #B1D8BE
                (224, 144, 89),   # Orange #E09059
                (198, 195, 190),  # Gray #C6C3BE
                (255, 252, 236)   # Cream #FFFCEC
            ]
            
            for series_idx, series in enumerate(chart.series):
                # Cycle through colors for multiple series
                series_color = base_colors[series_idx % len(base_colors)]
                
                try:
                    if chart_type == 'line':
                        series.format.line.color.rgb = RGBColor(*series_color)
                        series.format.line.width = Pt(2.5)
                        series.smooth = True
                        
                        try:
                            series.marker.style = 8
                            series.marker.size = 7
                            series.marker.format.fill.solid()
                            series.marker.format.fill.fore_color.rgb = RGBColor(*series_color)
                            series.marker.format.line.color.rgb = RGBColor(255, 255, 255)
                            series.marker.format.line.width = Pt(1.5)
                        except:
                            pass
                    
                    elif chart_type in ['column', 'bar']:
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = RGBColor(*series_color)
                    
                    elif chart_type == 'area':
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = RGBColor(*series_color)
                        series.format.line.color.rgb = RGBColor(*series_color)
                except Exception as e:
                    logger.debug(f"Series styling skipped: {e}")
            
            # CRITICAL: WHITE AXIS TEXT (X and Y) - THIS IS THE MOST IMPORTANT PART
            if chart_type != 'pie':
                # Category axis (X-axis) - WHITE TEXT
                try:
                    axis_rgb = self._get_rgb(color_config.get("axis_color"), (255, 255, 255))
                    cat_axis = chart.category_axis
                    cat_axis.visible = True
                    
                    # Method 1: Try standard approach
                    cat_axis.tick_labels.font.size = Pt(12)
                    cat_axis.tick_labels.font.bold = False
                    cat_axis.tick_labels.font.name = 'Calibri'
                    cat_axis.tick_labels.font.color.rgb = RGBColor(*axis_rgb)
                    
                    # Method 2: Force with XML (more reliable)
                    self._force_white_axis_text_xml(cat_axis)
                    
                    logger.info(f"  ✓ X-axis text: WHITE (forced via XML)")
                except Exception as e:
                    logger.error(f"❌ X-axis styling failed: {e}")
                    import traceback
                    traceback.print_exc()
                
                # Value axis (Y-axis) - WHITE TEXT
                try:
                    axis_rgb = self._get_rgb(color_config.get("axis_color"), (255, 255, 255))
                    val_axis = chart.value_axis
                    val_axis.visible = True
                    
                    # Method 1: Try standard approach
                    val_axis.tick_labels.font.size = Pt(12)
                    val_axis.tick_labels.font.bold = False
                    val_axis.tick_labels.font.name = 'Calibri'
                    val_axis.tick_labels.font.color.rgb = RGBColor(*axis_rgb)
                    
                    # Method 2: Force with XML (more reliable)
                    self._force_white_axis_text_xml(val_axis)
                    
                    # Gridlines
                    try:
                        if val_axis.has_major_gridlines:
                            grid_rgb = self._get_rgb(color_config.get("grid_color"), (92, 111, 122))
                            val_axis.major_gridlines.format.line.color.rgb = RGBColor(*grid_rgb)
                            val_axis.major_gridlines.format.line.width = Pt(0.75)
                    except:
                        pass
                    
                    logger.info(f"  ✓ Y-axis text: WHITE (forced via XML)")
                except Exception as e:
                    logger.error(f"❌ Y-axis styling failed: {e}")
                    import traceback
                    traceback.print_exc()
            
            logger.info("  ✅ Chart styling applied")
            
        except Exception as e:
            logger.error(f"❌ Chart styling error: {e}")
            import traceback
            traceback.print_exc()

    def _force_white_axis_text_xml(self, axis):
        """Force WHITE axis text using XML manipulation (fallback method)"""
        try:
            # Access the axis XML element
            axis_element = axis._element
            
            # Create white color XML
            white_color_xml = '''
            <c:txPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart" 
                     xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
                <a:bodyPr/>
                <a:lstStyle/>
                <a:p>
                    <a:pPr>
                        <a:defRPr sz="1200" b="0">
                            <a:solidFill>
                                <a:srgbClr val="FFFFFF"/>
                            </a:solidFill>
                            <a:latin typeface="Calibri"/>
                        </a:defRPr>
                    </a:pPr>
                    <a:endParaRPr lang="en-US"/>
                </a:p>
            </c:txPr>
            '''
            
            # Remove existing txPr if present
            existing_txPr = axis_element.find('{http://schemas.openxmlformats.org/drawingml/2006/chart}txPr')
            if existing_txPr is not None:
                axis_element.remove(existing_txPr)
            
            # Add new white color XML
            txPr = parse_xml(white_color_xml)
            axis_element.append(txPr)
            
            logger.info("  ✓ Forced WHITE axis text via XML")
            return True
        except Exception as e:
            logger.debug(f"XML axis color failed: {e}")
            return False
    
    def _get_rgb(self, hex_color: Optional[str], fallback: Optional[Tuple[int, int, int]] = None) -> Tuple[int, int, int]:
        """Convert hex to RGB with WHITE fallback"""
        if not hex_color:
            return fallback or (255, 255, 255)
        try:
            value = hex_color.lstrip("#")
            return tuple(int(value[i:i+2], 16) for i in (0, 2, 4))
        except Exception:
            return fallback or (255, 255, 255)