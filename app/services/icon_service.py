import json
import logging
from io import BytesIO
from typing import Optional, Dict, List
from pathlib import Path
from functools import lru_cache

import cairosvg
from config import settings

logger = logging.getLogger("icon_service")

# FIXED: Cache size limit to prevent memory leak
MAX_CACHE_SIZE = 100  # Max 100 rendered icons in memory


class IconService:
    """
    Icon service for handling presentation icons
    Supports SVG icons from assets with inline text rendering
    """
    
    def __init__(self, template_id: str = "standard"):
        """
        Initialize icon service with template awareness
        """
        self.template_id = template_id
        
        # Load icons data
        icons_path = Path(settings.ASSETS_DIR) / "icons.json"
        if not icons_path.exists():
            raise FileNotFoundError(f"Icons file not found: {icons_path}")
        
        try:
            with open(icons_path, 'r', encoding='utf-8') as f:
                self.icons_data = json.load(f)
            
            if "icons" not in self.icons_data:
                raise ValueError("Invalid icons.json structure: missing 'icons' key")
            
            logger.info(f"Loaded {len(self.icons_data['icons'])} icons from {icons_path}")
        except json.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON in icons.json: {e}")
        
        # FIXED: Load theme from specified template (not hardcoded 'standard')
        theme_path = Path(settings.TEMPLATES_DIR) / template_id / "theme.json"
        
        if not theme_path.exists():
            logger.warning(f" Theme file not found: {theme_path}, using default icon mapping")
            self.theme = {"icons": {"keyword_to_icon_map": {}}}
        else:
            try:
                with open(theme_path, 'r', encoding='utf-8') as f:
                    self.theme = json.load(f)
                logger.info(f"Loaded theme from {theme_path}")
            except json.JSONDecodeError as e:
                logger.error(f"Invalid theme.json: {e}")
                self.theme = {"icons": {"keyword_to_icon_map": {}}}
        
        # Get icon mapping
        self.icon_mapping = self.theme.get("icons", {}).get("keyword_to_icon_map", {})
        
        # FIXED: Cache with size limit (LRU eviction)
        self.cache: Dict[str, bytes] = {}
        self.cache_order: List[str] = []  # Track access order for LRU
        
        # Unicode fallback mapping
        self.unicode_fallback = {
            'briefcase': '💼',
            'rocket-launch': '🚀',
            'target': '🎯',
            'bullseye': '🎯',
            'trophy': '🏆',
            'package': '📦',
            'file-text': '📄',
            'certificate': '📜',
            'list-checks': '✅',
            'users-three': '👥',
            'graduation-cap': '🎓',
            'steps': '📊',
            'flow-arrow': '➡️',
            'gear': '⚙️',
            'calendar-check': '📅',
            'clock': '⏰',
            'currency-dollar': '💵',
            'shield-check': '🛡️',
            'check-circle': '✅',
            'circle': '⭕',
            'default': '●'
        }
        
        logger.info(f"IconService initialized (template: {template_id}, mappings: {len(self.icon_mapping)})")
    
    def get_icon(self, name: str) -> Optional[Dict]:
        """
        Get icon by exact name
        
        Args:
            name: Icon identifier
            
        Returns:
            Optional[Dict]: Icon data or None
        """
        if not name:
            return None
        
        for icon in self.icons_data['icons']:
            if icon.get('name') == name:
                return icon
        
        logger.debug(f"Icon not found: {name}")
        return None
    
    def get_unicode_icon(self, icon_name: str) -> str:
        """
        Get Unicode character for icon (fallback for inline text rendering)
        
        Args:
            icon_name: Icon identifier (e.g., 'briefcase', 'package')
            
        Returns:
            str: Unicode character or default bullet
        """
        if not icon_name:
            return '●'
        
        icon_name = icon_name.lower().strip()
        
        # Check Unicode fallback mapping
        if icon_name in self.unicode_fallback:
            return self.unicode_fallback[icon_name]
        
        # Check if icon exists in SVG library
        if self.get_icon(icon_name):
            return self.unicode_fallback.get(icon_name, '●')
        
        # Keyword-based fallback
        for keyword, mapped_icon in self.icon_mapping.items():
            if keyword in icon_name:
                return self.unicode_fallback.get(mapped_icon, '●')
        
        return '●'
    
    def search_by_keyword(self, keyword: str) -> str:
        """
        Search icon by keyword in tags
        
        Args:
            keyword: Search keyword
            
        Returns:
            str: Icon name or 'circle' as fallback
        """
        if not keyword:
            return 'circle'
        
        keyword = keyword.lower()
        
        # First check theme mapping
        if keyword in self.icon_mapping:
            return self.icon_mapping[keyword]
        
        # Then search in icon tags
        for icon in self.icons_data['icons']:
            tags = icon.get('tags', '').lower()
            if keyword in tags or keyword in icon.get('name', ''):
                return icon['name']
        
        logger.debug(f"No icon found for keyword: {keyword}")
        return 'circle'
    
    def auto_select_icon(self, title: str, content: str = "") -> str:
        """
        Intelligently select icon based on title and content
        Uses comprehensive keyword mapping from theme.json
        
        Args:
            title: Slide title
            content: Slide content
            
        Returns:
            str: Icon name
        """
        if not title:
            return 'circle'
        
        text = f"{title} {content}".lower()
        
        # Check each keyword in theme mapping
        for keyword, icon_name in self.icon_mapping.items():
            if keyword in text:
                # Verify icon exists
                if self.get_icon(icon_name):
                    return icon_name
        
        # Fallback: search by first word of title
        first_word = title.split()[0].lower() if title else ""
        if first_word in self.icon_mapping:
            return self.icon_mapping[first_word]
        
        # Check intelligent category keywords from theme
        intelligent = self.theme.get('icons', {}).get('intelligent_mapping', {})
        
        for category, keywords in intelligent.items():
            if any(kw in text for kw in keywords):
                # Map category to default icon
                category_icon_map = {
                    'strategy_keywords': 'target',
                    'finance_keywords': 'currency-dollar',
                    'growth_keywords': 'chart-line-up',
                    'technology_keywords': 'cpu',
                    'people_keywords': 'users-three',
                    'timeline_keywords': 'clock',
                    'security_keywords': 'shield-check',
                    'innovation_keywords': 'lightbulb',
                    'problem_keywords': 'warning-circle',
                    'solution_keywords': 'check-circle',
                    'goal_keywords': 'target',
                    'data_keywords': 'chart-pie-slice'
                }
                
                icon_name = category_icon_map.get(category, 'circle')
                if self.get_icon(icon_name):
                    return icon_name
        
        # Ultimate fallback
        return 'circle'
    
    def render_to_png(
        self, 
        icon_name: str, 
        size: int, 
        color: str
    ) -> Optional[BytesIO]:
        """
        Convert SVG icon to PNG with caching and size limit
        
        Args:
            icon_name: Icon identifier
            size: Size in pixels
            color: Hex color code (e.g., "#FFFFFF")
            
        Returns:
            Optional[BytesIO]: PNG image data or None
        """
        if not icon_name:
            logger.warning(" Empty icon name provided")
            return None
        
        cache_key = f"{icon_name}_{size}_{color}"
        
        # Check cache
        if cache_key in self.cache:
            logger.debug(f"Cache hit: {cache_key}")
            # Update access order (LRU)
            self.cache_order.remove(cache_key)
            self.cache_order.append(cache_key)
            return BytesIO(self.cache[cache_key])
        
        # Get icon
        icon = self.get_icon(icon_name)
        if not icon:
            logger.warning(f" Icon not found: {icon_name}, trying fallback")
            # Try fallback
            icon = self.get_icon('circle')
            if not icon:
                logger.error("Fallback icon 'circle' not found")
                return None
        
        # Replace color in SVG
        svg_content = icon.get('content', '')
        if not svg_content:
            logger.error(f"Empty SVG content for icon: {icon_name}")
            return None
        
        svg_content = svg_content.replace('currentColor', color)
        svg_content = svg_content.replace('fill="currentColor"', f'fill="{color}"')
        
        # Convert to PNG with high DPI
        try:
            png_data = cairosvg.svg2png(
                bytestring=svg_content.encode('utf-8'),
                output_width=size * 4,  # 4x for ultra-high quality
                output_height=size * 4
            )
            
            # FIXED: Cache with size limit (LRU eviction)
            if len(self.cache) >= MAX_CACHE_SIZE:
                # Remove oldest entry
                oldest_key = self.cache_order.pop(0)
                del self.cache[oldest_key]
                logger.debug(f"Evicted from cache: {oldest_key}")
            
            # Add to cache
            self.cache[cache_key] = png_data
            self.cache_order.append(cache_key)
            
            logger.debug(f"Rendered and cached: {cache_key} ({len(png_data)} bytes)")
            return BytesIO(png_data)
            
        except Exception as e:
            logger.error(f"SVG render error for {icon_name}: {e}")
            return None
    
    def render_inline_icon(
        self, 
        icon_name: str, 
        size: int, 
        color: str
    ) -> Optional[BytesIO]:
        """
        Render small icon for inline text use
        
        Args:
            icon_name: Icon identifier
            size: Size in pixels (typically 20-24 for inline)
            color: Hex color code
            
        Returns:
            Optional[BytesIO]: PNG image data or None
        """
        return self.render_to_png(icon_name, size, color)
    
    def get_icon_suggestions(self, text: str, limit: int = 5) -> List[str]:
        """
        Get multiple icon suggestions for given text
        
        Args:
            text: Text to analyze
            limit: Maximum number of suggestions
            
        Returns:
            List[str]: List of icon names
        """
        if not text:
            return ['circle']
        
        text_lower = text.lower()
        suggestions = []
        
        # Check keyword mapping
        for keyword, icon_name in self.icon_mapping.items():
            if keyword in text_lower and icon_name not in suggestions:
                suggestions.append(icon_name)
                if len(suggestions) >= limit:
                    break
        
        # Fill with search results if needed
        if len(suggestions) < limit:
            words = text_lower.split()
            for word in words:
                if len(word) > 3:  # Skip short words
                    icon_name = self.search_by_keyword(word)
                    if icon_name not in suggestions:
                        suggestions.append(icon_name)
                        if len(suggestions) >= limit:
                            break
        
        return suggestions if suggestions else ['circle']
    
    def add_icon_to_text(
        self, 
        text_frame, 
        icon_name: str, 
        text: str, 
        size: int, 
        color: str
    ):
        """
        Add icon inline with text in a paragraph
        
        Args:
            text_frame: PowerPoint text frame
            icon_name: Icon identifier
            text: Text to display after icon
            size: Font/icon size in points
            color: Hex color for icon and text
        """
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        
        # Get Unicode fallback
        icon_char = self.get_unicode_icon(icon_name)
        
        # Create paragraph with icon + text
        p = text_frame.add_paragraph()
        p.text = f"{icon_char}  {text}"
        p.font.size = Pt(size)
        p.font.bold = True
        
        # Parse color
        if color.startswith('#'):
            try:
                r = int(color[1:3], 16)
                g = int(color[3:5], 16)
                b = int(color[5:7], 16)
                p.font.color.rgb = RGBColor(r, g, b)
            except ValueError as e:
                logger.warning(f" Invalid color format: {color}, error: {e}")
        
        return p
    
    def clear_cache(self) -> None:
        """Clear the icon cache (useful for memory management)"""
        cache_size = len(self.cache)
        self.cache.clear()
        self.cache_order.clear()
        logger.info(f"Cleared icon cache ({cache_size} entries)")
    
    def get_cache_stats(self) -> Dict[str, int]:
        """Get cache statistics"""
        return {
            "size": len(self.cache),
            "max_size": MAX_CACHE_SIZE,
            "memory_bytes": sum(len(data) for data in self.cache.values())
        }
