import json
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from typing import Dict, Union
from pathlib import Path

from config import settings

class TableService:
    """Table generation with styling from constraints.json."""

    def __init__(self, template_id="arweqah"):
        constraints = self._load_constraints(template_id)
        tbl = constraints["table"]

        self.header_color = self._hex_to_rgb(tbl["header_bg"])
        self.header_text_color = self._hex_to_rgb(tbl["header_text"])
        self.alt_row_color = self._hex_to_rgb(tbl["alternate_row_bg"])
        self.text_color = self._hex_to_rgb(tbl["text"])
        self.border_color = self._hex_to_rgb(tbl["border"])
        self.border_width = tbl.get("border_width", 1)
        self.rounded_corners = tbl.get("rounded_corners", True)
        self.corner_radius = tbl.get("corner_radius", 8000)
        self.max_rows_per_slide = 8

    def add_table(self, slide, table_data: Union[Dict, object], position: Dict, size: Dict):
        """Add table using template constraints."""
        try:
            if isinstance(table_data, dict):
                headers = table_data.get('headers', [])
                rows = table_data.get('rows', [])
            else:
                headers = getattr(table_data, 'headers', [])
                rows = getattr(table_data, 'rows', [])

            if not headers or not rows:
                print("  Table missing headers or rows")
                return None
            headers = [str(h).strip() for h in headers]
            validated_rows = []
            for row in rows[:self.max_rows_per_slide]:
                clean_row = [str(cell).strip() if cell else '' for cell in row]
                if len(clean_row) < len(headers):
                    clean_row.extend([''] * (len(headers) - len(clean_row)))
                elif len(clean_row) > len(headers):
                    clean_row = clean_row[:len(headers)]
                validated_rows.append(clean_row)

            num_cols = len(headers)
            num_rows = len(validated_rows) + 1
            left, top = Inches(position['left']), Inches(position['top'])
            width, height = Inches(size['width']), Inches(size['height'])

            # Rounded rectangle bg
            bg_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left - Inches(0.03), top - Inches(0.03),
                width + Inches(0.06), height + Inches(0.06)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
            bg_shape.line.color.rgb = RGBColor(*self.border_color)
            bg_shape.line.width = Pt(self.border_width)
            try:
                bg_shape.adjustments[0] = 0.08
            except:
                pass

            # Table
            table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
            table = table_shape.table
            col_width_emu = int(width / num_cols)
            for col_idx in range(num_cols):
                table.columns[col_idx].width = col_width_emu

            # Header
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.word_wrap = True
                cell.text_frame.margin_left = Inches(0.1)
                cell.text_frame.margin_right = Inches(0.1)
                cell.text_frame.margin_top = Inches(0.08)
                cell.text_frame.margin_bottom = Inches(0.08)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*self.header_color)
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.bold = True
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = RGBColor(*self.header_text_color)
                paragraph.alignment = PP_ALIGN.CENTER
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                self._remove_cell_borders(cell)

            # Body
            for row_idx, row_data in enumerate(validated_rows, start=1):
                for col_idx, cell_value in enumerate(row_data):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = cell_value
                    cell.text_frame.word_wrap = True
                    cell.text_frame.margin_left = Inches(0.1)
                    cell.text_frame.margin_right = Inches(0.1)
                    cell.text_frame.margin_top = Inches(0.05)
                    cell.text_frame.margin_bottom = Inches(0.05)
                    if row_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*self.alt_row_color)
                    else:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    paragraph = cell.text_frame.paragraphs[0]
                    paragraph.font.size = Pt(11)
                    paragraph.font.color.rgb = RGBColor(*self.text_color)
                    paragraph.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER
                    paragraph.line_spacing = 1.2
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    self._add_subtle_inner_borders(cell, row_idx, col_idx, num_rows, num_cols)
            return table
        except Exception as e:
            print(f"  Table error: {e}")
            import traceback
            traceback.print_exc()
            return None

    def _remove_cell_borders(self, cell):
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                existing = tcPr.find(f'.//{{{tcPr.nsmap.get("a", "http://schemas.openxmlformats.org/drawingml/2006/main")}}}{border_name}')
                if existing is not None:
                    tcPr.remove(existing)
                ln = OxmlElement(f'a:{border_name}')
                ln.set('w', '0')
                noFill = OxmlElement('a:noFill')
                ln.append(noFill)
                tcPr.append(ln)
        except:
            pass

    def _add_subtle_inner_borders(self, cell, row_idx, col_idx, num_rows, num_cols):
        try:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            if col_idx < num_cols - 1:
                self._set_border(tcPr, 'lnR', self.border_color, width=6350)
            if row_idx < num_rows - 1:
                self._set_border(tcPr, 'lnB', self.border_color, width=6350)
        except:
            pass

    def _set_border(self, tcPr, border_name: str, color_rgb: tuple, width: int = 12700):
        try:
            ln = OxmlElement(f'a:{border_name}')
            ln.set('w', str(width))
            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            srgbClr.set('val', '%02x%02x%02x' % color_rgb)
            solidFill.append(srgbClr)
            ln.append(solidFill)
            tcPr.append(ln)
        except:
            pass

    def _load_constraints(self, template_id):
        path = Path(settings.TEMPLATES_DIR) / template_id / "constraints.json"
        with open(path, "r", encoding="utf-8") as file:
            return json.load(file)

    def _hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
