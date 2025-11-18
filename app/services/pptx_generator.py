import logging
import re
import json
from datetime import datetime
from pathlib import Path
from typing import Tuple, Optional, Dict, Any, List
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml import parse_xml
from pptx.util import Inches, Pt

from config import settings
from models.presentation import PresentationData, SlideContent, TableData, BulletPoint
from services.chart_service import ChartService
from services.icon_service import IconService
from utils.content_validator import validate_presentation

logger = logging.getLogger("pptx_generator")


class PptxGenerator:
    """FULLY FIXED PPTX Generator with Icons and Dynamic Layout"""

    MASTER_LAYOUT_INDEX = {
        "title": 0,
        "title_and_content": 1,
        "section_header": 2,
        "two_content": 3,
        "comparison": 4,
        "title_only": 5,
        "blank": 6,
        "content_with_caption": 7,
        "picture_with_caption": 8
    }

    def __init__(self, template_id: str = "arweqah", language: Optional[str] = None):
        self.template_id = template_id
        self.target_language = language

        self.template_dir = Path(settings.TEMPLATES_DIR) / template_id
        logger.info(f"🔍 Initializing ENHANCED generator: {self.template_dir}")

        if not self.template_dir.exists():
            raise FileNotFoundError(f"Template directory not found: {self.template_dir}")

        self.backgrounds_dir = self.template_dir / "Background"
        if not self.backgrounds_dir.exists():
            self.backgrounds_dir = self.template_dir / "backgrounds"

        # Load JSON
        self.config = self._load_json("config.json")
        self.theme = self._load_json("theme.json") if (self.template_dir / "theme.json").exists() else self._get_default_theme()
        self.constraints = self._load_json("constraints.json") if (self.template_dir / "constraints.json").exists() else {}
        self.layouts = self._load_json("layouts.json")

        # Initialize services
        try:
            self.icon_service = IconService(template_id=template_id)
            logger.info("✅ IconService loaded")
        except Exception as e:
            logger.warning(f"⚠️  IconService failed: {e}")
            self.icon_service = None

        try:
            self.chart_service = ChartService(template_id=template_id)
            logger.info("✅ ChartService loaded")
        except Exception as e:
            logger.warning(f"⚠️  ChartService failed: {e}")
            self.chart_service = None

        self.prs: Optional[Presentation] = None
        self.lang_config = None
        self.section_header_counter = 0

        logger.info(f"✅ ENHANCED Generator initialized")

    def _load_json(self, filename: str) -> Dict:
        json_path = self.template_dir / filename
        if not json_path.exists():
            raise FileNotFoundError(f"Required: {json_path}")
        with open(json_path, 'r', encoding='utf-8') as f:
            return json.load(f)

    def _get_default_theme(self) -> Dict:
        return {
            "colors": {
                "primary": "#F5F0E8",
                "text_primary": "#0D2026",
                "text_inverse": "#FFFCEC",
                "accent_teal": "#01415C"
            },
            "typography": {
                "font_families": {"primary": "Cairo", "body": "Tajawal"},
                "font_sizes": {"title": 44, "body": 18}
            }
        }

    # ========================================================================
    # LANGUAGE CONFIGURATION
    # ========================================================================

    def _detect_and_configure_language(self, presentation_data: PresentationData):
        """Detect and configure language with proper alignment"""
        lang_settings = self.config.get('language_settings', {})

        if not self.target_language:
            self.target_language = lang_settings.get('default', 'en')
            logger.info(f"🌐 Using default language: {self.target_language}")

        self.lang_config = lang_settings.get(self.target_language, {
            'rtl': False,
            'default_font': 'Arial',
            'heading_font': 'Arial',
            'alignment': 'left'
        })

        logger.info(f"   Language: {self.target_language}")
        logger.info(f"   RTL: {self.lang_config.get('rtl')}, Alignment: {self.lang_config.get('alignment')}")

    def get_text_alignment(self, force_center: bool = False) -> PP_ALIGN:
        """Get correct alignment based on language"""
        if force_center:
            return PP_ALIGN.CENTER
        
        if self.lang_config:
            if self.lang_config.get('rtl'):
                return PP_ALIGN.RIGHT
            else:
                return PP_ALIGN.LEFT
        
        return PP_ALIGN.LEFT

    # ========================================================================
    # DYNAMIC TEXT SIZING
    # ========================================================================

    def _get_dynamic_font_size(self, text: str, base_size: int, max_length: int = 100) -> int:
        """Calculate font size based on text length"""
        text_len = len(text)
        
        if text_len <= max_length * 0.5:
            return base_size
        elif text_len <= max_length * 0.75:
            return max(base_size - 4, 24)
        elif text_len <= max_length:
            return max(base_size - 8, 20)
        else:
            return max(base_size - 12, 18)

    def _calculate_text_height(self, text: str, font_size: int, width: float) -> float:
        """Estimate text height in inches"""
        chars_per_line = int(width * 96 / (font_size * 0.6))  # Rough estimate
        lines = max(1, len(text) / chars_per_line)
        line_height = font_size / 72  # Convert pt to inches
        return lines * line_height * 1.5  # 1.5 for line spacing

    # ========================================================================
    # LAYOUT SELECTION
    # ========================================================================

    def get_layout_for_content(self, content_type: str, slide_data=None) -> Dict:
        """Select correct layout with proper detection - works with config.json mappings"""
        content_mapping = self.config.get('content_type_mapping', {})
        layout_name = None

        # Get layout hint from slide data
        layout_hint = self._get_layout_hint(slide_data)
        
        # STEP 1: Try to map layout_hint directly
        if layout_hint:
            hint_lower = layout_hint.lower().strip().replace('_', '').replace('-', '')
            
            # Try exact match first (with normalization)
            if layout_hint in content_mapping:
                layout_name = content_mapping[layout_hint]
                logger.debug(f"   Exact match: {layout_hint} -> {layout_name}")
            # Try normalized match
            elif hint_lower in ['chartslide', 'chart']:
                layout_name = 'chart_slide'
                logger.debug(f"   Chart layout from hint: {layout_hint}")
            elif hint_lower in ['tableslide', 'table']:
                layout_name = 'table_slide'
                logger.debug(f"   Table layout from hint: {layout_hint}")
            elif hint_lower in ['agenda', 'agendaslide']:
                layout_name = 'agenda_slide'
                logger.debug(f"   Agenda layout from hint: {layout_hint}")
            elif 'fourbox' in hint_lower:
                if 'icon' in hint_lower:
                    layout_name = 'four_box_with_icons'
                else:
                    layout_name = 'four_boxes'
                logger.debug(f"   Four-box layout from hint: {layout_hint}")
            elif hint_lower in ['contentparagraph', 'paragraph']:
                layout_name = 'content_paragraph'
                logger.debug(f"   Paragraph layout from hint: {layout_hint}")
            elif hint_lower in ['twocontent', 'twocolumn', 'comparison']:
                layout_name = 'two_content'
                logger.debug(f"   Two-content layout from hint: {layout_hint}")

        # STEP 2: If no layout from hint, use content_type
        if not layout_name:
            if content_type == 'section':
                section_layouts = content_mapping.get('section', ['section_header_dark'])
                if isinstance(section_layouts, list):
                    layout_name = section_layouts[self.section_header_counter % len(section_layouts)]
                    self.section_header_counter += 1
                else:
                    layout_name = section_layouts
                logger.debug(f"   Section layout (rotating): {layout_name}")
            else:
                # Try direct lookup in content_mapping
                layout_name = content_mapping.get(content_type)
                if layout_name:
                    logger.debug(f"   Content type mapping: {content_type} -> {layout_name}")
                else:
                    # Fallback based on content type
                    if content_type == 'chart':
                        layout_name = 'chart_slide'
                    elif content_type == 'table':
                        layout_name = 'table_slide'
                    elif content_type == 'agenda':
                        layout_name = 'agenda_slide'
                    else:
                        layout_name = 'title_and_content'
                    logger.debug(f"   Fallback for {content_type}: {layout_name}")

        # STEP 3: Get layout config from layouts.json
        layout_config = self.layouts.get(layout_name)

        if not layout_config:
            logger.warning(f"⚠️  Layout '{layout_name}' not found in layouts.json, using title_and_content")
            layout_config = self.layouts.get('title_and_content', {})
            layout_name = 'title_and_content'

        logger.info(f"   ✓ Layout selected: '{layout_name}' (hint={layout_hint}, type={content_type})")
        return layout_config

    # ========================================================================
    # HELPER METHODS
    # ========================================================================

    def _get_layout_hint(self, slide_data) -> Optional[str]:
        """Extract layout hint with better detection"""
        if slide_data is None:
            return None

        hint = None
        if isinstance(slide_data, dict):
            hint = slide_data.get('layout_hint')
        else:
            hint = getattr(slide_data, 'layout_hint', None)

        if isinstance(hint, str):
            hint_clean = hint.strip().lower()
            
            # Normalize common variations
            if 'agenda' in hint_clean:
                return 'agenda'
            elif 'four' in hint_clean and 'box' in hint_clean:
                if 'icon' in hint_clean:
                    return 'four_box_with_icons'
                return 'four_boxes'
            elif 'table' in hint_clean:
                return 'table_slide'
            elif 'chart' in hint_clean:
                return 'chart_slide'
            elif 'two' in hint_clean and ('column' in hint_clean or 'content' in hint_clean):
                return 'two_content'
            
            return hint_clean
        
        return None

    def _get_data_value(self, data, key: Optional[str]):
        if not data or not key:
            return None
        if isinstance(data, dict):
            return data.get(key)
        return getattr(data, key, None)

    def _extract_text_value(self, data, element_id: str) -> str:
        text = self._get_data_value(data, element_id)

        if not text:
            if element_id == 'title' and hasattr(data, 'title'):
                text = data.title
            elif element_id == 'subtitle' and hasattr(data, 'subtitle'):
                text = getattr(data, 'subtitle', '')
            elif element_id == 'content' and hasattr(data, 'content'):
                text = getattr(data, 'content', '')

        if not text and hasattr(data, 'bullets') and data.bullets and element_id == 'content':
            text = " ".join(
                bullet.text for bullet in data.bullets if getattr(bullet, 'text', None)
            )

        if isinstance(text, list):
            text = "\n".join(str(item) for item in text if item)

        return (text or "").strip()

    def _coerce_bullet_points(self, value: Any) -> List[BulletPoint]:
        bullets: List[BulletPoint] = []
        if not value:
            return bullets

        if isinstance(value, list):
            for item in value:
                if isinstance(item, BulletPoint):
                    bullets.append(item)
                elif isinstance(item, dict):
                    try:
                        bullets.append(BulletPoint(**item))
                    except Exception:
                        text = item.get('text')
                        if text:
                            bullets.append(BulletPoint(text=text))
                elif item:
                    bullets.append(BulletPoint(text=str(item)))
        elif isinstance(value, str):
            bullets.append(BulletPoint(text=value))

        return bullets

    def _extract_bullet_items(self, data, element_id: Optional[str]) -> List[BulletPoint]:
        prioritized_value = None
        if element_id and element_id not in ('content', 'bullets'):
            prioritized_value = self._get_data_value(data, element_id)

        if prioritized_value:
            bullets = self._coerce_bullet_points(prioritized_value)
            if bullets:
                return bullets

        if hasattr(data, 'bullets') and data.bullets:
            return list(data.bullets)

        return []

    def _resolve_alignment_from_style(self, style: Dict, default_alignment: str = 'left') -> PP_ALIGN:
        alignment = style.get('alignment', default_alignment) if style else default_alignment
        rtl_preferred = style.get('rtl_support', False) if style else False
        return self._alignment_enum(alignment, rtl_preferred)

    def _alignment_enum(self, alignment: Optional[str], rtl_preferred: bool = False) -> PP_ALIGN:
        if rtl_preferred and self.lang_config:
            alignment = 'right' if self.lang_config.get('rtl') else 'left'

        alignment = (alignment or '').lower()

        if alignment == 'center':
            return PP_ALIGN.CENTER
        if alignment == 'right':
            return PP_ALIGN.RIGHT
        if alignment == 'justify':
            return PP_ALIGN.JUSTIFY
        return PP_ALIGN.LEFT

    def _get_slide_layout(self, layout_config: Dict):
        layout_index = layout_config.get('master_layout_index')
        if isinstance(layout_index, int) and 0 <= layout_index < len(self.prs.slide_layouts):
            return self.prs.slide_layouts[layout_index]

        layout_key = str(layout_config.get('master_layout', 'blank')).lower()
        mapped_index = self.MASTER_LAYOUT_INDEX.get(layout_key, self.MASTER_LAYOUT_INDEX['blank'])
        mapped_index = min(mapped_index, len(self.prs.slide_layouts) - 1)
        return self.prs.slide_layouts[mapped_index]

    def _clear_default_placeholders(self, slide):
        for shape in list(slide.shapes):
            if getattr(shape, "is_placeholder", False):
                element = shape._element
                element.getparent().remove(element)

    def get_color_hex(self, color_path: str) -> str:
        try:
            return self.theme['colors'].get(color_path, '#000000')
        except:
            return '#000000'

    def get_color_rgb(self, color_path: str) -> Tuple[int, int, int]:
        hex_color = self.get_color_hex(color_path).lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def hex_to_rgb(self, hex_color: str) -> Tuple[int, int, int]:
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    def get_font(self, font_key: str) -> str:
        """Always return valid font name"""
        if self.lang_config:
            if font_key in ['title_font', 'heading_font']:
                font = self.lang_config.get('heading_font', 'Cairo')
            else:
                font = self.lang_config.get('default_font', 'Tajawal')

            if font and isinstance(font, str):
                return font

        return 'Arial'

    def get_font_size(self, size_key: str) -> int:
        try:
            return self.theme['typography']['font_sizes'].get(size_key, 18)
        except:
            return 18

    def get_dimension(self, dim_key: str) -> float:
        try:
            return float(self.config['slide_dimensions'].get(dim_key, 13.333 if dim_key == 'width' else 7.5))
        except:
            return 13.333 if dim_key == 'width' else 7.5

    def _scrub_title(self, title: str) -> str:
        if not title:
            return ""
        t = title.strip()
        t = re.sub(r"\(\s*continued\s*\)", "", t, flags=re.IGNORECASE)
        return re.sub(r"\s{2,}", " ", t).strip()

    # ========================================================================
    # BACKGROUND & PAGE NUMBER
    # ========================================================================

    def _add_background(self, slide, bg_config: Dict) -> None:
        """Add background image"""
        try:
            if bg_config.get('type') == 'image':
                bg_path = self.template_dir / bg_config.get('path', '')

                if not bg_path.exists():
                    logger.warning(f"⚠️  Background not found: {bg_path}")
                    return

                width = Inches(self.get_dimension('width'))
                height = Inches(self.get_dimension('height'))

                picture = slide.shapes.add_picture(
                    str(bg_path), Inches(0), Inches(0),
                    width=width, height=height
                )

                slide.shapes._spTree.remove(picture._element)
                slide.shapes._spTree.insert(2, picture._element)

                logger.debug(f"✅ Background: {bg_path.name}")

        except Exception as e:
            logger.warning(f"⚠️  Background error: {e}")

    def add_page_number(self, slide, page_num: int) -> None:
        """Add page number diamond"""
        try:
            diamond = slide.shapes.add_shape(
                MSO_SHAPE.DIAMOND,
                Inches(0.2), Inches(6.8),
                Inches(0.4), Inches(0.4)
            )

            diamond.fill.solid()
            diamond.fill.fore_color.rgb = RGBColor(*self.get_color_rgb('accent_teal'))
            diamond.line.fill.background()

            tf = diamond.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = str(page_num)
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.bold = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        except Exception as e:
            logger.warning(f"⚠️  Page number failed: {e}")

    # ========================================================================
    # MAIN GENERATION
    # ========================================================================

    def generate(self, presentation_data: PresentationData) -> str:
        """Generate presentation"""
        logger.info("🎨 Starting ENHANCED generation...")

        self._detect_and_configure_language(presentation_data)

        presentation_data.slides = validate_presentation(presentation_data.slides)

        self.prs = Presentation()
        self.prs.slide_width = Inches(self.get_dimension('width'))
        self.prs.slide_height = Inches(self.get_dimension('height'))

        for slide_data in presentation_data.slides:
            if slide_data.title:
                slide_data.title = self._scrub_title(slide_data.title)

        # Title slide
        try:
            self._create_title_slide_dynamic(presentation_data)
        except Exception as e:
            logger.error(f"❌ Title slide: {e}")

        # Content slides
        for idx, slide_data in enumerate(presentation_data.slides):
            logger.info(f"🔨 Slide {idx + 2}: {slide_data.title[:50]}...")

            try:
                layout_type = (slide_data.layout_type or "content").lower()
                layout_hint = self._get_layout_hint(slide_data)

                if layout_type == "section":
                    content_type = "section"
                elif layout_type == "two_column":
                    content_type = "two_column"
                elif layout_hint:
                    content_type = layout_hint
                elif slide_data.table_data:
                    content_type = "table"
                elif slide_data.chart_data:
                    content_type = "chart"
                elif slide_data.bullets:
                    content_type = "bullets"
                else:
                    content_type = "content"

                self._create_slide_from_json(content_type, slide_data, page_num=idx + 2)

            except Exception as e:
                logger.error(f"❌ Slide error: {e}")
                logger.exception(e)

        output_path = self._get_output_path(presentation_data.title)
        self.prs.save(output_path)

        logger.info(f"✅ Generated: {output_path}")
        logger.info(f"   Slides: {len(self.prs.slides)}, Language: {self.target_language}")

        return output_path

    # ========================================================================
    # ENHANCED TITLE SLIDE WITH DYNAMIC SIZING
    # ========================================================================

    def _create_title_slide_dynamic(self, presentation_data: PresentationData) -> None:
        """Create title slide without 'None' in subtitle"""
        layout_config = self.layouts.get('title_slide', {})
        slide_layout = self._get_slide_layout(layout_config)
        slide = self.prs.slides.add_slide(slide_layout)
        self._clear_default_placeholders(slide)

        bg_config = layout_config.get('background', {})
        if bg_config.get('type') == 'image':
            self._add_background(slide, bg_config)

        elements = layout_config.get('elements', [])
        for element in elements:
            if element.get('type') == 'image' and element.get('id') == 'logo':
                self._add_logo(slide, element)

        # Title
        title = presentation_data.title or "Untitled Presentation"
        if len(title) > 80:
            title = title[:77] + "..."
        
        title_element = next((e for e in elements if e.get('id') == 'title'), None)
        
        if title_element:
            pos = title_element.get('position', {})
            size = title_element.get('size', {})
            style = title_element.get('style', {})

            base_font_size = style.get('font_size', 44)
            if len(title) > 60:
                font_size = 36
            elif len(title) > 40:
                font_size = 40
            else:
                font_size = base_font_size

            char_per_line = 50 if font_size >= 40 else 60
            lines = max(1, len(title) / char_per_line)
            required_height = lines * 0.6
            actual_height = max(size.get('height', 1.5), required_height)
            
            textbox = slide.shapes.add_textbox(
                Inches(pos.get('left', 1.5)),
                Inches(pos.get('top', 2.8)),
                Inches(size.get('width', 10.33)),
                Inches(actual_height)
            )

            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(font_size)
            p.font.name = style.get('font', self.get_font('heading_font'))
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(style.get('color', '#FFFCEC')))
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.15

        # Subtitle - CRITICAL FIX: Remove "None"
        subtitle_parts = []
        if presentation_data.subtitle and presentation_data.subtitle != "None":
            subtitle_parts.append(presentation_data.subtitle)
        if presentation_data.author and presentation_data.author != "None":
            subtitle_parts.append(presentation_data.author)
        
        subtitle = "\n".join(subtitle_parts).strip()
        
        subtitle_element = next((e for e in elements if e.get('id') == 'subtitle'), None)
        
        if subtitle_element and subtitle:
            pos = subtitle_element.get('position', {})
            size = subtitle_element.get('size', {})
            style = subtitle_element.get('style', {})

            subtitle_top = pos.get('top', 4.5)
            if lines > 1.5:
                subtitle_top += 0.3

            textbox = slide.shapes.add_textbox(
                Inches(pos.get('left', 2.0)),
                Inches(subtitle_top),
                Inches(size.get('width', 9.33)),
                Inches(size.get('height', 1.0))
            )

            tf = textbox.text_frame
            tf.word_wrap = True
            tf.clear()

            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(style.get('font_size', 24))
            p.font.name = style.get('font', self.get_font('body_font'))
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(style.get('color', '#FFFCEC')))
            p.alignment = PP_ALIGN.CENTER


    # ========================================================================
    # ENHANCED AGENDA SLIDE WITH DECORATIVE LINE
    # ========================================================================

    def _create_agenda_slide_enhanced(self, slide, layout_config: Dict, data) -> None:
        """Enhanced agenda with LARGER font and PERFECT icon alignment"""
        # Background
        bg_config = layout_config.get('background', {})
        if bg_config.get('type') == 'image':
            self._add_background(slide, bg_config)
        
        # Add "Agenda" text on dark side (centered)
        agenda_textbox = slide.shapes.add_textbox(
            Inches(7.5),
            Inches(3.0),
            Inches(4.5),
            Inches(1.5)
        )
        
        tf = agenda_textbox.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        p = tf.paragraphs[0]
        p.text = "Agenda"
        p.font.size = Pt(54)  # Larger
        p.font.name = self.get_font('heading_font')
        p.font.bold = True
        p.font.color.rgb = RGBColor(*self.hex_to_rgb('#FFFCEC'))
        p.alignment = PP_ALIGN.CENTER
        
        # Logo (bottom right)
        logo_path = self.template_dir / "logo.png"
        if logo_path.exists():
            try:
                slide.shapes.add_picture(
                    str(logo_path),
                    Inches(11.5),
                    Inches(6.3),
                    width=Inches(1.5),
                    height=Inches(0.8)
                )
            except:
                pass
        
        # Content bullets (left side) - MAX 5 ITEMS
        bullets = self._extract_bullet_items(data, 'content')
        
        if not bullets:
            bullets = self._extract_bullet_items(data, 'bullets')
        
        if not bullets:
            logger.warning("⚠️  Agenda slide has no content")
            return
        
        # Limit to 5 items
        bullets = bullets[:5]
        
        # Calculate vertical centering
        num_items = len(bullets)
        item_height = 0.85  # Increased spacing
        total_height = num_items * item_height
        start_y = (7.5 - total_height) / 2
        
        text_color = '#0D2026'
        text_rgb = self.hex_to_rgb(text_color)
        font_name = self.get_font('body_font')
        
        # Render each agenda item with icon
        for idx, bullet in enumerate(bullets):
            # Calculate Y position for this item
            item_y = start_y + (idx * item_height)
            
            # Add icon (PERFECT alignment)
            if self.icon_service:
                icon_name = self.icon_service.auto_select_icon(bullet.text or "", "")
                try:
                    icon_data = self.icon_service.render_to_png(
                        icon_name,
                        size=56,  # Larger icon
                        color=text_color
                    )
                    if icon_data:
                        icon_left = Inches(0.8)
                        icon_top = Inches(item_y + 0.15)
                        
                        slide.shapes.add_picture(
                            icon_data,
                            icon_left,
                            icon_top,
                            width=Inches(0.45),
                            height=Inches(0.45)
                        )
                except Exception as e:
                    logger.debug(f"Icon skip: {e}")
            
            # Add text (aligned with icon)
            text_box = slide.shapes.add_textbox(
                Inches(1.4),  # Start after icon
                Inches(item_y),
                Inches(4.8),
                Inches(0.7)
            )
            
            tf = text_box.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = tf.margin_right = Inches(0.05)
            
            p = tf.paragraphs[0]
            p.text = (bullet.text or "").strip()
            p.font.size = Pt(20)  # LARGER FONT
            p.font.name = font_name
            p.font.color.rgb = RGBColor(*text_rgb)
            p.font.bold = False
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.3
        
        logger.info(f"✅ Agenda created with {len(bullets)} items (larger font, centered)")

    # ========================================================================
    # SLIDE CREATION FROM JSON WITH ICONS
    # ========================================================================

    def _create_slide_from_json(self, content_type: str, data, page_num: int = None) -> None:
        """FIXED: Create slide with ALL content types rendering correctly"""
        
        # Get layout hint
        layout_hint = self._get_layout_hint(data)
        
        # Check for agenda
        if layout_hint and 'agenda' in layout_hint.lower():
            content_type = 'agenda'
            logger.info(f"   Detected AGENDA slide from layout_hint: {layout_hint}")
        
        layout_config = self.get_layout_for_content(content_type, data)
        slide_layout = self._get_slide_layout(layout_config)
        slide = self.prs.slides.add_slide(slide_layout)
        self._clear_default_placeholders(slide)
        
        # Background
        bg_config = layout_config.get('background', {})
        if bg_config.get('type') == 'image':
            self._add_background(slide, bg_config)
        
        # Special: Agenda
        if content_type == 'agenda' or (layout_hint and 'agenda' in layout_hint.lower()):
            self._create_agenda_slide_enhanced(slide, layout_config, data)
            if page_num:
                self.add_page_number(slide, page_num)
            return
        
        elements = layout_config.get('elements', [])
        
        # Special: Section headers
        if content_type == 'section':
            title_text = self._extract_text_value(data, 'title')
            
            if len(title_text) > 60:
                title_text = title_text[:57] + "..."
                data.title = title_text
            
            if self.icon_service and title_text:
                if any(word in title_text.lower() for word in ['thank', 'thanks', 'شكر']):
                    icon_name = 'hand-waving'
                else:
                    icon_name = self.icon_service.auto_select_icon(title_text, "")
                
                self._add_centered_section_icon(slide, icon_name, title_text, layout_config)
            
            for element in elements:
                if element.get('type') == 'text' and element.get('id') == 'title':
                    self._add_text_master(slide, element, data, content_type)
            
            return
        
        # ========================================================================
        # CRITICAL FIX: DETECT CONTENT TYPE FROM DATA
        # ========================================================================
        
        has_chart = False
        has_table = False
        has_paragraph = False
        has_bullets = False
        
        chart_data = getattr(data, 'chart_data', None) or getattr(data, 'chart', None)
        table_data = getattr(data, 'table_data', None) or getattr(data, 'table', None)
        content_text = getattr(data, 'content', None)
        bullets = getattr(data, 'bullets', None)
        
        if chart_data:
            has_chart = True
            logger.info(f"   → Chart slide detected")
        if table_data:
            has_table = True
            logger.info(f"   → Table slide detected")
        if content_text and len(str(content_text).strip()) > 0:
            has_paragraph = True
            logger.info(f"   → Paragraph slide detected")
        if bullets and len(bullets) > 0:
            has_bullets = True
            logger.info(f"   → Bullets slide detected")
        
        # ========================================================================
        # PROCESS ELEMENTS - FIXED RENDERING LOGIC
        # ========================================================================
        
        title_element = None
        content_rendered = False
        
        for element in elements:
            try:
                element_type = element.get('type')
                element_id = element.get('id', 'unknown')
                
                # Logo
                if element_type == 'image' and element_id == 'logo':
                    self._add_logo(slide, element)
                
                # Title with icon
                elif element_type == 'text' and element_id == 'title':
                    title_element = element
                    title_text = self._extract_text_value(data, 'title')
                    
                    if len(title_text) > 80:
                        title_text = title_text[:77] + "..."
                        data.title = title_text
                    
                    if self.icon_service and title_text and not (has_chart or has_table):
                        icon_name = self.icon_service.auto_select_icon(title_text, "")
                        self._add_icon_to_title(slide, icon_name, element)
                    
                    self._add_text_master(slide, element, data, content_type)
                
                # Title underline
                elif element_type == 'line' and element_id == 'title_line':
                    if title_element:
                        self._add_title_underline(slide, title_element)
                
                # ========================================================================
                # CRITICAL FIX: RENDER BASED ON DATA, NOT JUST ELEMENT TYPE
                # ========================================================================
                
                # CHART: Render if we detected chart data AND this is a chart element
                elif element_type == 'chart' and has_chart and not content_rendered:
                    logger.info(f"   → Rendering CHART")
                    
                    position = element.get('position', {'left': 1.5, 'top': 2.0})
                    size = element.get('size', {'width': 10.33, 'height': 5.0})
                    style = element.get('style', {})
                    
                    bg_color = style.get('background_color', '#0D2026')
                    bg_rgb = self.hex_to_rgb(bg_color)
                    
                    # Convert chart_data to dict
                    if hasattr(chart_data, 'dict'):
                        chart_dict = chart_data.dict()
                    elif isinstance(chart_data, dict):
                        chart_dict = chart_data
                    else:
                        chart_dict = chart_data.__dict__
                    
                    # Ensure visible text - use style colors with fallbacks
                    chart_dict['font_color'] = style.get('font_color', '#FFFFFF')
                    chart_dict['data_label_color'] = style.get('data_label_color', '#FFFFFF')
                    chart_dict['axis_color'] = style.get('axis_color', '#FFFFFF')
                    chart_dict['axis_label_color'] = style.get('axis_label_color', '#FFFFFF')
                    chart_dict['grid_color'] = style.get('grid_color', '#5C6F7A')
                    chart_dict['legend_font_color'] = style.get('legend_font_color', '#FFFFFF')
                    
                    # Call ChartService
                    if self.chart_service:
                        self.chart_service.add_native_chart(
                            slide=slide,
                            chart_data=chart_dict,
                            position=position,
                            size=size,
                            background_rgb=bg_rgb
                        )
                        content_rendered = True
                        logger.info(f"   ✅ Chart rendered")
                    else:
                        logger.warning(f"   ⚠️ ChartService not available")
                
                # TABLE: Render if we detected table data AND this is a table element
                elif element_type == 'table' and has_table and not content_rendered:
                    logger.info(f"   → Rendering TABLE")
                    self._add_table_master(slide, element, data)
                    content_rendered = True
                    logger.info(f"   ✅ Table rendered")
                
                # PARAGRAPH: Render if we have paragraph content
                elif element_type == 'text_paragraph' and has_paragraph and not content_rendered:
                    logger.info(f"   → Rendering PARAGRAPH")
                    self._add_paragraph_text(slide, element, data)
                    content_rendered = True
                    logger.info(f"   ✅ Paragraph rendered")
                
                # BULLETS: Render if we have bullets (or if element type is bullets/text_bullets)
                elif (element_type in ['bullets', 'text_bullets'] or element_id == 'content') and has_bullets and not content_rendered:
                    logger.info(f"   → Rendering BULLETS")
                    self._add_bullets_master(slide, element, data)
                    content_rendered = True
                    logger.info(f"   ✅ Bullets rendered")
                
                # FOUR-BOX: Render if layout indicates boxes
                elif element_type == 'boxes' and element_id == 'content_boxes':
                    logger.info(f"   → Rendering FOUR-BOX")
                    self._add_content_box_with_icon_enhanced(slide, element, data)
                    content_rendered = True
                
                # Shapes
                elif element_type == 'shape':
                    self._add_shape_element(slide, element)
                
                # Standalone icons
                elif element_type == 'icon' and element_id != 'icon':  # Skip title icon
                    self._add_icon_master(slide, element, data)
                
                # Content boxes
                elif element_type == 'content_box':
                    self._add_content_box_master(slide, element, data)
            
            except Exception as e:
                logger.warning(f"⚠️  Element '{element_id}' error: {e}")
                import traceback
                traceback.print_exc()
        
        # ========================================================================
        # FALLBACK: If content detected but not rendered, try generic rendering
        # ========================================================================
        if not content_rendered:
            if has_chart and chart_data:
                logger.warning(f"⚠️  Chart detected but not rendered - attempting fallback")
                # Find any chart element
                chart_element = next((e for e in elements if e.get('type') == 'chart'), None)
                if chart_element:
                    try:
                        position = chart_element.get('position', {'left': 1.5, 'top': 2.0})
                        size = chart_element.get('size', {'width': 10.33, 'height': 5.0})
                        style = chart_element.get('style', {})
                        
                        bg_color = style.get('background_color', '#0D2026')
                        bg_rgb = self.hex_to_rgb(bg_color)
                        
                        if hasattr(chart_data, 'dict'):
                            chart_dict = chart_data.dict()
                        elif isinstance(chart_data, dict):
                            chart_dict = chart_data
                        else:
                            chart_dict = chart_data.__dict__
                        
                        chart_dict['font_color'] = style.get('font_color', '#FFFFFF')
                        chart_dict['data_label_color'] = style.get('data_label_color', '#FFFFFF')
                        chart_dict['axis_color'] = style.get('axis_color', '#FFFFFF')
                        chart_dict['axis_label_color'] = style.get('axis_label_color', '#FFFFFF')
                        chart_dict['grid_color'] = style.get('grid_color', '#5C6F7A')
                        chart_dict['legend_font_color'] = style.get('legend_font_color', '#FFFFFF')
                        
                        if self.chart_service:
                            self.chart_service.add_native_chart(
                                slide=slide,
                                chart_data=chart_dict,
                                position=position,
                                size=size,
                                background_rgb=bg_rgb
                            )
                            content_rendered = True
                            logger.info(f"   ✅ Chart rendered (fallback)")
                    except Exception as e:
                        logger.error(f"   ❌ Chart fallback failed: {e}")
            
            elif has_table and table_data:
                logger.warning(f"⚠️  Table detected but not rendered - attempting fallback")
                # Find any table element
                table_element = next((e for e in elements if e.get('type') == 'table'), None)
                if table_element:
                    try:
                        self._add_table_master(slide, table_element, data)
                        content_rendered = True
                        logger.info(f"   ✅ Table rendered (fallback)")
                    except Exception as e:
                        logger.error(f"   ❌ Table fallback failed: {e}")
            
            elif has_paragraph and content_text:
                logger.warning(f"⚠️  Paragraph detected but not rendered - attempting fallback")
                # Find any text_paragraph or content element
                para_element = next((e for e in elements if e.get('type') in ['text_paragraph', 'text_bullets'] or e.get('id') == 'content'), None)
                if para_element:
                    try:
                        self._add_paragraph_text(slide, para_element, data)
                        content_rendered = True
                        logger.info(f"   ✅ Paragraph rendered (fallback)")
                    except Exception as e:
                        logger.error(f"   ❌ Paragraph fallback failed: {e}")
            
            elif has_bullets and bullets:
                logger.warning(f"⚠️  Bullets detected but not rendered - attempting fallback")
                # Find any bullets or content element
                bullets_element = next((e for e in elements if e.get('type') in ['bullets', 'text_bullets'] or e.get('id') == 'content'), None)
                if bullets_element:
                    try:
                        self._add_bullets_master(slide, bullets_element, data)
                        content_rendered = True
                        logger.info(f"   ✅ Bullets rendered (fallback)")
                    except Exception as e:
                        logger.error(f"   ❌ Bullets fallback failed: {e}")
        
        # Final warning if still not rendered
        if not content_rendered and (has_chart or has_table or has_paragraph or has_bullets):
            logger.error(f"❌ CONTENT NOT RENDERED! has_chart={has_chart}, has_table={has_table}, has_paragraph={has_paragraph}, has_bullets={has_bullets}")
            logger.error(f"   Available elements: {[e.get('id') + ':' + e.get('type') for e in elements]}")
        
        # Add page number
        if page_num:
            self.add_page_number(slide, page_num)



    # ========================================================================
    # ICON INTEGRATION
    # ========================================================================

    def _add_icon_to_title(self, slide, icon_name: str, title_element: Dict) -> None:
        """Add icon next to title with matching color"""
        if not self.icon_service:
            return
        
        try:
            pos = title_element.get('position', {})
            style = title_element.get('style', {})
            
            # Get text color from style to match icon color
            text_color = style.get('color', '#01415C')
            
            # Icon position: left of title
            icon_left = pos.get('left', 1.0) - 0.7
            icon_top = pos.get('top', 0.6) + 0.05
            
            icon_data = self.icon_service.render_to_png(
                icon_name,
                size=80,
                color=text_color  # MATCH TEXT COLOR
            )
            
            if icon_data:
                slide.shapes.add_picture(
                    icon_data,
                    Inches(icon_left),
                    Inches(icon_top),
                    width=Inches(0.5),
                    height=Inches(0.5)
                )
                logger.debug(f"✅ Icon added to title: {icon_name} (color: {text_color})")
        except Exception as e:
            logger.warning(f"⚠️  Icon render failed: {e}")

    def _add_centered_section_icon(self, slide, icon_name: str, title_text: str, layout_config: Dict) -> None:
        """
        Add icon CENTERED ABOVE title for section headers (not beside)
        This creates better visual balance for section dividers
        """
        if not self.icon_service:
            return
        
        try:
            # Get title element to extract color
            elements = layout_config.get('elements', [])
            title_element = next((e for e in elements if e.get('id') == 'title'), None)
            
            if not title_element:
                return
            
            title_style = title_element.get('style', {})
            text_color = title_style.get('color', '#FFFCEC')
            
            # Get slide dimensions
            slide_width = self.get_dimension('width')
            
            # Icon size
            icon_size = 1.2  # inches
            
            # Position: Centered horizontally, above title
            icon_left = (slide_width - icon_size) / 2
            icon_top = 2.5  # Fixed position above title
            
            icon_data = self.icon_service.render_to_png(
                icon_name,
                size=int(icon_size * 96),
                color=text_color  # Match title color
            )
            
            if icon_data:
                slide.shapes.add_picture(
                    icon_data,
                    Inches(icon_left),
                    Inches(icon_top),
                    width=Inches(icon_size),
                    height=Inches(icon_size)
                )
                logger.info(f"✅ Centered section icon: {icon_name} (color: {text_color})")
        
        except Exception as e:
            logger.warning(f"⚠️  Centered section icon failed: {e}")
    
    def _add_title_underline(self, slide, title_element: Dict) -> None:
        """
        Add horizontal line below title - LEFT-ALIGNED starting from left margin
        """
        try:
            # **FIX: Start from left margin (1.0"), not from title position**
            line_left = 1.0  # Fixed left margin
            line_top = 1.5   # Below title
            line_width = 11.33  # Full content width
            
            # Create thin rectangle as line
            line_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(line_left),
                Inches(line_top),
                Inches(line_width),
                Inches(0.02)  # Very thin line
            )
            
            # Style the line
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(*self.get_color_rgb('accent_teal'))
            line_shape.line.fill.background()  # No border
            
            logger.debug("✅ Added title underline (left-aligned)")
        
        except Exception as e:
            logger.warning(f"⚠️  Title underline failed: {e}")



    def _add_content_box_with_icon_enhanced(self, slide, element: Dict, data) -> None:
        """
        Enhanced four-box layout with icon INSIDE box (top-center) and text below
        Creates consistent, professional appearance across all 4 boxes
        """
        # Extract box configuration
        box_config = element
        bullets = self._extract_bullet_items(data, 'content')
        
        if not bullets or len(bullets) == 0:
            return
        
        # Get layout configuration
        layout_type = box_config.get('layout', '2x2')  # Default 2x2 grid
        position = box_config.get('position', {})
        size = box_config.get('size', {})
        box_style = box_config.get('box_style', {})
        
        # Box dimensions
        box_width = box_style.get('width', 5.4)
        box_height = box_style.get('height', 2.2)
        gap_h = box_style.get('gap_horizontal', 0.53)
        gap_v = box_style.get('gap_vertical', 0.4)
        
        # Colors for 4 boxes
        colors = box_style.get('colors', ['#B1D8BE', '#F9D462', '#C6C3BE', '#E09059'])
        text_color = box_style.get('text_color', '#0D2026')
        font_name = box_style.get('font', self.get_font('body_font'))
        font_size = box_style.get('font_size', 16)
        icon_size = box_style.get('icon_size', 0.6)
        
        # Base position
        base_left = position.get('left', 1.0)
        base_top = position.get('top', 2.0)
        
        # Render each of the 4 boxes
        for idx, bullet in enumerate(bullets[:4]):  # Max 4 boxes
            # Calculate grid position (2x2)
            row = idx // 2
            col = idx % 2
            
            # Calculate box position
            box_left = Inches(base_left + col * (box_width + gap_h))
            box_top = Inches(base_top + row * (box_height + gap_v))
            box_w = Inches(box_width)
            box_h = Inches(box_height)
            
            # 1. Create colored box
            box_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                box_left, box_top, box_w, box_h
            )
            
            box_color = colors[idx % len(colors)]
            box_rgb = self.hex_to_rgb(box_color)
            
            box_shape.fill.solid()
            box_shape.fill.fore_color.rgb = RGBColor(*box_rgb)
            box_shape.line.fill.background()  # No border
            
            # 2. Add icon INSIDE box (top-center)
            if self.icon_service:
                icon_name = self.icon_service.auto_select_icon(bullet.text or "", "")
                
                try:
                    icon_data = self.icon_service.render_to_png(
                        icon_name,
                        size=int(icon_size * 96),
                        color=text_color  # Match text color
                    )
                    
                    if icon_data:
                        # Center icon horizontally within box
                        icon_left = box_left + (box_w - Inches(icon_size)) / 2
                        icon_top = box_top + Inches(0.25)  # Small padding from top
                        
                        slide.shapes.add_picture(
                            icon_data,
                            icon_left,
                            icon_top,
                            Inches(icon_size),
                            Inches(icon_size)
                        )
                except Exception as e:
                    logger.debug(f"Icon render skipped for box {idx + 1}: {e}")
            
            # 3. Add text BELOW icon (inside box)
            text_top = box_top + Inches(icon_size + 0.35)  # Below icon with gap
            text_height = box_h - Inches(icon_size + 0.5)  # Remaining space
            
            text_box = slide.shapes.add_textbox(
                box_left + Inches(0.2),  # Small padding from box edges
                text_top,
                box_w - Inches(0.4),  # Account for padding
                text_height
            )
            
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.clear()
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = tf.margin_right = Inches(0.1)
            
            p = tf.paragraphs[0]
            p.text = (bullet.text or "").strip()
            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = 1.3
            
            logger.debug(f"✅ Four-box item {idx + 1}: {bullet.text[:30]}...")

    # ========================================================================
    # OTHER ELEMENT METHODS
    # ========================================================================

    def _add_logo(self, slide, element: Dict) -> None:
        """Add logo to slide"""
        try:
            logo_path = self.template_dir / element.get('path', 'logo.png')
            if not logo_path.exists():
                return
            
            pos = element.get('position', {})
            size = element.get('size', {})
            
            slide.shapes.add_picture(
                str(logo_path),
                Inches(pos.get('left', 0.5)),
                Inches(pos.get('top', 6.2)),
                width=Inches(size.get('width', 2.0)),
                height=Inches(size.get('height', 0.8))
            )
        except Exception as e:
            logger.warning(f"⚠️  Logo failed: {e}")

    def _apply_solid_background(self, slide, color_hex: str) -> None:
        """Apply solid color background"""
        try:
            color = color_hex.lstrip("#")
            bg_xml = f'''
            <p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"
                xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <p:bgPr>
                <a:solidFill>
                <a:srgbClr val="{color}"/>
                </a:solidFill>
            </p:bgPr>
            </p:bg>
            '''
            bg_element = parse_xml(bg_xml)
            slide._element.insert(0, bg_element)
        except Exception as e:
            logger.error(f"❌ Background failed: {e}")

    def _add_shape_element(self, slide, element: Dict) -> None:
        """Render decorative shapes from JSON layouts"""
        pos = element.get('position', {})
        size = element.get('size', {})
        style = element.get('style', {})

        shape_type = (element.get('shape') or 'rectangle').lower()
        shape_map = {
            'circle': MSO_SHAPE.OVAL,
            'rectangle': MSO_SHAPE.RECTANGLE
        }
        shape_enum = shape_map.get(shape_type, MSO_SHAPE.RECTANGLE)

        try:
            shape = slide.shapes.add_shape(
                shape_enum,
                Inches(pos.get('left', 0)),
                Inches(pos.get('top', 0)),
                Inches(size.get('width', 1)),
                Inches(size.get('height', 1))
            )

            fill_color = style.get('fill_color')
            if fill_color:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(fill_color))
                opacity = style.get('opacity')
                if opacity is not None:
                    try:
                        transparency = max(0.0, min(1.0, 1.0 - float(opacity)))
                        shape.fill.fore_color.transparency = transparency
                    except Exception:
                        pass
            else:
                shape.fill.background()

            shape.line.fill.background()
        except Exception as exc:
            logger.debug(f"Shape render skipped: {exc}")

    def _add_text_master(self, slide, element: Dict, data, content_type: str) -> None:
        """Add text with support for 'content' field"""
        pos = element.get('position', {})
        size = element.get('size', {})
        style = element.get('style', {})
        element_id = element.get('id', '')

        # CRITICAL: Extract from both 'title' and 'content' fields
        if element_id == 'title':
            text = self._extract_text_value(data, 'title')
        elif element_id == 'content':
            text = self._extract_text_value(data, 'content')
        else:
            text = self._extract_text_value(data, element_id)

        if not text:
            return

        textbox = slide.shapes.add_textbox(
            Inches(pos.get('left', 0)),
            Inches(pos.get('top', 0)),
            Inches(size.get('width', 5)),
            Inches(size.get('height', 1))
        )

        tf = textbox.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP

        padding = float(style.get('padding', 0))
        margin_value = Inches(max(padding, 0))
        tf.margin_left = tf.margin_right = margin_value
        tf.margin_top = tf.margin_bottom = margin_value

        line_spacing = max(style.get('line_spacing', 1.45), 1.35)
        
        if content_type in ['title', 'section']:
            alignment = self._resolve_alignment_from_style(style)
        else:
            alignment = self.get_text_alignment()
        
        text_color = style.get('color')
        if not text_color:
            if content_type in ['title', 'section']:
                text_color = '#FFFCEC'
            else:
                text_color = '#0D2026'
        
        font_name = style.get('font')
        if not font_name and element_id == 'title':
            font_name = self.get_font('heading_font')
        elif not font_name:
            font_name = self.get_font('body_font')

        paragraphs = [line.strip() for line in str(text).splitlines() if line.strip()]
        if not paragraphs:
            paragraphs = [str(text).strip()]

        for idx, line in enumerate(paragraphs):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = line
            paragraph.font.size = Pt(style.get('font_size', 18))
            paragraph.font.name = font_name
            if style.get('bold'):
                paragraph.font.bold = True

            if text_color:
                rgb = self.hex_to_rgb(text_color)
                paragraph.font.color.rgb = RGBColor(*rgb)

            paragraph.line_spacing = line_spacing
            paragraph.space_after = Pt(style.get('space_after', 6))
            paragraph.alignment = alignment

    def _add_paragraph_text(self, slide, element: Dict, data) -> None:
        """FIXED: Render paragraph content"""
        pos = element.get('position', {})
        size = element.get('size', {})
        style = element.get('style', {})

        # Get text from 'content' field
        text = getattr(data, 'content', None)
        
        if not text or len(str(text).strip()) == 0:
            if hasattr(data, 'bullets') and data.bullets:
                text = "\n\n".join(b.text for b in data.bullets if getattr(b, 'text', None))

        text = (text or "").strip()
        if not text:
            logger.warning(f"⚠️  No paragraph text found")
            return

        logger.info(f"   → Paragraph text: {len(text)} chars")

        textbox = slide.shapes.add_textbox(
            Inches(pos.get('left', 1.5)),
            Inches(pos.get('top', 2.0)),
            Inches(size.get('width', 10.33)),
            Inches(size.get('height', 4.8))
        )

        tf = textbox.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.2)

        # Split into paragraphs
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
        if not paragraphs:
            paragraphs = [p.strip() for p in text.split('. ') if p.strip()]
        if not paragraphs:
            paragraphs = [text]

        font_name = style.get('font', self.get_font('body_font'))
        text_color = style.get('color', '#0D2026')

        for idx, chunk in enumerate(paragraphs):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = chunk
            paragraph.font.size = Pt(18)
            paragraph.font.name = font_name
            paragraph.line_spacing = 1.6
            paragraph.space_after = Pt(12)
            paragraph.alignment = PP_ALIGN.LEFT

            rgb = self.hex_to_rgb(text_color)
            paragraph.font.color.rgb = RGBColor(*rgb)
        
        logger.info(f"   ✅ Paragraph added: {len(paragraphs)} blocks")

    def _add_bullets_master(self, slide, element: Dict, data) -> None:
        """Add bullets with VISIBLE bullet symbols"""
        element_id = element.get('id')
        bullets = self._extract_bullet_items(data, element_id)
        if not bullets:
            return
        
        pos = element.get('position', {})
        size = element.get('size', {})
        style = element.get('style', {})
        bullet_style_cfg = element.get('bullet_style', {})
        
        textbox = slide.shapes.add_textbox(
            Inches(pos.get('left', 0)),
            Inches(pos.get('top', 0)),
            Inches(size.get('width', 5)),
            Inches(size.get('height', 3))
        )
        
        tf = textbox.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Inches(style.get('padding', 0.1))
        
        font_size = style.get('font_size', 18)
        line_spacing = max(style.get('line_spacing', 1.75), 1.5)
        
        font_name = style.get('font') or self.get_font('body_font')
        text_color = style.get('color', '#0D2026')
        text_rgb = self.hex_to_rgb(text_color)
        text_alignment = self._resolve_alignment_from_style(style)
        
        # Get bullet character
        bullet_char = bullet_style_cfg.get('bullet_char', '●')
        
        for idx, bullet in enumerate(bullets[:6]):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            
            # CRITICAL: Add bullet symbol manually
            bullet_text = (bullet.text or "").replace("**", "").strip()
            paragraph.text = f"{bullet_char}  {bullet_text}"
            
            paragraph.font.bold = style.get('bold', False)
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = font_name
            paragraph.font.color.rgb = RGBColor(*text_rgb)
            paragraph.alignment = text_alignment
            paragraph.line_spacing = line_spacing
            paragraph.space_before = Pt(bullet_style_cfg.get('spacing_before', 6))
            paragraph.space_after = Pt(bullet_style_cfg.get('spacing_after', 6))
            
            # Sub-bullets
            if getattr(bullet, 'sub_bullets', None):
                sub_bullet_char = bullet_style_cfg.get('sub_bullet_char', '○')
                for sub in bullet.sub_bullets[:3]:
                    sp = tf.add_paragraph()
                    sub_text = str(sub).strip()
                    sp.text = f"    {sub_bullet_char}  {sub_text}"
                    sp.font.size = Pt(max(font_size - 2, 14))
                    sp.font.name = font_name
                    sp.font.color.rgb = RGBColor(*text_rgb)
                    sp.alignment = text_alignment
                    sp.line_spacing = line_spacing

    def _add_icon_master(self, slide, element: Dict, data) -> None:
        """
        FIXED:
        - Supports standalone icons (id == "icon")
        - Avoids int() conversion on IDs
        - Auto-selects icon based on slide title or box content
        """

        if not self.icon_service:
            return

        element_id = element.get("id", "")
        pos = element.get("position", {})
        size = element.get("size", {})

        # -------------------------------
        # 1) Determine the text to derive icon from
        # -------------------------------
        content = ""

        # Case A: Standalone icon → use slide title
        if element_id == "icon":
            if hasattr(data, "title") and data.title:
                content = data.title

        # Case B: Box icons like box1_icon, box2_icon
        elif element_id.endswith("_icon"):
            box_id = element_id.replace("_icon", "")
            if hasattr(data, "bullets") and data.bullets:
                try:
                    idx = int(box_id.replace("box", "")) - 1
                    if 0 <= idx < len(data.bullets):
                        content = data.bullets[idx].text
                except:
                    # Fallback: use title
                    content = getattr(data, "title", "")

        # Fallback safeguard
        if not content:
            content = getattr(data, "title", "") or ""

        if not content.strip():
            return

        # -------------------------------
        # 2) Resolve icon
        # -------------------------------
        icon_name = self.icon_service.auto_select_icon(content, "")

        try:
            icon_data = self.icon_service.render_to_png(
                icon_name,
                size=int(size.get("width", 0.6) * 96),
                color=element.get("color", "#0D2026")
            )

            if not icon_data:
                return

            slide.shapes.add_picture(
                icon_data,
                Inches(pos.get("left", 0)),
                Inches(pos.get("top", 0)),
                width=Inches(size.get("width", 0.6)),
                height=Inches(size.get("height", 0.6))
            )

            logger.debug(f"✅ Icon added: {icon_name} (content='{content[:30]}...')")

        except Exception as e:
            logger.warning(f"⚠️ Icon render failed: {e}")


    def _add_content_box_master(self, slide, element: Dict, data) -> None:
        """Add colored content box (text only, icon separate)"""
        box_id = element.get('id', '')
        pos = element.get('position', {})
        size = element.get('size', {})
        style = element.get('style', {})
        
        content = ""
        if isinstance(data, dict):
            content = data.get(box_id, '')
        elif hasattr(data, 'bullets') and data.bullets:
            box_num = int(box_id.replace('box', '')) - 1
            if box_num < len(data.bullets):
                content = data.bullets[box_num].text
        
        if not content:
            return
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(pos.get('left', 0)),
            Inches(pos.get('top', 0)),
            Inches(size.get('width', 3)),
            Inches(size.get('height', 2))
        )
        
        bg_color = style.get('background_color', '#FFFFFF')
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(bg_color))
        box.line.fill.background()
        
        padding = style.get('padding', 0.3)
        textbox = slide.shapes.add_textbox(
            Inches(pos.get('left') + padding),
            Inches(pos.get('top') + padding),
            Inches(size.get('width') - 2*padding),
            Inches(size.get('height') - 2*padding)
        )
        
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = tf.margin_right = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(style.get('font_size', 16))
        p.font.name = style.get('font', self.get_font('body_font'))
        text_alignment_style = dict(style)
        if 'text_alignment' in style:
            text_alignment_style['alignment'] = style.get('text_alignment')
        p.alignment = self._resolve_alignment_from_style(text_alignment_style, default_alignment='left')
        p.line_spacing = max(style.get('line_spacing', 1.45), 1.35)
        p.space_after = Pt(style.get('space_after', 4))
        
        text_color = style.get('text_color', '#0D2026')
        p.font.color.rgb = RGBColor(*self.hex_to_rgb(text_color))

    def _add_table_master(self, slide, element: Dict, data) -> None:
        """Add table - FIXED"""
        table_data_obj = getattr(data, 'table_data', None) or getattr(data, 'table', None)
        
        if not table_data_obj:
            logger.warning(f"⚠️  No table_data found")
            return
        
        pos = element.get('position', {})
        size = element.get('size', {})
        style = element.get('style', {})
        
        # Convert to dict
        if hasattr(table_data_obj, 'dict'):
            table_dict = table_data_obj.dict()
        elif isinstance(table_data_obj, dict):
            table_dict = table_data_obj
        else:
            table_dict = table_data_obj.__dict__
        
        headers = table_dict.get("headers", [])
        rows = table_dict.get("rows", [])
        
        if not headers or not rows:
            logger.warning(f"⚠️  Table has no data")
            return
        
        logger.info(f"   → Table: {len(headers)} cols, {len(rows)} rows")
        
        num_cols = len(headers)
        num_rows = len(rows) + 1  # +1 for header
        
        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(pos.get('left', 1.0)),
            Inches(pos.get('top', 2.0)),
            Inches(size.get('width', 11.33)),
            Inches(size.get('height', 4.5))
        )
        
        table = table_shape.table
        
        # Header styling
        header_bg = style.get('header_bg', '#01415C')
        header_text = style.get('header_text', '#FFFCEC')
        header_alignment = PP_ALIGN.CENTER
        header_font = style.get('header_font', self.get_font('heading_font'))
        
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(header_bg))
            
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(16)
                paragraph.font.color.rgb = RGBColor(*self.hex_to_rgb(header_text))
                paragraph.font.name = header_font
                paragraph.alignment = header_alignment
        
        # Body rows
        body_text = style.get('body_text', '#0D2026')
        body_alignment = PP_ALIGN.LEFT
        body_font = style.get('body_font', self.get_font('body_font'))
        alt_bg = style.get('alternate_row_bg', '#E8E3D8')
        
        for row_idx, row in enumerate(rows, start=1):
            for col_idx, cell_value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_value)
                
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(alt_bg))
                
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.font.color.rgb = RGBColor(*self.hex_to_rgb(body_text))
                    paragraph.font.name = body_font
                    paragraph.alignment = body_alignment
        
        logger.info(f"✅ Table rendered successfully")

    def _add_chart_master(self, slide, element: Dict, data) -> None:
        """Add chart with VISIBLE text - FIXED"""
        if not self.chart_service:
            logger.warning("⚠️  ChartService not available")
            return
        
        chart_data_obj = getattr(data, 'chart_data', None) or getattr(data, 'chart', None)
        
        if not chart_data_obj:
            logger.warning(f"⚠️  No chart_data found")
            return
        
        pos = element.get('position', {})
        size = element.get('size', {})
        style = element.get('style', {})
        
        try:
            # Convert to dict
            if hasattr(chart_data_obj, 'dict'):
                chart_data = chart_data_obj.dict()
            elif isinstance(chart_data_obj, dict):
                chart_data = chart_data_obj
            else:
                chart_data = chart_data_obj.__dict__
            
            logger.info(f"   → Chart type: {chart_data.get('chart_type')}")
            logger.info(f"   → Categories: {chart_data.get('categories')}")
            logger.info(f"   → Series: {chart_data.get('series')}")
            
            # Ensure text is visible on dark background
            chart_data['font_color'] = '#FFFCEC'
            chart_data['data_label_color'] = '#FFFCEC'
            chart_data['axis_color'] = '#FFFCEC'
            chart_data['axis_label_color'] = '#FFFCEC'
            chart_data['grid_color'] = '#5C6F7A'
            chart_data['legend_font_color'] = '#FFFCEC'
            chart_data['title_color'] = '#FFFCEC'
            
            bg_color = style.get('background_color', '#0D2026')
            bg_rgb = self.hex_to_rgb(bg_color)
            
            self.chart_service.add_native_chart(
                slide=slide,
                chart_data=chart_data,
                position={"left": pos.get('left', 1.5), "top": pos.get('top', 2.0)},
                size={"width": size.get('width', 10.33), "height": size.get('height', 5.0)},
                background_rgb=bg_rgb
            )
            
            logger.info(f"✅ Chart rendered successfully")
        except Exception as e:
            logger.error(f"❌ Chart error: {e}")
            logger.exception(e)

    def _get_output_path(self, title: str) -> str:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = (title or "Presentation").strip()
        safe_title = re.sub(r'[<>:"/\\\\|?*]', '', safe_title).replace(" ", "_")[:50]
        filename = f"{self.template_id}_{safe_title}_{timestamp}.pptx"
        output_dir = Path(settings.OUTPUT_DIR)
        output_dir.mkdir(exist_ok=True, parents=True)
        return str(output_dir / filename)