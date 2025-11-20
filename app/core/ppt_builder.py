import json
import logging
from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger("ppt_builder")

def json_to_pptx(json_text: str, out_path: str) -> str:
    """
    Accepts OpenAI JSON slides (array of objects).
    Writes a simple PPTX (blank layout) with title + bullets per slide.
    """
    try:
        slides = json.loads(json_text)
        prs = Presentation()
        for s in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            shapes = slide.shapes

            title = s.get("title") or ""
            content = s.get("content") or []

            if title:
                tbox = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                tf = tbox.text_frame
                p = tf.paragraphs[0]
                p.text = title
                p.font.bold = True
                p.font.size = Pt(30)

            top = 1.5
            for item in content:
                # each item could be a string or list (we simplify to string)
                text = item if isinstance(item, str) else ", ".join(map(str, item)) if isinstance(item, list) else str(item)
                tbox = shapes.add_textbox(Inches(0.8), Inches(top), Inches(8), Inches(0.5))
                tf = tbox.text_frame
                p = tf.paragraphs[0]
                p.text = text
                p.font.size = Pt(18)
                top += 0.6

        prs.save(out_path)
        logger.info(f"PPTX written: {out_path}")
        return out_path
    except Exception as e:
        logger.exception("json_to_pptx failed")
        raise
