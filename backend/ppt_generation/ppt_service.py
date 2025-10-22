import os
import uuid as uuid_module
import subprocess
from typing import Dict, Any, List, Tuple, Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
import matplotlib
matplotlib.use('Agg') 
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np

from routes.config import CHARTS_DIR
from routes.logging import logger


# Set seaborn style globally
sns.set_theme(style="whitegrid")
sns.set_palette("husl")


def analyze_template_layouts(template_path: str) -> Tuple[str, Dict[int, Dict]]:
    if not template_path or not Path(template_path).exists():
        logger.warning(f"Template not found: {template_path}")
        return "Default layouts: TITLE_ONLY, SINGLE_CONTENT", {}
    
    try:
        logger.info(f"Analyzing template: {template_path}")
        
        prs = Presentation(template_path)
        
        analysis = (
            f"Template: {Path(template_path).name}\n"
            f"Dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"\n"
            f"Total Layouts: {len(prs.slide_layouts)}\n\n"
        )
        
        layout_details: Dict[int, Dict] = {}
        
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = []
            title_idx = None
            content_indices: List[int] = []
            
            for ph in layout.placeholders:
                ph_type = ph.placeholder_format.type
                ph_idx = ph.placeholder_format.idx
                ph_name = getattr(ph, "name", "")
                
                placeholders.append({
                    "idx": ph_idx,
                    "type": ph_type,
                    "name": ph_name
                })
                
                if ph_type == 1:
                    title_idx = ph_idx
                
                if ph_type in (2, 4, 7, 14):
                    content_indices.append(ph_idx)
            
            layout_details[i] = {
                "name": layout.name,
                "placeholders": placeholders,
                "title_idx": title_idx,
                "content_indices": content_indices,
            }
            
            ph_names = [f"{p['name']} (idx={p['idx']}, type={p['type']})" for p in placeholders]
            analysis += f"Layout {i}: {layout.name}\n"
            analysis += f"  Placeholders: {', '.join(ph_names)}\n"
            analysis += f"  Title: {title_idx}, Content: {content_indices}\n\n"
        
        logger.info(f"✅ Template analyzed: {len(prs.slide_layouts)} layouts found")
        
        return analysis, layout_details
        
    except Exception as e:
        logger.error(f"❌ Template analysis failed: {e}")
        return "Template analysis failed. Using default layouts.", {}


def generate_chart(chart_data: Dict[str, Any], output_path: str) -> str:
    try:
        chart_type = chart_data.get("chart_type", "bar").lower()
        title = chart_data.get("title", "Chart")
        data = chart_data.get("data", {})
        x_label = chart_data.get("x_label", "")
        y_label = chart_data.get("y_label", "")
        
        logger.info(f"Generating {chart_type} chart: {title}")
        
        fig, ax = plt.subplots(figsize=(12, 7), facecolor='none')
        ax.patch.set_alpha(0)
        
        if chart_type in ["bar", "column"]:
            labels = data.get("labels", [])
            if "series" in data:
                series = data["series"]
                x = np.arange(len(labels))
                width = 0.8 / len(series)
                
                for i, s in enumerate(series):
                    offset = (i - len(series)/2 + 0.5) * width
                    ax.bar(x + offset, s["values"], width, label=s["name"], alpha=0.8)
                
                ax.set_xticks(x)
                ax.set_xticklabels(labels, rotation=45, ha='right', fontsize=11)
                ax.legend(fontsize=11, framealpha=0.9)
            else:
                values = data.get("values", [])
                colors = sns.color_palette("husl", len(labels))
                ax.bar(labels, values, color=colors, alpha=0.8, edgecolor='white', linewidth=1.5)
                plt.xticks(rotation=45, ha='right', fontsize=11)
        
        elif chart_type == "line":
            labels = data.get("labels", [])
            
            if "series" in data:
                for s in data["series"]:
                    ax.plot(labels, s["values"], marker='o', label=s["name"],
                           linewidth=3, markersize=8, alpha=0.8)
                ax.legend(fontsize=11, framealpha=0.9)
            else:
                values = data.get("values", [])
                ax.plot(labels, values, marker='o', linewidth=3, markersize=8,
                       color=sns.color_palette("husl")[0], alpha=0.8)
            
            plt.xticks(rotation=45, ha='right', fontsize=11)
            ax.grid(True, alpha=0.3, linestyle='--')
        
        elif chart_type == "pie":
            labels = data.get("labels", [])
            values = data.get("values", [])
            colors = sns.color_palette("pastel", len(labels))
            
            wedges, texts, autotexts = ax.pie(
                values, labels=labels, autopct='%1.1f%%',
                startangle=90, colors=colors,
                textprops={'fontsize': 12, 'weight': 'bold'},
                wedgeprops={'edgecolor': 'white', 'linewidth': 2}
            )
            
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontsize(11)
            
            ax.axis('equal')
        
        elif chart_type == "scatter":
            if "series" in data:
                for s in data["series"]:
                    x_vals = s.get("x_values", [])
                    y_vals = s.get("y_values", [])
                    ax.scatter(x_vals, y_vals, label=s["name"], s=150, alpha=0.7,
                              edgecolors='white', linewidth=1.5)
                ax.legend(fontsize=11, framealpha=0.9)
            else:
                x_vals = data.get("x_values", [])
                y_vals = data.get("y_values", [])
                ax.scatter(x_vals, y_vals, s=150, alpha=0.7,
                          color=sns.color_palette("husl")[0],
                          edgecolors='white', linewidth=1.5)
            
            ax.grid(True, alpha=0.3, linestyle='--')
        
        ax.set_title(title, fontsize=16, fontweight='bold', pad=20, color='#2c3e50')
        if x_label:
            ax.set_xlabel(x_label, fontsize=13, fontweight='semibold', color='#34495e')
        if y_label:
            ax.set_ylabel(y_label, fontsize=13, fontweight='semibold', color='#34495e')
        
        ax.tick_params(axis='both', which='major', labelsize=11, colors='#34495e')
        
        if chart_type not in ["pie"]:
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color('#7f8c8d')
            ax.spines['bottom'].set_color('#7f8c8d')
        
        plt.tight_layout(pad=1.5)
        plt.savefig(output_path, dpi=300, bbox_inches='tight',
                   facecolor='none', edgecolor='none', transparent=True)
        plt.close()
        
        logger.info(f"✅ Chart saved: {output_path}")
        
        return output_path
        
    except Exception as e:
        logger.error(f"❌ Error generating chart: {e}")
        raise


def clear_all_slides(prs: Presentation):
    xml_slides = prs.slides._sldIdLst
    
    while len(xml_slides):
        rId = xml_slides[0].rId
        prs.part.drop_rel(rId)
        del xml_slides[0]
    
    logger.info("Cleared all existing slides from template")


def create_ppt_from_json(
    slides: List[Dict[str, Any]],
    template_path: str,
    output_path: str
) -> str:
    try:
        logger.info(f"Creating PPT from {len(slides)} slides...")
        if template_path and Path(template_path).exists():
            prs = Presentation(template_path)
            clear_all_slides(prs)
            logger.info(f"Using template: {template_path}")
        else:
            prs = Presentation()
            logger.info("Using blank presentation")
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        for i, slide_data in enumerate(slides, start=1):
            idx = int(slide_data.get("layout_index", 1))
            if idx >= len(prs.slide_layouts):
                logger.warning(f"Layout index {idx} invalid (max {len(prs.slide_layouts)-1}). Using 1.")
                idx = 1
            
            layout = prs.slide_layouts[idx]
            slide = prs.slides.add_slide(layout)
            
            title = slide_data.get("title", "") or ""
            layout_type = slide_data.get("layout_type", "SINGLE_CONTENT")
            
            set_title = False
            for shp in slide.placeholders:
                if getattr(shp.placeholder_format, "type", None) in (1, 3):  
                    try:
                        shp.text = title
                        set_title = True
                        break
                    except Exception:
                        pass
            
            if not set_title and getattr(slide.shapes, "title", None):
                try:
                    slide.shapes.title.text = title
                    set_title = True
                except Exception:
                    pass
            
            if not set_title:
                logger.warning(f"Could not set title on slide {i}")
            
            if layout_type == "CHART" and "chart" in slide_data:
                chart_spec = slide_data["chart"]
                chart_filename = f"chart_{i}_{uuid_module.uuid4().hex[:8]}.png"
                chart_path = os.path.join(CHARTS_DIR, chart_filename)
                
                try:
                    generate_chart(chart_spec, chart_path)
                    shapes_to_remove = []
                    for shape in slide.shapes:
                        if hasattr(shape, 'placeholder_format'):
                            ph_type = shape.placeholder_format.type
                            if ph_type in (2, 4, 7, 14): 
                                shapes_to_remove.append(shape)
                    
                    for shape in shapes_to_remove:
                        sp = shape.element
                        sp.getparent().remove(sp)
                    
                    top_margin = Inches(1.8)
                    left_margin = Inches(0.8)
                    right_margin = Inches(0.8)
                    bottom_margin = Inches(0.5)
                    
                    available_width = slide_width - left_margin - right_margin
                    available_height = slide_height - top_margin - bottom_margin
                    slide.shapes.add_picture(
                        chart_path,
                        left_margin,
                        top_margin,
                        width=available_width,
                        height=available_height
                    )
                    
                    logger.info(f"✅ Chart added to slide {i}")
                    
                except Exception as e:
                    logger.error(f"❌ Failed to generate chart for slide {i}: {e}")
            else:
                content_arrays: List[List[str]] = slide_data.get("content") or []
                
                if content_arrays:
                    content_ph = []
                    for shp in slide.placeholders:
                        if getattr(shp.placeholder_format, "type", None) in (2, 4, 7, 14):
                            if hasattr(shp, "text_frame"):
                                content_ph.append((shp.placeholder_format.idx, shp))
                    
                    content_ph.sort(key=lambda x: x[0])
                    
                    for a_idx, bullets in enumerate(content_arrays):
                        if a_idx < len(content_ph):
                            _, shp = content_ph[a_idx]
                            tf = shp.text_frame
                            tf.clear()
                            
                            first = True
                            for b in (bullets or []):
                                if first:
                                    p = tf.paragraphs[0]
                                    first = False
                                else:
                                    p = tf.add_paragraph()
                                
                                p.text = str(b)
                                p.level = 0
                        else:
                            logger.warning(f"Slide {i}: content[{a_idx}] has no matching placeholder")
        
        prs.save(output_path)
        logger.info(f"✅ PPT saved: {output_path}")
        
        return output_path
        
    except Exception as e:
        logger.error(f"❌ Error creating PPT: {e}")
        raise


def convert_ppt_to_pdf(ppt_path: str, pdf_path: str) -> str:
    try:
        logger.info(f"Converting PPT to PDF: {ppt_path}")
        soffice_paths = [
            'soffice',  
            r'C:\Program Files\LibreOffice\program\soffice.exe', 
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',  
            '/usr/bin/soffice',  
            '/usr/local/bin/soffice',  
        ]
        
        for soffice_path in soffice_paths:
            try:
                result = subprocess.run([
                    soffice_path,
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', str(Path(pdf_path).parent),
                    ppt_path
                ], capture_output=True, text=True, timeout=60)
                
                if result.returncode == 0 and os.path.exists(pdf_path):
                    logger.info(f"[OK] PDF converted using LibreOffice: {pdf_path}")
                    return pdf_path
            except FileNotFoundError:
                continue 
            except Exception as e:
                logger.debug(f"LibreOffice attempt failed at {soffice_path}: {e}")
                continue
        
        logger.warning("LibreOffice not found in any standard location")
        
        if os.name == 'nt':
            try:
                import comtypes.client
                
                logger.info("Attempting PDF conversion using Microsoft PowerPoint...")
                
                powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                powerpoint.Visible = 1
                
                ppt_abs = os.path.abspath(ppt_path)
                pdf_abs = os.path.abspath(pdf_path)
                
                deck = powerpoint.Presentations.Open(ppt_abs)
                deck.SaveAs(pdf_abs, 32)  
                deck.Close()
                powerpoint.Quit()
                
                if os.path.exists(pdf_path):
                    logger.info(f"[OK] PDF converted using PowerPoint: {pdf_path}")
                    return pdf_path
                
            except ImportError:
                logger.debug("comtypes not installed")
            except Exception as e:
                logger.debug(f"PowerPoint COM automation failed: {e}")
        
        logger.warning("=" * 60)
        logger.warning("NO PDF CONVERSION TOOL AVAILABLE")
        logger.warning("PDF generation skipped - PPTX file is available")
        logger.warning("")
        logger.warning("To enable PDF conversion:")
        logger.warning("  Windows: Download LibreOffice from https://www.libreoffice.org/download/")
        logger.warning("  Linux:   sudo apt install libreoffice -y")
        logger.warning("  macOS:   brew install libreoffice")
        logger.warning("=" * 60)
        
        return ""  
        
    except Exception as e:
        logger.error(f"[ERROR] Error converting PPT to PDF: {e}")
        return "" 
