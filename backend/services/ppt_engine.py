import io
import os
import uuid
import subprocess
import tempfile
from typing import Dict, Any, Tuple, List, Optional

import numpy as np
import matplotlib.pyplot as plt
import matplotlib
import matplotlib.patheffects
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from core.logger import get_logger
from core.config import settings
from core.utils import short_kb

logger = get_logger("ppt")
matplotlib.use('Agg')
sns.set_theme(style="whitegrid")
sns.set_context("talk", font_scale=1.2)

# Create output directory
os.makedirs(settings.OUTPUT_CHARTS_DIR, exist_ok=True)

def analyze_template_layouts(template_bytes: bytes) -> Tuple[str, Dict[int, Dict[str, Any]]]:
    """Analyze PPTX template and return layout details."""
    prs = Presentation(io.BytesIO(template_bytes))
    details: Dict[int, Dict[str, Any]] = {}
    
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = getattr(layout, "placeholders", [])
        phs: List[Dict[str, Any]] = []
        
        for ph in placeholders:
            try:
                phs.append({
                    "idx": ph.placeholder_format.idx,
                    "type": ph.placeholder_format.type,
                    "name": getattr(ph, "name", "")
                })
            except Exception:
                continue
        
        title_idx = next((p["idx"] for p in phs if p["type"] == 1), None)
        content_indices = [p["idx"] for p in phs if p["type"] in (2, 4, 7, 14)]
        
        details[i] = {
            "name": getattr(layout, "name", f"Layout {i}"),
            "title_idx": title_idx,
            "content_indices": content_indices,
        }
    
    logger.info(f"template layouts detected: {len(prs.slide_layouts)}")
    return (f"Template layouts: {len(prs.slide_layouts)} found.", details)


def extract_dominant_colors_from_template(template_bytes: bytes) -> List[tuple]:
    """Extract dominant colors from template."""
    try:
        prs = Presentation(io.BytesIO(template_bytes))
        colors = []
        
        slide_master = prs.slide_master
        
        for shape in slide_master.shapes:
            try:
                if hasattr(shape, 'fill') and shape.fill.type == 1:
                    if hasattr(shape.fill.fore_color, 'rgb'):
                        rgb = shape.fill.fore_color.rgb
                        colors.append((rgb[0], rgb[1], rgb[2]))
                
                if hasattr(shape, 'line') and hasattr(shape.line, 'color'):
                    if hasattr(shape.line.color, 'rgb'):
                        rgb = shape.line.color.rgb
                        colors.append((rgb[0], rgb[1], rgb[2]))
            except Exception:
                continue
        
        unique_colors = []
        for color in colors:
            if sum(color) < 100 or sum(color) > 700:
                continue
            if color not in unique_colors:
                unique_colors.append(color)
        
        if unique_colors:
            logger.info(f"Extracted {len(unique_colors)} colors from template")
            return unique_colors[:6]
        
        logger.info("Using default professional palette")
        return [
            (68, 114, 196),    # Blue
            (237, 125, 49),    # Orange
            (165, 165, 165),   # Gray
            (255, 192, 0),     # Yellow
            (91, 155, 213),    # Light Blue
            (112, 173, 71),    # Green
        ]
    except Exception as e:
        logger.error(f"Color extraction failed: {e}")
        return [(68, 114, 196), (237, 125, 49), (165, 165, 165), (255, 192, 0)]


def rgb_to_matplotlib(rgb_tuple: tuple) -> tuple:
    """Convert RGB (0-255) to matplotlib (0-1)"""
    return tuple(c / 255.0 for c in rgb_tuple)


def generate_chart_png(chart_spec: Dict[str, Any], cache, template_colors: List[tuple]) -> str:
    """
    Generate professional chart with transparent background.
    """
    title = chart_spec.get("title", "Chart")
    chart_type = str(chart_spec.get("chart_type", "bar")).lower()
    data = chart_spec.get("data", {}) or {}
    path = os.path.join(settings.OUTPUT_CHARTS_DIR, f"chart_{uuid.uuid4().hex[:8]}.png")
    
    labels = data.get("labels", []) or []
    colors_mpl = [rgb_to_matplotlib(c) for c in template_colors]
    sns.set_palette(sns.color_palette(colors_mpl))
    
    # Create figure
    fig, ax = plt.subplots(figsize=(12, 7))
    fig.patch.set_alpha(0.0)
    ax.patch.set_alpha(0.0)
    text_color = colors_mpl[0]
    if sum(template_colors[0]) > 600:  # If color is too light
        text_color = (0.2, 0.2, 0.2)  # Use dark gray
    
    try:
        if chart_type in ("bar", "column"):
            values = data.get("values", []) or []
            series = data.get("series", [])
            
            if series:
                import pandas as pd
                plot_data = []
                for label in labels:
                    for s in series:
                        series_values = s.get("values", [])
                        series_name = s.get("name", "")
                        idx = labels.index(label)
                        if idx < len(series_values):
                            plot_data.append({
                                'Category': label,
                                'Value': series_values[idx],
                                'Series': series_name
                            })
                
                df = pd.DataFrame(plot_data)
                sns.barplot(data=df, x='Category', y='Value', hue='Series', ax=ax,
                           palette=colors_mpl, edgecolor='white', linewidth=2)
                ax.legend(title='', fontsize=13, frameon=True, facecolor='white',
                         edgecolor=text_color, framealpha=0.9)
            else:
                bars = sns.barplot(x=labels, y=values, ax=ax,
                                  palette=colors_mpl[:len(labels)],
                                  edgecolor='white', linewidth=2.5)
                for i, (label, value) in enumerate(zip(labels, values)):
                    ax.text(i, value, f'{int(value):,}', ha='center', va='bottom',
                           fontsize=15, fontweight='bold', color=text_color,
                           bbox=dict(boxstyle='round,pad=0.5', facecolor='white',
                                    edgecolor=text_color, alpha=0.9, linewidth=1.5))
            
            ax.tick_params(labelsize=14, colors=text_color, width=2)
            ax.set_xlabel(chart_spec.get("x_label", ""), fontsize=16, fontweight='bold', color=text_color)
            ax.set_ylabel(chart_spec.get("y_label", ""), fontsize=16, fontweight='bold', color=text_color)
            ax.grid(True, alpha=0.4, linestyle='--', axis='y', color=text_color, linewidth=1.5)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color(text_color)
            ax.spines['left'].set_linewidth(2)
            ax.spines['bottom'].set_color(text_color)
            ax.spines['bottom'].set_linewidth(2)
            
        elif chart_type == "line":
            values = data.get("values", []) or []
            
            line = ax.plot(range(len(labels)), values, marker='o', linewidth=4, markersize=16,
                          color=colors_mpl[0],
                          markerfacecolor=colors_mpl[1] if len(colors_mpl) > 1 else colors_mpl[0],
                          markeredgecolor='white', markeredgewidth=3)[0]
            
            ax.set_xticks(range(len(labels)))
            ax.set_xticklabels(labels, fontsize=14, color=text_color, fontweight='bold')
            for i, value in enumerate(values):
                ax.text(i, value, f'{value}', ha='center', va='bottom',
                       fontsize=14, fontweight='bold', color=text_color,
                       bbox=dict(boxstyle='round,pad=0.4', facecolor='white',
                                edgecolor=text_color, alpha=0.9, linewidth=1.5))
            
            ax.tick_params(labelsize=14, colors=text_color, width=2)
            ax.set_xlabel(chart_spec.get("x_label", ""), fontsize=16, fontweight='bold', color=text_color)
            ax.set_ylabel(chart_spec.get("y_label", ""), fontsize=16, fontweight='bold', color=text_color)
            ax.grid(True, alpha=0.4, linestyle='--', color=text_color, linewidth=1.5)
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            ax.spines['left'].set_color(text_color)
            ax.spines['left'].set_linewidth(2)
            ax.spines['bottom'].set_color(text_color)
            ax.spines['bottom'].set_linewidth(2)
            
        elif chart_type == "pie":
            values = data.get("values", []) or []
            if not labels or not values or len(labels) != len(values):
                labels = labels or [f"Segment {i+1}" for i in range(len(values))]
                labels = labels[:len(values)]
            
            wedges, texts, autotexts = ax.pie(values, labels=labels, autopct='%1.1f%%',
                                              colors=colors_mpl[:len(values)], startangle=90,
                                              textprops={'fontsize': 14, 'fontweight': 'bold', 'color': text_color},
                                              wedgeprops={'edgecolor': 'white', 'linewidth': 3, 'antialiased': True})
            
            for autotext in autotexts:
                autotext.set_color('white')
                autotext.set_fontsize(13)
                autotext.set_fontweight('bold')
                autotext.set_path_effects([
                    matplotlib.patheffects.withStroke(linewidth=3, foreground='black', alpha=0.5)
                ])
            
            for text in texts:
                text.set_fontsize(13)
                text.set_fontweight('bold')
                text.set_color(text_color)
        else:
            ax.text(0.5, 0.5, f"Unsupported chart: {chart_type}",
                   ha="center", va="center", fontsize=16, color='red', fontweight='bold')
        title_text = ax.set_title(
            title,
            fontsize=22,
            fontweight='bold',
            pad=20,
            color=text_color
        )
        title_text.set_path_effects([
            matplotlib.patheffects.withStroke(linewidth=4, foreground='white', alpha=0.8)
        ])
        
        fig.tight_layout()
        fig.savefig(path, bbox_inches="tight", facecolor='none', edgecolor='none',
                   transparent=True, dpi=150)
        
        logger.info(f"chart generated '{title}' ({chart_type}) with transparent background")
        
    except Exception as e:
        logger.error(f"Chart generation error: {e}")
        ax.text(0.5, 0.5, f"Chart Error: {str(e)[:50]}", ha="center", va="center",
               fontsize=14, color='red', fontweight='bold')
        fig.savefig(path, bbox_inches="tight", transparent=True, dpi=150)
    finally:
        plt.close(fig)
    
    try:
        cache.add_temp_file(path)
    except Exception:
        pass
    
    return path

def format_text_frame(text_frame, content_list: List[str], template_colors: List[tuple]):
    """Format text with template colors."""
    text_frame.clear()
    primary_color = template_colors[0] if template_colors else (68, 114, 196)
    
    for bullet_idx, bullet_text in enumerate(content_list):
        if bullet_text:
            if bullet_idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = str(bullet_text)
            p.level = 0
            p.alignment = PP_ALIGN.LEFT
            
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(18)
                run.font.bold = False
                run.font.color.rgb = RGBColor(*primary_color)


def build_ppt_from_slides(slides: list, template_bytes: bytes, cache) -> bytes:
    """Build PowerPoint with centered charts and creative formatting."""
    prs = Presentation(io.BytesIO(template_bytes))
    
    logger.info("Analyzing template colors...")
    template_colors = extract_dominant_colors_from_template(template_bytes)
    logger.info(f"Template color palette: {template_colors}")
    
    # Remove existing slides
    xml_slides = getattr(prs.slides, "_sldIdLst", None)
    while xml_slides is not None and len(xml_slides):
        rId = xml_slides[0].rId
        prs.part.drop_rel(rId)
        del xml_slides[0]
    
    for idx, s in enumerate(slides, start=1):
        layout_idx = s.get("layout_index", 1)
        
        if not isinstance(layout_idx, int) or layout_idx < 0:
            layout_idx = 1
        if layout_idx >= len(prs.slide_layouts):
            layout_idx = len(prs.slide_layouts) - 1
        
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        title = s.get("title", "") or ""
        
        # Format title
        try:
            if slide.shapes.title:
                slide.shapes.title.text = title
                if slide.shapes.title.text_frame.paragraphs:
                    title_para = slide.shapes.title.text_frame.paragraphs[0]
                    title_para.alignment = PP_ALIGN.CENTER
                    if title_para.runs:
                        title_run = title_para.runs[0]
                        title_run.font.size = Pt(32)
                        title_run.font.bold = True
                        if template_colors:
                            title_run.font.color.rgb = RGBColor(*template_colors[0])
        except Exception:
            pass
        
        # Handle CHART slides 
        if str(s.get("layout_type", "")).upper() == "CHART" and "chart" in s:
            try:
                chart_path = generate_chart_png(s["chart"], cache, template_colors)
                
                content_placeholder = None
                for shape in slide.placeholders:
                    if shape.placeholder_format.type in (2, 7, 14):
                        content_placeholder = shape
                        break
                
                if content_placeholder:
                    ph_left = content_placeholder.left
                    ph_top = content_placeholder.top
                    ph_width = content_placeholder.width
                    ph_height = content_placeholder.height
                    
                    sp = content_placeholder.element
                    sp.getparent().remove(sp)
                    chart_width = int(ph_width * 0.9)
                    chart_height = int(ph_height * 0.85)
                    chart_left = ph_left + int((ph_width - chart_width) / 2)
                    chart_top = ph_top + int((ph_height - chart_height) / 2)
                    
                    slide.shapes.add_picture(chart_path, chart_left, chart_top,
                                            width=chart_width, height=chart_height)
                else:
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    
                    chart_width = Inches(8)
                    chart_height = Inches(4.5)
                    
                    left = int((slide_width - chart_width) / 2)
                    top = int((slide_height - chart_height) / 2) + Inches(0.5)
                    
                    slide.shapes.add_picture(chart_path, left, top,
                                            width=chart_width, height=chart_height)
            except Exception as e:
                logger.warning(f"chart add failed on slide #{idx}: {e}")
        
        elif "content" in s:
            content_arrays = s.get("content", [])
            non_empty_content = [arr for arr in content_arrays if arr]
            
            if non_empty_content:
                placeholders = [p for p in slide.placeholders
                               if p.placeholder_format.type in (2, 7, 14)]
                
                for pidx, content_list in enumerate(non_empty_content):
                    if pidx < len(placeholders):
                        try:
                            ph = placeholders[pidx]
                            tf = ph.text_frame
                            format_text_frame(tf, content_list, template_colors)
                        except Exception as e:
                            logger.warning(f"content add failed on slide #{idx}, placeholder {pidx}: {e}")
        
        logger.info(f"slide added #{idx} (layout_index={layout_idx}, title='{title}')")
    
    out = io.BytesIO()
    prs.save(out)
    ppt_bytes = out.getvalue()
    
    logger.info(f"PowerPoint built (slides={len(slides)})")
    logger.info(f"PowerPoint ready in memory (size={short_kb(len(ppt_bytes))} KB)")
    
    return ppt_bytes

def convert_ppt_to_pdf(ppt_path: str, pdf_path: str) -> str:
    """Convert PPTX to PDF using LibreOffice."""
    try:
        from pathlib import Path
        logger.info(f"Converting PPT to PDF: {ppt_path}")
        
        ppt_path = os.path.abspath(ppt_path)
        desired_pdf_path = os.path.abspath(pdf_path)
        outdir = str(Path(desired_pdf_path).parent)
        os.makedirs(outdir, exist_ok=True)
        
        expected_pdf = str(Path(outdir) / (Path(ppt_path).stem + ".pdf"))
        
        soffice_paths = [
            "soffice",
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "/usr/bin/soffice",
            "/usr/local/bin/soffice",
        ]
        
        for soffice_path in soffice_paths:
            try:
                result = subprocess.run([soffice_path, "--headless", "--convert-to", "pdf",
                                        "--outdir", outdir, ppt_path],
                                       capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0 and os.path.exists(expected_pdf):
                    final_pdf = expected_pdf
                    if os.path.abspath(expected_pdf) != desired_pdf_path:
                        try:
                            os.replace(expected_pdf, desired_pdf_path)
                            final_pdf = desired_pdf_path
                        except Exception:
                            pass
                    logger.info(f"[OK] PDF converted: {final_pdf}")
                    return final_pdf
            except FileNotFoundError:
                continue
            except Exception as e:
                logger.debug(f"LibreOffice attempt failed: {e}")
        
        logger.warning("LibreOffice not found; PDF conversion skipped")
        return ""
    except Exception as e:
        logger.error(f"Error converting PPT to PDF: {e}")
        return ""


def export_pdf_from_ppt_bytes(ppt_bytes: bytes) -> Optional[bytes]:
    """Convert PPTX bytes to PDF bytes."""
    with tempfile.TemporaryDirectory() as tdir:
        ppt_path = os.path.join(tdir, "input.pptx")
        pdf_path = os.path.join(tdir, "output.pdf")
        
        with open(ppt_path, "wb") as f:
            f.write(ppt_bytes)
        
        final_pdf_path = convert_ppt_to_pdf(ppt_path, pdf_path)
        
        if final_pdf_path and os.path.exists(final_pdf_path):
            with open(final_pdf_path, "rb") as f:
                data = f.read()
            logger.info("PDF bytes ready in memory")
            return data
    return None
